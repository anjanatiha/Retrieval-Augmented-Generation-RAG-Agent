"""document_loader.py — DocumentLoader class: owns all ingestion."""

import csv as _csv
import os
import re
import sys
import tempfile
from typing import List
from urllib.parse import urlparse

from src.rag.config import (
    DOCS_ROOT, DOC_FOLDERS, EXT_TO_TYPE,
    TXT_CHUNK_SIZE, TXT_CHUNK_OVERLAP,
    PDF_CHUNK_SENTENCES, DOCX_CHUNK_PARAS,
    PPTX_CHUNK_SLIDES, HTML_CHUNK_SENTENCES,
)

__all__ = ['DocumentLoader']


class DocumentLoader:
    """Owns all document ingestion: reading files, parsing 9 formats, chunking, URL fetching.

    State:
        doc_folders:  dict mapping type → folder path (from config)
        ext_to_type:  dict mapping file extension → document type
        chunk_sizes:  dict of all chunk-size constants (txt, pdf, docx, pptx, html)

    Public API:
        ensure_folders()        — create ./docs subfolders if missing
        scan_all_files()        — find every file under DOCS_ROOT, flag misplaced ones
        chunk_all_documents()   — scan + dispatch all files → list of chunk dicts
        chunk_url(url)          — fetch a URL, detect type, return chunks
    """

    def __init__(self) -> None:
        """Bind config constants to instance state so chunkers don't import config directly."""
        self.doc_folders = DOC_FOLDERS
        self.ext_to_type = EXT_TO_TYPE
        self.chunk_sizes = {
            'txt_chunk_size':        TXT_CHUNK_SIZE,
            'txt_chunk_overlap':     TXT_CHUNK_OVERLAP,
            'pdf_chunk_sentences':   PDF_CHUNK_SENTENCES,
            'docx_chunk_paras':      DOCX_CHUNK_PARAS,
            'pptx_chunk_slides':     PPTX_CHUNK_SLIDES,
            'html_chunk_sentences':  HTML_CHUNK_SENTENCES,
        }

    # ------------------------------------------------------------------ Public

    def ensure_folders(self) -> None:
        """Create DOCS_ROOT and all subfolders if they don't exist."""
        os.makedirs(DOCS_ROOT, exist_ok=True)
        for folder in self.doc_folders.values():
            if not os.path.exists(folder):
                os.makedirs(folder)
                print(f"  Created folder: {folder}/")

    def scan_all_files(self) -> List[dict]:
        """
        Returns a list of (filepath, detected_type, canonical_folder, is_misplaced).
        Scans every subfolder under DOCS_ROOT so no file is missed.
        """
        found = []
        for folder_type, folder_path in self.doc_folders.items():
            if not os.path.isdir(folder_path):
                continue
            for fname in os.listdir(folder_path):
                ext = os.path.splitext(fname)[1].lower()
                detected_type = self.ext_to_type.get(ext)
                if detected_type is None:
                    continue  # unsupported extension, skip silently
                filepath      = os.path.join(folder_path, fname)
                canonical_dir = self.doc_folders[detected_type]
                is_misplaced  = (folder_path != canonical_dir)
                found.append({
                    'filepath':      filepath,
                    'filename':      fname,
                    'detected_type': detected_type,
                    'found_in':      folder_path,
                    'canonical_dir': canonical_dir,
                    'is_misplaced':  is_misplaced,
                })
                if is_misplaced:
                    print(f"  [MISPLACED] '{fname}' found in '{folder_path}/' "
                          f"— detected as '{detected_type.upper()}', "
                          f"canonical folder is '{canonical_dir}/'. Processing anyway.")
        return found

    def chunk_all_documents(self) -> List[dict]:
        """
        Scans every subfolder under ./docs/, auto-detects file type,
        and routes each file to the correct chunker.
        Files in the wrong folder are still processed (with a notice).
        """
        print("\nLoading documents...")
        file_list = self.scan_all_files()

        if not file_list:
            print(f"\nNo supported documents found under '{DOCS_ROOT}/'")
            print("Supported types: PDF, TXT, DOCX, XLSX, XLS, PPTX, CSV, MD, HTML")
            print("Place files in the matching subfolder (or any subfolder — smart detection handles the rest).")
            sys.exit(1)

        all_chunks  = []
        type_counts = {}

        for file_info in file_list:
            t      = file_info['detected_type']
            chunks = self._dispatch_chunker(file_info)
            all_chunks.extend(chunks)
            type_counts[t] = type_counts.get(t, 0) + len(chunks)
            flag = " [MISPLACED — processed anyway]" if file_info['is_misplaced'] else ""
            print(f"  [{t.upper()}] '{file_info['filename']}': {len(chunks)} chunks{flag}")

        print("\n  Chunk summary by type:")
        for t, count in sorted(type_counts.items()):
            print(f"    {t.upper():<8} {count} chunks")
        print(f"  Total: {len(all_chunks)} chunks\n")

        return all_chunks

    def chunk_url(self, url: str) -> List[dict]:
        """
        Fetches a URL and routes it to the correct chunker based on:
        1. Content-Type header from the response  (most reliable)
        2. File extension in the URL path         (fallback)
        3. PDF magic bytes sniff
        4. Default to HTML
        """
        try:
            import requests
        except ImportError:
            print("  [WARNING] requests not installed. pip install requests")
            return []

        url = url.strip()
        print(f"\n  Fetching URL: {url}")
        try:
            resp = requests.get(url, timeout=30, headers={
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) '
                              'AppleWebKit/537.36 (KHTML, like Gecko) '
                              'Chrome/120.0.0.0 Safari/537.36'
            }, allow_redirects=True, stream=True)
            resp.raise_for_status()
            content = resp.content
            text    = None
        except Exception as e:
            print(f"  [ERROR] Could not fetch URL: {e}")
            return []

        # ── Type detection ──────────────────────────────────────────
        # Priority 1: Content-Type header (most reliable for any URL)
        content_type = resp.headers.get('Content-Type', '').lower().split(';')[0].strip()
        dtype = None

        CONTENT_TYPE_MAP = {
            'application/pdf':                                          'pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
            'application/msword':                                       'docx',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
            'application/vnd.ms-excel':                                 'xlsx',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
            'application/vnd.ms-powerpoint':                            'pptx',
            'text/csv':                                                 'csv',
            'text/plain':                                               'txt',
            'text/markdown':                                            'md',
            'text/html':                                                'html',
            'application/xhtml+xml':                                    'html',
        }
        dtype = CONTENT_TYPE_MAP.get(content_type)

        # Fuzzy fallback for non-standard content-type values
        if dtype is None:
            if 'pdf'          in content_type: dtype = 'pdf'
            elif 'word'       in content_type: dtype = 'docx'
            elif 'excel'      in content_type or 'spreadsheet' in content_type: dtype = 'xlsx'
            elif 'powerpoint' in content_type or 'presentation' in content_type: dtype = 'pptx'
            elif 'csv'        in content_type: dtype = 'csv'

        # Priority 2: file extension in the URL path (ignores query strings)
        if dtype is None:
            parsed_path = urlparse(url).path.lower()
            clean_path  = re.sub(r'[?#].*$', '', parsed_path)
            ext         = os.path.splitext(clean_path)[1]
            dtype       = self.ext_to_type.get(ext)

        # Priority 3: sniff first bytes for PDF magic number
        if dtype is None and content[:4] == b'%PDF':
            dtype = 'pdf'

        # Priority 4: default to HTML
        if dtype is None:
            dtype = 'html'

        # ── Source label ─────────────────────────────────────────────
        parsed      = urlparse(url)
        source_name = (parsed.netloc + parsed.path).rstrip('/')
        if parsed.query:
            source_name += '?' + parsed.query[:20] + ('...' if len(parsed.query) > 20 else '')
        if len(source_name) > 60:
            source_name = source_name[:57] + '...'
        if not source_name:
            source_name = url[:60]

        print(f"  Detected type: {dtype.upper()} (content-type: {content_type})")

        # ── Dispatch to chunker ───────────────────────────────────────
        if dtype in ('pdf', 'docx', 'xlsx', 'pptx', 'xls'):
            suffix = {'pdf':'.pdf','docx':'.docx','xlsx':'.xlsx','pptx':'.pptx','xls':'.xls'}[dtype]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                file_info = {
                    'filepath':      tmp_path,
                    'filename':      source_name,
                    'detected_type': dtype,
                    'is_misplaced':  False,
                }
                chunks = self._dispatch_chunker(file_info)
            finally:
                try: os.unlink(tmp_path)
                except: pass
            if not chunks:
                print(f"  [WARNING] No chunks extracted from {dtype.upper()} at URL. "
                      f"File may be empty, password-protected, or corrupt.")
            return chunks

        # Text formats
        try:
            text = content.decode(resp.encoding or 'utf-8', errors='replace')
        except Exception:
            text = content.decode('utf-8', errors='replace')

        if dtype == 'csv':
            with tempfile.NamedTemporaryFile(delete=False, suffix='.csv', mode='w',
                                             encoding='utf-8', errors='replace') as tmp:
                tmp.write(text)
                tmp_path = tmp.name
            try:
                chunks = self._chunk_csv(tmp_path, source_name)
            finally:
                try: os.unlink(tmp_path)
                except: pass
            return chunks

        if dtype == 'txt':
            lines = [l.strip() for l in text.splitlines() if l.strip()]
            return [{'text': l, 'source': source_name, 'start_line': i+1,
                     'end_line': i+1, 'type': 'txt'} for i, l in enumerate(lines)]

        if dtype == 'md':
            with tempfile.NamedTemporaryFile(delete=False, suffix='.md', mode='w',
                                              encoding='utf-8', errors='replace') as tmp:
                tmp.write(text)
                tmp_path = tmp.name
            try:
                chunks = self._chunk_md(tmp_path, source_name)
            finally:
                try: os.unlink(tmp_path)
                except: pass
            return chunks

        # HTML / webpage
        try:
            from bs4 import BeautifulSoup
            text = BeautifulSoup(text, 'html.parser').get_text(separator=' ', strip=True)
        except ImportError:
            text = re.sub(r'<[^>]+>', ' ', text)

        sents  = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
        chunks = []
        for i in range(0, len(sents), HTML_CHUNK_SENTENCES):
            window = sents[i: i + HTML_CHUNK_SENTENCES]
            if window:
                chunks.append({'text': ' '.join(window), 'source': source_name,
                               'start_line': i+1, 'end_line': i+len(window), 'type': 'html'})
        return chunks

    # ----------------------------------------------------------------- Private

    def _dispatch_chunker(self, file_info: dict) -> List[dict]:
        """Picks the right chunker based on detected_type (not folder)."""
        fp  = file_info['filepath']
        fn  = file_info['filename']
        ext = os.path.splitext(fn)[1].lower()
        t   = file_info['detected_type']

        if t == 'txt':
            return self._chunk_txt(fp, fn)
        elif t == 'md':
            return self._chunk_md(fp, fn)
        elif t == 'pdf':
            return self._chunk_pdf(fp, fn)
        elif t == 'docx':
            return self._chunk_docx(fp, fn)
        elif t == 'xlsx':
            if ext == '.xls':
                return self._chunk_xls(fp, fn)
            elif ext == '.csv':
                return self._chunk_csv(fp, fn)
            else:
                return self._chunk_xlsx(fp, fn)
        elif t == 'csv':
            return self._chunk_csv(fp, fn)
        elif t == 'pptx':
            return self._chunk_pptx(fp, fn)
        elif t == 'html':
            return self._chunk_html(fp, fn)
        else:
            print(f"  [SKIP] No chunker for type '{t}' ({fn})")
            return []

    def _chunk_txt(self, filepath: str, filename: str) -> List[dict]:
        # Splits on non-empty lines; overlap lets sliding window preserve sentence context.
        chunk_size = self.chunk_sizes['txt_chunk_size']
        overlap    = self.chunk_sizes['txt_chunk_overlap']
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            lines = [l.strip() for l in f.readlines() if l.strip()]
        step   = max(1, chunk_size - overlap)
        chunks = []
        for i in range(0, len(lines), step):
            window = lines[i: i + chunk_size]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'txt',
            })
        return chunks

    def _chunk_md(self, filepath: str, filename: str) -> List[dict]:
        """Strips markdown syntax and chunks like plain text."""
        chunk_size = self.chunk_sizes['txt_chunk_size']
        overlap    = self.chunk_sizes['txt_chunk_overlap']
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()
        clean = re.sub(r'#{1,6}\s*', '', raw)
        clean = re.sub(r'[*_`~]{1,3}', '', clean)
        clean = re.sub(r'!\[.*?\]\(.*?\)', '', clean)
        clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)
        lines = [l.strip() for l in clean.splitlines() if l.strip()]
        step  = max(1, chunk_size - overlap)
        chunks = []
        for i in range(0, len(lines), step):
            window = lines[i: i + chunk_size]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'md',
            })
        return chunks

    def _chunk_pdf(self, filepath: str, filename: str) -> List[dict]:
        # Page-level isolation: start_line == end_line == page number so source labels are exact.
        sentences_per_chunk = self.chunk_sizes['pdf_chunk_sentences']
        try:
            import fitz
        except ImportError:
            print("  [WARNING] pymupdf not installed — skipping PDF. Install: pip install pymupdf")
            return []

        chunks   = []
        page_num = 0
        try:
            doc = fitz.open(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        for page_num, page in enumerate(doc, start=1):
            raw = page.get_text()
            if not raw.strip():
                continue
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', raw) if s.strip()]
            for i in range(0, len(sentences), sentences_per_chunk):
                window = sentences[i: i + sentences_per_chunk]
                if not window:
                    continue
                chunks.append({
                    'text':       ' '.join(window),
                    'source':     filename,
                    'start_line': page_num,
                    'end_line':   page_num,
                    'type':       'pdf',
                })
        doc.close()
        return chunks

    def _chunk_docx(self, filepath: str, filename: str) -> List[dict]:
        paras_per_chunk = self.chunk_sizes['docx_chunk_paras']
        try:
            from docx import Document
        except ImportError:
            print("  [WARNING] python-docx not installed — skipping DOCX. Install: pip install python-docx")
            return []

        chunks = []
        try:
            doc = Document(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        # Collect body paragraphs
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables (resumes often use tables for layout)
        for table in doc.tables:
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                # Deduplicate merged cells (python-docx repeats merged cell text)
                seen = []
                for cell in row_cells:
                    if cell not in seen:
                        seen.append(cell)
                if seen:
                    paras.append(' | '.join(seen))

        for i in range(0, len(paras), paras_per_chunk):
            window = paras[i: i + paras_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'docx',
            })
        return chunks

    def _chunk_xlsx(self, filepath: str, filename: str) -> List[dict]:
        """Each non-empty row becomes a chunk (joined as key=value pairs)."""
        try:
            import openpyxl
        except ImportError:
            print("  [WARNING] openpyxl not installed — skipping XLSX. Install: pip install openpyxl")
            return []

        chunks = []
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        for sheet in wb.sheetnames:
            ws     = wb[sheet]
            rows   = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            header = [str(c) if c is not None else f"col{i}" for i, c in enumerate(rows[0])]
            for row_idx, row in enumerate(rows[1:], start=2):
                cells = [str(v) for v in row if v is not None and str(v).strip()]
                if not cells:
                    continue
                pairs = '; '.join(
                    f"{header[i] if i < len(header) else f'col{i}'}={str(row[i])}"
                    for i in range(len(row)) if row[i] is not None and str(row[i]).strip()
                )
                chunks.append({
                    'text':       f"[{sheet}] {pairs}",
                    'source':     filename,
                    'start_line': row_idx,
                    'end_line':   row_idx,
                    'type':       'xlsx',
                })
        wb.close()
        return chunks

    def _chunk_xls(self, filepath: str, filename: str) -> List[dict]:
        """Fallback for old .xls via xlrd."""
        try:
            import xlrd
        except ImportError:
            print("  [WARNING] xlrd not installed — skipping XLS. Install: pip install xlrd")
            return []

        chunks = []
        try:
            wb = xlrd.open_workbook(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        for sheet in wb.sheets():
            if sheet.nrows < 2:
                continue
            header = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
            for row_idx in range(1, sheet.nrows):
                pairs = '; '.join(
                    f"{header[c] if c < len(header) else f'col{c}'}={sheet.cell_value(row_idx, c)}"
                    for c in range(sheet.ncols)
                    if str(sheet.cell_value(row_idx, c)).strip()
                )
                if not pairs:
                    continue
                chunks.append({
                    'text':       f"[{sheet.name}] {pairs}",
                    'source':     filename,
                    'start_line': row_idx + 1,
                    'end_line':   row_idx + 1,
                    'type':       'xlsx',
                })
        return chunks

    def _chunk_csv(self, filepath: str, filename: str) -> List[dict]:
        """Each CSV row becomes a chunk."""
        chunks = []
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                reader = _csv.DictReader(f)
                for row_idx, row in enumerate(reader, start=2):
                    pairs = '; '.join(f"{k}={v}" for k, v in row.items() if v and str(v).strip())
                    if not pairs:
                        continue
                    chunks.append({
                        'text':       pairs,
                        'source':     filename,
                        'start_line': row_idx,
                        'end_line':   row_idx,
                        'type':       'csv',
                    })
        except Exception as e:
            print(f"  [ERROR] Could not read '{filename}': {e}")
        return chunks

    def _chunk_pptx(self, filepath: str, filename: str) -> List[dict]:
        # Collects text from every text shape on each slide; one slide = one chunk by default.
        slides_per_chunk = self.chunk_sizes['pptx_chunk_slides']
        try:
            from pptx import Presentation
        except ImportError:
            print("  [WARNING] python-pptx not installed — skipping PPTX. Install: pip install python-pptx")
            return []

        chunks = []
        try:
            prs = Presentation(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        slide_texts = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(' '.join(texts))

        for i in range(0, len(slide_texts), slides_per_chunk):
            window = slide_texts[i: i + slides_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'pptx',
            })
        return chunks

    def _chunk_html(self, filepath: str, filename: str) -> List[dict]:
        # BeautifulSoup strips all tags; sentence-window chunking then produces coherent text.
        sentences_per_chunk = self.chunk_sizes['html_chunk_sentences']
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            print("  [WARNING] beautifulsoup4 not installed — skipping HTML. Install: pip install beautifulsoup4")
            return []

        chunks = []
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        text      = soup.get_text(separator=' ', strip=True)
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
        for i in range(0, len(sentences), sentences_per_chunk):
            window = sentences[i: i + sentences_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'html',
            })
        return chunks

    def _truncate_chunk(self, text: str, max_words: int = 300, max_chars: int = 1200) -> str:
        """Truncate to 300 words OR 1200 chars, whichever is shorter."""
        words     = text.split()
        truncated = ' '.join(words[:max_words]) if len(words) > max_words else text
        return truncated[:max_chars] if len(truncated) > max_chars else truncated

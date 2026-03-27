"""test_contracts.py — Contract tests for chunkers, URL ingestion, and search internals (HF Space).

Asserts output *shape*, required keys, and value types — not exact values.
This guards against accidental API breakage across model and provider changes.

Scope (this file): chunkers (_chunk_*), chunk_url(), _hybrid_retrieve(), _rerank().
run_pipeline() and agent.run() contracts are in test_contracts_pipeline.py.

HF differences from local:
  - Embeddings via sentence-transformers (_get_st_model), patched by conftest.py.
  - LLM calls via _llm_call (InferenceClient), patched by conftest.py.
  - _rerank uses cross-encoder (also patched here where needed).
  - ChromaDB is EphemeralClient (in-memory).
  - conversation_history (not conversation) is the attribute name.
"""

import os
import sys
from unittest.mock import MagicMock, patch

import pytest
import chromadb
from rank_bm25 import BM25Okapi

# ── make src importable from huggingface/ ────────────────────────────────────
HF_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if HF_ROOT not in sys.path:
    sys.path.insert(0, HF_ROOT)

from tests.conftest import sample_chunks, make_store_with_chunks


# ---------------------------------------------------------------------------
# Chunk-dict contract helper
# ---------------------------------------------------------------------------

CHUNK_KEYS = {
    'text':       str,
    'source':     str,
    'start_line': int,
    'end_line':   int,
    'type':       str,
}


def _assert_chunk_contract(chunk: dict, label: str = '') -> None:
    """Assert *chunk* satisfies the 5-key chunk-dict contract.

    Args:
        chunk: The chunk dict to validate.
        label: Optional identifier used in assertion messages.
    """
    prefix = f"[{label}] " if label else ""
    for key, expected_type in CHUNK_KEYS.items():
        assert key in chunk, f"{prefix}Missing key '{key}'"
        assert chunk[key] is not None, f"{prefix}Key '{key}' is None"
        assert isinstance(chunk[key], expected_type), (
            f"{prefix}Key '{key}': expected {expected_type.__name__}, "
            f"got {type(chunk[key]).__name__}"
        )


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def loader():
    """Fresh DocumentLoader instance for chunker contract tests."""
    from src.rag.document_loader import DocumentLoader
    return DocumentLoader()


@pytest.fixture
def store():
    """VectorStore backed by 5 cat-facts chunks.

    conftest.py autouse fixtures patch _get_st_model and _llm_call globally,
    so no additional patching is needed here.
    """
    chunks = sample_chunks(5)
    return make_store_with_chunks(chunks)


# ---------------------------------------------------------------------------
# 1. Chunker contracts — every _chunk_* method
# ---------------------------------------------------------------------------

class TestChunkerContracts:
    """Every chunker returns a list of dicts satisfying the 5-key contract."""

    def test_chunk_txt_contract(self, loader, tmp_path):
        """_chunk_txt chunks satisfy the 5-key contract."""
        f = tmp_path / 'test.txt'
        f.write_text("Cats sleep 16 hours.\nDogs are loyal companions.\n")
        chunks = loader._chunk_txt(str(f), 'test.txt')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"txt[{i}]")

    def test_chunk_md_contract(self, loader, tmp_path):
        """_chunk_md chunks satisfy the 5-key contract."""
        f = tmp_path / 'test.md'
        f.write_text("# Animals\nCats sleep 16 hours.\nDogs are loyal.\n")
        chunks = loader._chunk_md(str(f), 'test.md')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"md[{i}]")

    def test_chunk_pdf_contract(self, loader, tmp_path):
        """_chunk_pdf chunks satisfy the 5-key contract."""
        import fitz
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Cats sleep sixteen hours. They hunt at night.")
        path = str(tmp_path / 'test.pdf')
        doc.save(path); doc.close()
        chunks = loader._chunk_pdf(path, 'test.pdf')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"pdf[{i}]")

    def test_chunk_docx_contract(self, loader, tmp_path):
        """_chunk_docx chunks satisfy the 5-key contract."""
        from docx import Document
        doc = Document()
        doc.add_paragraph("Cats sleep sixteen hours a day.")
        doc.add_paragraph("They are excellent nocturnal hunters.")
        path = str(tmp_path / 'test.docx')
        doc.save(path)
        chunks = loader._chunk_docx(path, 'test.docx')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"docx[{i}]")

    def test_chunk_xlsx_contract(self, loader, tmp_path):
        """_chunk_xlsx chunks satisfy the 5-key contract."""
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(['animal', 'fact'])
        ws.append(['cat', 'sleeps 16 hours'])
        path = str(tmp_path / 'test.xlsx')
        wb.save(path)
        chunks = loader._chunk_xlsx(path, 'test.xlsx')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"xlsx[{i}]")

    def test_chunk_csv_contract(self, loader, tmp_path):
        """_chunk_csv chunks satisfy the 5-key contract."""
        f = tmp_path / 'test.csv'
        f.write_text("animal,fact\ncat,sleeps 16 hours\ndog,loyal\n")
        chunks = loader._chunk_csv(str(f), 'test.csv')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"csv[{i}]")

    def test_chunk_pptx_contract(self, loader, tmp_path):
        """_chunk_pptx chunks satisfy the 5-key contract."""
        from pptx import Presentation
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text               = "Cat Facts"
        slide.placeholders[1].text_frame.text = "Cats sleep 16 hours a day."
        path  = str(tmp_path / 'test.pptx')
        prs.save(path)
        chunks = loader._chunk_pptx(path, 'test.pptx')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"pptx[{i}]")

    def test_chunk_html_contract(self, loader, tmp_path):
        """_chunk_html chunks satisfy the 5-key contract."""
        f = tmp_path / 'test.html'
        f.write_text(
            "<html><body><p>Cats sleep 16 hours a day. "
            "They are nocturnal hunters.</p></body></html>"
        )
        chunks = loader._chunk_html(str(f), 'test.html')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"html[{i}]")


# ---------------------------------------------------------------------------
# 2. chunk_url contract
# ---------------------------------------------------------------------------

class TestChunkUrlContract:
    """chunk_url() returns a list; each item satisfies the 5-key contract."""

    def test_chunk_url_html_contract(self, loader):
        """chunk_url for an HTML page returns chunk dicts with all required keys."""
        mock_resp             = MagicMock()
        mock_resp.headers     = {'Content-Type': 'text/html'}
        mock_resp.content     = (
            b"<html><body><p>Cats sleep 16 hours a day. "
            b"They are nocturnal hunters.</p></body></html>"
        )
        mock_resp.encoding            = 'utf-8'
        mock_resp.raise_for_status    = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url('http://example.com/page')

        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"url_html[{i}]")

    def test_chunk_url_txt_contract(self, loader):
        """chunk_url for a plain-text URL returns chunk dicts with all required keys."""
        mock_resp             = MagicMock()
        mock_resp.headers     = {'Content-Type': 'text/plain'}
        mock_resp.content     = b"Cats sleep 16 hours a day.\nDogs are loyal.\n"
        mock_resp.encoding    = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url('http://example.com/file.txt')

        assert isinstance(chunks, list)
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"url_txt[{i}]")


# ---------------------------------------------------------------------------
# 3. _hybrid_retrieve contract
# ---------------------------------------------------------------------------

class TestHybridRetrieveContract:
    """_hybrid_retrieve returns list of (dict, float) 2-tuples."""

    def test_returns_list(self, store):
        """_hybrid_retrieve always returns a list."""
        result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        assert isinstance(result, list)

    def test_tuple_shape(self, store):
        """Each element is a (dict, float) 2-tuple."""
        result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        for item in result:
            assert len(item) == 2
            entry, score = item
            assert isinstance(entry, dict)
            assert isinstance(score, float)

    def test_entry_has_required_keys(self, store):
        """Each entry dict in hybrid retrieve has all 5 required keys."""
        result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        for entry, _ in result:
            _assert_chunk_contract(entry, '_hybrid_retrieve entry')

    def test_score_is_float(self, store):
        """All returned hybrid scores are finite floats.

        Note: cosine distance converted to similarity can produce negative values
        when the dense store has only a few near-identical embeddings, so we only
        assert the type here — not the sign.
        """
        result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        for _, score in result:
            assert isinstance(score, float), f"Score {score!r} is not a float"


# ---------------------------------------------------------------------------
# 4. _rerank contract
# ---------------------------------------------------------------------------

class TestRerankContract:
    """_rerank returns list of (dict, float, float) 3-tuples."""

    def test_returns_list(self, store):
        """_rerank always returns a list."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        # Patch cross-encoder predict to return fixed scores
        with patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8, 0.7]
            result = store._rerank('cats sleep', candidates, top_n=2)
        assert isinstance(result, list)

    def test_triple_shape(self, store):
        """Each element is a (dict, float, float) 3-tuple."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        with patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8, 0.7]
            result = store._rerank('cats sleep', candidates, top_n=2)
        for item in result:
            assert len(item) == 3
            entry, sim, rerank_score = item
            assert isinstance(entry, dict)
            assert isinstance(sim,         float)
            assert isinstance(rerank_score, float)

    def test_entry_has_required_keys(self, store):
        """Each entry in rerank result has all 5 required chunk keys."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        with patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8, 0.7]
            result = store._rerank('cats sleep', candidates, top_n=2)
        for entry, _, _ in result:
            _assert_chunk_contract(entry, '_rerank entry')

    def test_sim_is_bounded(self, store):
        """sim value for each result is in [0, 1]."""
        candidates = [(c, 0.75) for c in store.chunks[:2]]
        with patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8, 0.7]
            result = store._rerank('cats sleep', candidates, top_n=2)
        for _, sim, _ in result:
            assert 0.0 <= sim <= 1.0, f"sim {sim} out of [0, 1]"

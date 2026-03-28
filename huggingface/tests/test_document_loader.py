"""test_document_loader.py — File-type chunker tests for the HF Space version.

Covers file-type chunkers: TXT, CSV, MD, PDF, DOCX, XLSX, XLS, PPTX.
HTML chunker, URL ingestion, truncate_chunk, and _dispatch_chunker
are in test_url_ingestion.py.

Mock strategy:
  - requests.get is mocked only where needed (none in this file).
  - Never mock: fitz, python-docx, openpyxl, xlrd, python-pptx, beautifulsoup4,
    BM25Okapi, or truncate_chunk itself.

HF differences from local:
  - No DOC_FOLDERS / scan_all_files / ensure_folders in HF version.
  - HTML chunker filters lines shorter than 40 chars (boilerplate filter).
  - Chunker functions are imported directly from src.rag.chunkers.
"""

import os
import sys
import tempfile

import pytest

# ── make src importable from huggingface/ ────────────────────────────────────
HF_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if HF_ROOT not in sys.path:
    sys.path.insert(0, HF_ROOT)

# ── Text-based chunker functions ─────────────────────────────────────────────
from src.rag.chunkers import (
    chunk_txt, chunk_md, chunk_csv, truncate_chunk,
)

# ── Binary-format chunker functions (moved to binary_chunkers.py) ─────────────
from src.rag.binary_chunkers import (
    chunk_pdf, chunk_docx, chunk_xlsx, chunk_xls, chunk_pptx,
)


# ═══════════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _write_tmp(content: str, suffix: str) -> str:
    """Write *content* to a temporary file with the given *suffix* and return its path.

    The caller is responsible for deleting the file after the test.
    """
    f = tempfile.NamedTemporaryFile(delete=False, suffix=suffix,
                                    mode='w', encoding='utf-8')
    f.write(content)
    f.close()
    return f.name


# ═══════════════════════════════════════════════════════════════════════════════
# TXT, MD, and CSV (plain-text chunkers)
# ═══════════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderTxt:
    """Tests for plain-text, Markdown, and CSV chunkers from src.rag.chunkers."""

    def test_chunk_txt_basic(self):
        """Three non-empty lines → three txt chunks with correct source and type."""
        path = _write_tmp("Line one\nLine two\nLine three\n", '.txt')
        try:
            chunks = chunk_txt(path, 'test.txt')
        finally:
            os.unlink(path)
        assert len(chunks) == 3
        assert all(c['type'] == 'txt' for c in chunks)
        assert all(c['source'] == 'test.txt' for c in chunks)

    def test_chunk_txt_empty_lines_skipped(self):
        """Blank lines are silently dropped — only non-empty lines produce chunks."""
        path = _write_tmp("Hello\n\n\nWorld\n", '.txt')
        try:
            chunks = chunk_txt(path, 'test.txt')
        finally:
            os.unlink(path)
        assert len(chunks) == 2

    def test_chunk_txt_start_line_is_1(self):
        """The first chunk from a txt file has start_line equal to 1."""
        path = _write_tmp("a\nb\nc\n", '.txt')
        try:
            chunks = chunk_txt(path, 'test.txt')
        finally:
            os.unlink(path)
        assert chunks[0]['start_line'] == 1

    def test_chunk_md_strips_syntax(self):
        """Markdown headings, bold/italic/code markers, and link syntax are removed."""
        md = "# Heading\n**bold** and _italic_ and `code`\n[link](http://x.com)\n"
        path = _write_tmp(md, '.md')
        try:
            chunks = chunk_md(path, 'test.md')
        finally:
            os.unlink(path)
        assert chunks
        combined = ' '.join(c['text'] for c in chunks)
        assert '#' not in combined
        assert '**' not in combined
        assert '[link]' not in combined

    def test_chunk_md_type_is_md(self):
        """Chunks produced by chunk_md always carry type='md'."""
        path = _write_tmp("Some markdown content\n", '.md')
        try:
            chunks = chunk_md(path, 'test.md')
        finally:
            os.unlink(path)
        assert all(c['type'] == 'md' for c in chunks)

    def test_chunk_csv_key_value_pairs(self):
        """Each CSV data row becomes one chunk formatted as 'col=value; ...' pairs."""
        csv_content = "name,age,city\nAlice,30,NYC\nBob,25,LA\n"
        path = _write_tmp(csv_content, '.csv')
        try:
            chunks = chunk_csv(path, 'data.csv')
        finally:
            os.unlink(path)
        assert len(chunks) == 2
        assert 'name=Alice' in chunks[0]['text']
        assert 'age=30' in chunks[0]['text']
        assert all(c['type'] == 'csv' for c in chunks)

    def test_chunk_csv_empty_rows_skipped(self):
        """Rows where all values are empty are not added to the chunk list."""
        csv_content = "name,age\nAlice,30\n,\nBob,25\n"
        path = _write_tmp(csv_content, '.csv')
        try:
            chunks = chunk_csv(path, 'data.csv')
        finally:
            os.unlink(path)
        texts = [c['text'] for c in chunks]
        assert any('Alice' in t for t in texts)
        assert any('Bob' in t for t in texts)


# ═══════════════════════════════════════════════════════════════════════════════
# PDF
# ═══════════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderPdf:
    """Tests for the PDF chunker using real fitz (pymupdf)."""

    def _make_pdf(self, text="Hello world. This is a test sentence. Another sentence here."):
        """Create a minimal single-page PDF in a temp file and return its path.

        Skips the test if PyMuPDF is not installed.
        """
        try:
            import fitz
        except ImportError:
            pytest.skip("pymupdf not installed")
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), text)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        doc.save(tmp.name)
        doc.close()
        tmp.close()
        return tmp.name

    def test_chunk_pdf_returns_chunks(self):
        """Real PDF → at least one chunk with type='pdf' and correct source."""
        path = self._make_pdf()
        try:
            chunks = chunk_pdf(path, 'test.pdf')
        finally:
            os.unlink(path)
        assert len(chunks) > 0
        assert all(c['type'] == 'pdf' for c in chunks)
        assert all(c['source'] == 'test.pdf' for c in chunks)

    def test_chunk_pdf_page_number_in_start_line(self):
        """start_line on each PDF chunk equals the 1-based page number."""
        path = self._make_pdf()
        try:
            chunks = chunk_pdf(path, 'test.pdf')
        finally:
            os.unlink(path)
        assert chunks[0]['start_line'] == 1

    def test_chunk_pdf_bad_file_returns_empty(self):
        """A file with non-PDF content returns an empty list without raising."""
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf',
                                          mode='w', encoding='utf-8')
        tmp.write("not a pdf")
        tmp.close()
        try:
            chunks = chunk_pdf(tmp.name, 'bad.pdf')
        finally:
            os.unlink(tmp.name)
        assert chunks == []


# ═══════════════════════════════════════════════════════════════════════════════
# DOCX
# ═══════════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderDocx:
    """Tests for the DOCX chunker including table extraction and merged-cell deduplication."""

    def _make_docx(self, paragraphs=None, table_rows=None):
        """Create a minimal DOCX file with given paragraphs and optional table rows.

        Skips the test if python-docx is not installed.
        """
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")
        doc = Document()
        for p in (paragraphs or ["First paragraph.", "Second paragraph.", "Third paragraph."]):
            doc.add_paragraph(p)
        if table_rows:
            table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
            for i, row_data in enumerate(table_rows):
                for j, val in enumerate(row_data):
                    table.rows[i].cells[j].text = val
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        doc.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_chunk_docx_basic(self):
        """DOCX with three paragraphs → at least one chunk with type='docx'."""
        path = self._make_docx()
        try:
            chunks = chunk_docx(path, 'test.docx')
        finally:
            os.unlink(path)
        assert len(chunks) > 0
        assert all(c['type'] == 'docx' for c in chunks)

    def test_chunk_docx_table_rows_extracted(self):
        """Table row values appear in the combined chunk text."""
        path = self._make_docx(table_rows=[["Name", "Age"], ["Alice", "30"]])
        try:
            chunks = chunk_docx(path, 'test.docx')
        finally:
            os.unlink(path)
        all_text = ' '.join(c['text'] for c in chunks)
        assert 'Alice' in all_text

    def test_chunk_docx_merged_cells_deduplicated(self):
        """Merged cells repeat the same text — should appear only once per row."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")
        doc = Document()
        table = doc.add_table(rows=1, cols=3)
        # Simulate merged cell by setting the same text in all three cells
        for cell in table.rows[0].cells:
            cell.text = "MergedCell"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        doc.save(tmp.name)
        tmp.close()
        try:
            chunks = chunk_docx(tmp.name, 'test.docx')
        finally:
            os.unlink(tmp.name)
        all_text = ' '.join(c['text'] for c in chunks)
        # Deduplication means "MergedCell" appears exactly once
        assert all_text.count('MergedCell') == 1


# ═══════════════════════════════════════════════════════════════════════════════
# XLSX
# ═══════════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderXlsx:
    """Tests for the XLSX chunker using openpyxl."""

    def _make_xlsx(self, rows=None):
        """Create a minimal XLSX workbook with the given rows and return its path.

        Skips the test if openpyxl is not installed.
        """
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl not installed")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        data = rows or [["Name", "Age", "City"], ["Alice", 30, "NYC"], ["Bob", 25, "LA"]]
        for row in data:
            ws.append(row)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
        wb.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_chunk_xlsx_key_value_pairs(self):
        """Each data row becomes a chunk with header=value pairs in the text."""
        path = self._make_xlsx()
        try:
            chunks = chunk_xlsx(path, 'data.xlsx')
        finally:
            os.unlink(path)
        assert len(chunks) == 2
        assert any('Name=Alice' in c['text'] for c in chunks)
        assert all(c['type'] == 'xlsx' for c in chunks)

    def test_chunk_xlsx_row_number_in_start_line(self):
        """start_line for the first data row is 2 because row 1 is the header."""
        path = self._make_xlsx()
        try:
            chunks = chunk_xlsx(path, 'data.xlsx')
        finally:
            os.unlink(path)
        # Row 1 is the header, so the first data row has start_line == 2
        assert chunks[0]['start_line'] == 2

    def test_chunk_xlsx_sheet_name_in_text(self):
        """Sheet name is prefixed in brackets in each chunk's text field."""
        path = self._make_xlsx()
        try:
            chunks = chunk_xlsx(path, 'data.xlsx')
        finally:
            os.unlink(path)
        assert all('[Sheet1]' in c['text'] for c in chunks)

    def test_chunk_xlsx_empty_rows_skipped(self):
        """Rows where all cells are None are skipped; only valid rows produce chunks."""
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl not installed")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Name", "Age"])
        ws.append([None, None])
        ws.append(["Bob", 25])
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
        wb.save(tmp.name)
        tmp.close()
        try:
            chunks = chunk_xlsx(tmp.name, 'data.xlsx')
        finally:
            os.unlink(tmp.name)
        assert len(chunks) == 1


# ═══════════════════════════════════════════════════════════════════════════════
# XLS (legacy)
# ═══════════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderXls:
    """Tests for the legacy .xls chunker (xlrd path)."""

    def test_chunk_xls_basic(self):
        """Legacy .xls file is readable via xlrd and produces correct chunks."""
        try:
            import xlrd   # noqa: F401
        except ImportError:
            pytest.skip("xlrd not installed")
        try:
            import xlwt
        except ImportError:
            pytest.skip("xlwt not installed (needed to create .xls for testing)")
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Sheet1")
        ws.write(0, 0, "Name");  ws.write(0, 1, "Score")
        ws.write(1, 0, "Alice"); ws.write(1, 1, 95)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.xls')
        wb.save(tmp.name)
        tmp.close()
        try:
            chunks = chunk_xls(tmp.name, 'data.xls')
        finally:
            os.unlink(tmp.name)
        assert len(chunks) >= 1
        assert any('Alice' in c['text'] for c in chunks)


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderPptx:
    """Tests for the PPTX chunker using python-pptx."""

    def _make_pptx(self, slides_text=None):
        """Create a minimal PPTX presentation with one textbox per slide.

        Skips the test if python-pptx is not installed.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        for text in (slides_text or ["Slide one content.", "Slide two content."]):
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox.text_frame.text = text
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_chunk_pptx_one_chunk_per_slide(self):
        """Three slides → three chunks when PPTX_CHUNK_SLIDES is 1."""
        path = self._make_pptx(["Slide one.", "Slide two.", "Slide three."])
        try:
            chunks = chunk_pptx(path, 'deck.pptx')
        finally:
            os.unlink(path)
        assert len(chunks) == 3
        assert all(c['type'] == 'pptx' for c in chunks)

    def test_chunk_pptx_slide_number_in_start_line(self):
        """start_line on each PPTX chunk equals the 1-based slide index."""
        path = self._make_pptx(["Slide A.", "Slide B."])
        try:
            chunks = chunk_pptx(path, 'deck.pptx')
        finally:
            os.unlink(path)
        assert chunks[0]['start_line'] == 1
        assert chunks[1]['start_line'] == 2

    def test_chunk_pptx_text_content_present(self):
        """Slide text content appears in the chunk text field."""
        path = self._make_pptx(["Hello from slide one"])
        try:
            chunks = chunk_pptx(path, 'deck.pptx')
        finally:
            os.unlink(path)
        assert 'Hello from slide one' in chunks[0]['text']

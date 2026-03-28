"""test_combinations.py — Combination matrix tests for the HF Space.

Scope (this file): chat/agent mode × doc types, URL fetch × content types, chunker contract.
Truncation, classification, and source-label combinations are in test_combinations_analysis.py.

HF differences from local:
  - Embeddings via sentence-transformers (patched in conftest.py).
  - LLM calls via _llm_call (patched in conftest.py).
  - _rerank uses cross-encoder (patched per test where needed).
  - ChromaDB is EphemeralClient (in-memory via conftest.py unique-collection fixture).
  - _expand_query is disabled (always returns [original_query]).

Mock strategy (per CLAUDE.md):
  Always mock: InferenceClient (_llm_call), _get_st_model — both via conftest.py autouse
  Never mock:  fitz, python-docx, openpyxl, python-pptx, beautifulsoup4, BM25Okapi
"""

import os
import sys
from unittest.mock import MagicMock, patch

import pytest

# ── make src importable from huggingface/ ────────────────────────────────────
HF_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if HF_ROOT not in sys.path:
    sys.path.insert(0, HF_ROOT)

from tests.conftest import sample_chunks, make_store_with_chunks

from src.rag.chunkers import (
    chunk_txt, chunk_md, chunk_csv, chunk_html,
)
from src.rag.binary_chunkers import (
    chunk_pdf, chunk_docx, chunk_xlsx, chunk_pptx,
)

# Map doc type string → chunker function (used instead of getattr on loader)
_CHUNKER_MAP = {
    'txt':  chunk_txt,
    'md':   chunk_md,
    'csv':  chunk_csv,
    'html': chunk_html,
    'pdf':  chunk_pdf,
    'docx': chunk_docx,
    'xlsx': chunk_xlsx,
    'pptx': chunk_pptx,
}


# ---------------------------------------------------------------------------
# Binary file fixture helpers (real parsers — never mocked)
# ---------------------------------------------------------------------------

def _make_pdf_file(tmp_path, filename='test.pdf'):
    """Write a real single-page PDF to tmp_path and return the absolute path."""
    import fitz
    doc  = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Cats sleep sixteen hours a day. They are nocturnal hunters.")
    path = str(tmp_path / filename)
    doc.save(path); doc.close()
    return path


def _make_docx_file(tmp_path, filename='test.docx'):
    """Write a real DOCX file to tmp_path and return the absolute path."""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Cats sleep sixteen hours a day.")
    doc.add_paragraph("They are excellent nocturnal hunters.")
    path = str(tmp_path / filename)
    doc.save(path)
    return path


def _make_xlsx_file(tmp_path, filename='test.xlsx'):
    """Write a real XLSX file to tmp_path and return the absolute path."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(['animal', 'fact'])
    ws.append(['cat', 'sleeps 16 hours a day'])
    path = str(tmp_path / filename)
    wb.save(path)
    return path


def _make_pptx_file(tmp_path, filename='test.pptx'):
    """Write a real PPTX file to tmp_path and return the absolute path."""
    from pptx import Presentation
    prs   = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text               = "Cat Facts"
    slide.placeholders[1].text_frame.text = "Cats sleep 16 hours a day."
    path  = str(tmp_path / filename)
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# 1. Chat mode × all 8 doc types
# ---------------------------------------------------------------------------

class TestChatModeAllDocTypes:
    """run_pipeline returns a non-empty response for every supported doc type."""

    @pytest.mark.parametrize("doc_type,text_content,filename", [
        (
            "txt",
            "Cats sleep 16 hours a day.\nDogs are loyal companions.",
            "test.txt",
        ),
        (
            "md",
            "# Animals\n\nCats sleep 16 hours.\nDogs are loyal companions.",
            "test.md",
        ),
        (
            "csv",
            "animal,fact\ncat,sleeps 16 hours\ndog,loyal companion",
            "test.csv",
        ),
        (
            "html",
            "<html><body><p>Cats sleep 16 hours a day. They are nocturnal hunters. Dogs are loyal companions. Cats have sharp retractable claws. Dogs bark to communicate.</p></body></html>",
            "test.html",
        ),
    ])
    def test_text_doc_type_pipeline(self, doc_type, text_content, filename, tmp_path):
        """Pipeline returns a non-empty response for text-based doc types.

        Args:
            doc_type:     One of txt, md, csv, html.
            text_content: Raw string written to the temp file.
            filename:     Filename passed to the chunker.
        """
        f = tmp_path / filename
        f.write_text(text_content)
        chunker = _CHUNKER_MAP[doc_type]
        chunks  = chunker(str(f), filename)
        assert len(chunks) >= 1, f"No chunks for {doc_type}"

        store = make_store_with_chunks(chunks)
        with patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8] * len(chunks)
            result = store.run_pipeline('How long do cats sleep?')

        assert isinstance(result['response'], str)
        assert len(result['response']) > 0, f"Empty response for doc_type={doc_type}"

    @pytest.mark.parametrize("doc_type,make_fn", [
        ("pdf",  _make_pdf_file),
        ("docx", _make_docx_file),
        ("xlsx", _make_xlsx_file),
        ("pptx", _make_pptx_file),
    ])
    def test_binary_doc_type_pipeline(self, doc_type, make_fn, tmp_path):
        """Pipeline returns a non-empty response for binary doc types.

        Args:
            doc_type: One of pdf, docx, xlsx, pptx.
            make_fn:  Fixture helper that creates a real binary temp file.
        """
        path     = make_fn(tmp_path)
        filename = os.path.basename(path)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(path, filename)
        assert len(chunks) >= 1, f"No chunks for {doc_type}"

        store = make_store_with_chunks(chunks)
        with patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8] * len(chunks)
            result = store.run_pipeline('How long do cats sleep?')

        assert isinstance(result['response'], str)
        assert len(result['response']) > 0, f"Empty response for doc_type={doc_type}"


# ---------------------------------------------------------------------------
# 2. Agent mode × all 8 doc types
# ---------------------------------------------------------------------------

class TestAgentModeAllDocTypes:
    """agent.run() returns a non-empty answer for every supported doc type."""

    @pytest.mark.parametrize("doc_type,text_content,filename", [
        ("txt",  "Cats sleep 16 hours a day.\nDogs are loyal.",         "test.txt"),
        ("md",   "# Animals\nCats sleep 16 hours.\nDogs are loyal.",    "test.md"),
        ("csv",  "animal,fact\ncat,sleeps 16 hours\ndog,loyal",         "test.csv"),
        ("html", "<html><body><p>Cats sleep 16 hours a day. They are nocturnal hunters. Dogs are loyal companions. Cats have sharp claws. Dogs bark loudly.</p></body></html>","test.html"),
    ])
    def test_text_doc_type_agent(self, doc_type, text_content, filename, tmp_path):
        """Agent returns a non-empty answer for text-based doc types.

        Args:
            doc_type:     One of txt, md, csv, html.
            text_content: Raw string written to the temp file.
            filename:     Filename passed to the chunker.
        """
        from src.rag.agent import Agent

        f = tmp_path / filename
        f.write_text(text_content)
        chunker = _CHUNKER_MAP[doc_type]
        chunks  = chunker(str(f), filename)
        assert len(chunks) >= 1

        store = make_store_with_chunks(chunks)
        agent = Agent(store)

        with patch.object(store, '_llm_chat',
                          return_value='TOOL: finish(Cats sleep 16 hours a day.)'), \
             patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8] * len(chunks)
            result = agent.run('How long do cats sleep?', streamlit_mode=True)

        assert isinstance(result['answer'], str)
        assert len(result['answer']) > 0, f"Empty agent answer for doc_type={doc_type}"

    @pytest.mark.parametrize("doc_type,make_fn", [
        ("pdf",  _make_pdf_file),
        ("docx", _make_docx_file),
        ("xlsx", _make_xlsx_file),
        ("pptx", _make_pptx_file),
    ])
    def test_binary_doc_type_agent(self, doc_type, make_fn, tmp_path):
        """Agent returns a non-empty answer for binary doc types.

        Args:
            doc_type: One of pdf, docx, xlsx, pptx.
            make_fn:  Fixture helper that creates a real binary temp file.
        """
        from src.rag.agent import Agent

        path     = make_fn(tmp_path)
        filename = os.path.basename(path)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(path, filename)
        assert len(chunks) >= 1

        store = make_store_with_chunks(chunks)
        agent = Agent(store)

        with patch.object(store, '_llm_chat',
                          return_value='TOOL: finish(Cats sleep 16 hours a day.)'), \
             patch('src.rag.vector_store._get_cross_encoder') as mock_ce:
            mock_ce.return_value.predict.return_value = [0.8] * len(chunks)
            result = agent.run('How long do cats sleep?', streamlit_mode=True)

        assert isinstance(result['answer'], str)
        assert len(result['answer']) > 0, f"Empty agent answer for doc_type={doc_type}"


# ---------------------------------------------------------------------------
# 3. URL fetch × all URL content types
# ---------------------------------------------------------------------------

class TestUrlFetchContentTypes:
    """chunk_url returns non-empty chunks for every supported content-type."""

    @pytest.mark.parametrize("content_type,url_path,expected_type", [
        ("text/html",    "http://example.com/page",     "html"),
        ("text/plain",   "http://example.com/file.txt", "txt"),
        ("text/csv",     "http://example.com/data.csv", "csv"),
        ("text/markdown","http://example.com/doc.md",   "md"),
        ("application/json", "http://example.com/page", "html"),  # fallback to html
    ])
    def test_text_content_type_produces_chunks(self, content_type, url_path, expected_type):
        """chunk_url for text-based Content-Type returns chunks with correct shape.

        Args:
            content_type:  HTTP Content-Type header to inject.
            url_path:      URL string to pass to chunk_url.
            expected_type: Expected doc type in returned chunks.
        """
        from src.rag.document_loader import DocumentLoader
        loader = DocumentLoader()

        body = (
            b"<html><body><p>Cats sleep 16 hours a day. "
            b"They are nocturnal hunters.</p></body></html>"
            if expected_type in ('html',)
            else b"Cats sleep 16 hours a day.\nDogs are loyal.\n"
        )
        mock_resp                  = MagicMock()
        mock_resp.headers          = {'Content-Type': content_type}
        mock_resp.content          = body
        mock_resp.encoding         = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url(url_path)

        assert isinstance(chunks, list)
        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"[{content_type}] chunk {i} missing key '{key}'"

    @pytest.mark.parametrize("content_type,url_path,expected_type", [
        ("application/pdf",          "http://example.com/doc",      "pdf"),
        ("application/octet-stream", "http://example.com/d.xlsx",   "xlsx"),
        ("application/octet-stream", "http://example.com/d.pptx",   "pptx"),
        ("application/octet-stream", "http://example.com/d.docx",   "docx"),
        ("application/octet-stream", "http://example.com/data.csv", "csv"),
    ])
    def test_binary_content_type_produces_chunks(self, content_type, url_path, expected_type, tmp_path):
        """chunk_url for binary Content-Type runs real parsers on actual file bytes.

        Args:
            content_type:  HTTP Content-Type header to inject.
            url_path:      URL string.
            expected_type: Expected doc type (selects fixture helper).
        """
        from src.rag.document_loader import DocumentLoader

        if expected_type == 'pdf':
            path = _make_pdf_file(tmp_path)
        elif expected_type == 'xlsx':
            path = _make_xlsx_file(tmp_path)
        elif expected_type == 'pptx':
            path = _make_pptx_file(tmp_path)
        elif expected_type == 'docx':
            path = _make_docx_file(tmp_path)
        elif expected_type == 'csv':
            path = str(tmp_path / 'data.csv')
            with open(path, 'w') as fh:
                fh.write("animal,fact\ncat,sleeps 16 hours\n")
        else:
            pytest.skip(f"No fixture for {expected_type}")

        with open(path, 'rb') as fh:
            real_bytes = fh.read()

        loader                     = DocumentLoader()
        mock_resp                  = MagicMock()
        mock_resp.headers          = {'Content-Type': content_type}
        mock_resp.content          = real_bytes
        mock_resp.encoding         = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url(url_path)

        assert isinstance(chunks, list)
        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"[{content_type}] chunk {i} missing key '{key}'"


# ---------------------------------------------------------------------------
# 4. URL type detection × detection priority
# ---------------------------------------------------------------------------

class TestUrlTypeDetection:
    """URL type is correctly inferred by each of the 4 detection priorities."""

    @pytest.mark.parametrize("content_type,url_path,expected_type", [
        ("text/html",            "http://example.com/page",      "html"),
        ("application/pdf",      "http://example.com/doc",       "pdf"),
        ("application/octet-stream", "http://example.com/d.csv", "csv"),
        ("text/plain",           "http://example.com/f.txt",     "txt"),
        ("application/octet-stream", "http://example.com/d.xlsx","xlsx"),
        ("application/octet-stream", "http://example.com/d.pptx","pptx"),
    ])
    def test_url_type_by_content_type(self, content_type, url_path, expected_type, tmp_path):
        """Detected URL type matches expected_type for each Content-Type header.

        Args:
            content_type:  HTTP Content-Type to inject.
            url_path:      URL string.
            expected_type: Expected 'type' in returned chunks.
        """
        from src.rag.document_loader import DocumentLoader
        loader = DocumentLoader()

        # Build real bytes for binary types so parsers don't fail
        if expected_type == 'pdf':
            path = _make_pdf_file(tmp_path)
            with open(path, 'rb') as fh: body = fh.read()
        elif expected_type == 'xlsx':
            path = _make_xlsx_file(tmp_path)
            with open(path, 'rb') as fh: body = fh.read()
        elif expected_type == 'pptx':
            path = _make_pptx_file(tmp_path)
            with open(path, 'rb') as fh: body = fh.read()
        elif expected_type in ('csv', 'txt'):
            body = b"animal,fact\ncat,sleeps 16 hours\n"
        else:
            body = b"<html><body><p>Cats sleep 16 hours a day. They are nocturnal hunters. Dogs are loyal companions. Cats have sharp claws. Dogs bark loudly.</p></body></html>"

        mock_resp                  = MagicMock()
        mock_resp.headers          = {'Content-Type': content_type}
        mock_resp.content          = body
        mock_resp.encoding         = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url(url_path)

        if chunks:
            detected = chunks[0]['type']
            if expected_type == 'html':
                assert detected in ('html', 'webpage'), (
                    f"Expected html/webpage, got {detected}"
                )
            else:
                assert detected == expected_type, (
                    f"content_type={content_type!r}: "
                    f"expected {expected_type!r}, got {detected!r}"
                )


# ---------------------------------------------------------------------------
# 5. Chunker contract × every doc type (parametrized)
# ---------------------------------------------------------------------------

class TestChunkerContractParametrized:
    """Every chunker returns dicts with all 5 required keys — parametrized across types."""

    @pytest.mark.parametrize("doc_type,text_content,filename", [
        ("txt",  "Cats sleep 16 hours.\nDogs are loyal.",             "test.txt"),
        ("md",   "# Animals\nCats sleep 16 hours.\nDogs loyal.",      "test.md"),
        ("csv",  "animal,fact\ncat,sleeps 16 hours\ndog,loyal",       "test.csv"),
        ("html", "<html><body><p>Cats sleep 16 hours a day. They are nocturnal hunters. Dogs are loyal companions. Cats have sharp claws. Dogs bark loudly.</p></body></html>", "test.html"),
    ])
    def test_text_chunker_5_key_contract(self, doc_type, text_content, filename, tmp_path):
        """Text-format chunkers produce dicts with all 5 required keys.

        Args:
            doc_type:     One of txt, md, csv, html.
            text_content: Raw string written to the temp file.
            filename:     Filename passed to the chunker.
        """
        f = tmp_path / filename
        f.write_text(text_content)
        chunker = _CHUNKER_MAP[doc_type]
        chunks  = chunker(str(f), filename)

        assert isinstance(chunks, list),   f"_chunk_{doc_type} did not return a list"
        assert len(chunks) >= 1,           f"_chunk_{doc_type} returned empty list"

        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"_chunk_{doc_type} chunk[{i}] missing key '{key}'"
            assert isinstance(c['text'],       str), f"chunk[{i}]['text'] not str"
            assert isinstance(c['source'],     str), f"chunk[{i}]['source'] not str"
            assert isinstance(c['start_line'], int), f"chunk[{i}]['start_line'] not int"
            assert isinstance(c['end_line'],   int), f"chunk[{i}]['end_line'] not int"
            assert isinstance(c['type'],       str), f"chunk[{i}]['type'] not str"

    @pytest.mark.parametrize("doc_type,make_fn", [
        ("pdf",  _make_pdf_file),
        ("docx", _make_docx_file),
        ("xlsx", _make_xlsx_file),
        ("pptx", _make_pptx_file),
    ])
    def test_binary_chunker_5_key_contract(self, doc_type, make_fn, tmp_path):
        """Binary-format chunkers produce dicts with all 5 required keys.

        Args:
            doc_type: One of pdf, docx, xlsx, pptx.
            make_fn:  Fixture helper that creates a real binary temp file.
        """
        path     = make_fn(tmp_path)
        filename = os.path.basename(path)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(path, filename)

        assert isinstance(chunks, list),   f"_chunk_{doc_type} did not return a list"
        assert len(chunks) >= 1,           f"_chunk_{doc_type} returned empty list"

        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"_chunk_{doc_type} chunk[{i}] missing key '{key}'"
            assert isinstance(c['text'],       str), f"chunk[{i}]['text'] not str"
            assert isinstance(c['source'],     str), f"chunk[{i}]['source'] not str"
            assert isinstance(c['start_line'], int), f"chunk[{i}]['start_line'] not int"
            assert isinstance(c['end_line'],   int), f"chunk[{i}]['end_line'] not int"
            assert isinstance(c['type'],       str), f"chunk[{i}]['type'] not str"

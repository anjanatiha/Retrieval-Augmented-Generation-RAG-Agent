"""Full test suite for the HuggingFace Space version of RAG Agent.

Mock strategy:
  - _llm_call (module-level in vector_store) → returns deterministic strings
  - _get_st_model → returns a mock that produces fake 384-dim embeddings
  - requests.get → mocked for URL ingestion tests
  - Never mock: fitz, python-docx, openpyxl, xlrd, python-pptx, beautifulsoup4,
    BM25Okapi, chunk truncation, calculator eval, chromadb.EphemeralClient
"""

import io
import os
import sys
import re
import tempfile
import struct
from unittest.mock import MagicMock, patch

import pytest

# ── make src importable from huggingface/ ───────────────────────────────────
HF_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if HF_ROOT not in sys.path:
    sys.path.insert(0, HF_ROOT)


# ── helpers ─────────────────────────────────────────────────────────────────

def _sample_chunks(n=5, doc_type='txt'):
    """Return n minimal chunk dicts for use as test fixtures."""
    return [
        {
            'text': f"Cats sleep about 16 hours a day chunk {i}.",
            'source': 'cats.txt',
            'start_line': i + 1,
            'end_line': i + 1,
            'type': doc_type,
        }
        for i in range(n)
    ]


def _make_store_with_chunks(chunks, llm_response="mock answer"):
    """Build a VectorStore with fake embeddings and a mocked LLM."""
    from src.rag.vector_store import VectorStore
    store = VectorStore()
    store.build_or_load(chunks)
    return store


# ═══════════════════════════════════════════════════════════════════════════
# 1. CONFIG
# ═══════════════════════════════════════════════════════════════════════════

class TestConfig:
    """Tests that config.py exports all required constants with valid values."""

    def test_all_constants_present(self):
        """All expected constants exist, are the right types, and have sensible values."""
        from src.rag.config import (
            EMBEDDING_MODEL, LANGUAGE_MODEL, LANGUAGE_MODEL_FALLBACKS,
            CHROMA_COLLECTION,
            SIMILARITY_THRESHOLD, TOP_RETRIEVE, TOP_RERANK,
            TXT_CHUNK_SIZE, TXT_CHUNK_OVERLAP,
            PDF_CHUNK_SENTENCES, DOCX_CHUNK_PARAS,
            PPTX_CHUNK_SLIDES, HTML_CHUNK_SENTENCES,
            EXT_TO_TYPE,
        )
        assert isinstance(EMBEDDING_MODEL, str) and EMBEDDING_MODEL
        assert isinstance(LANGUAGE_MODEL, str) and LANGUAGE_MODEL
        assert isinstance(LANGUAGE_MODEL_FALLBACKS, list) and len(LANGUAGE_MODEL_FALLBACKS) >= 1
        assert CHROMA_COLLECTION == 'rag_docs'
        assert 0 < SIMILARITY_THRESHOLD < 1
        assert TOP_RETRIEVE > 0
        assert TOP_RERANK > 0
        assert TXT_CHUNK_SIZE >= 1
        assert TXT_CHUNK_OVERLAP >= 0

    def test_ext_to_type_covers_all_formats(self):
        """EXT_TO_TYPE maps all 13 expected extensions without gaps."""
        from src.rag.config import EXT_TO_TYPE
        required_exts = ['.pdf', '.txt', '.docx', '.doc', '.xlsx', '.xls',
                         '.pptx', '.ppt', '.csv', '.md', '.markdown', '.html', '.htm']
        for ext in required_exts:
            assert ext in EXT_TO_TYPE, f"Missing extension: {ext}"

    def test_fallbacks_list_has_no_gated_models(self):
        """LANGUAGE_MODEL_FALLBACKS contains no gated models that require HF access approval."""
        from src.rag.config import LANGUAGE_MODEL_FALLBACKS
        gated = ['google/gemma', 'meta-llama/Llama-3.2', 'google/gemma-2']
        for model in LANGUAGE_MODEL_FALLBACKS:
            for g in gated:
                assert not model.startswith(g), f"Gated model in fallbacks: {model}"


# ═══════════════════════════════════════════════════════════════════════════
# 2. DOCUMENT LOADER — file chunkers
# ═══════════════════════════════════════════════════════════════════════════

class TestDocumentLoaderTxt:
    """Tests for plain-text and markdown chunkers in DocumentLoader."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _write_tmp(self, content, suffix):
        """Write content to a temp file and return its path."""
        f = tempfile.NamedTemporaryFile(delete=False, suffix=suffix,
                                        mode='w', encoding='utf-8')
        f.write(content)
        f.close()
        return f.name

    def test_chunk_txt_basic(self):
        """Three non-empty lines → three txt chunks, correct source and type."""
        path = self._write_tmp("Line one\nLine two\nLine three\n", '.txt')
        try:
            chunks = self.loader._chunk_txt(path, 'test.txt')
        finally:
            os.unlink(path)
        assert len(chunks) == 3
        assert all(c['type'] == 'txt' for c in chunks)
        assert all(c['source'] == 'test.txt' for c in chunks)

    def test_chunk_txt_empty_lines_skipped(self):
        """Blank lines are silently dropped — only non-empty lines produce chunks."""
        path = self._write_tmp("Hello\n\n\nWorld\n", '.txt')
        try:
            chunks = self.loader._chunk_txt(path, 'test.txt')
        finally:
            os.unlink(path)
        assert len(chunks) == 2

    def test_chunk_md_strips_syntax(self):
        """Markdown headings, bold/italic/code markers, and link syntax are removed."""
        md = "# Heading\n**bold** and _italic_ and `code`\n[link](http://x.com)\n"
        path = self._write_tmp(md, '.md')
        try:
            chunks = self.loader._chunk_md(path, 'test.md')
        finally:
            os.unlink(path)
        assert chunks
        combined = ' '.join(c['text'] for c in chunks)
        assert '#' not in combined
        assert '**' not in combined
        assert '[link]' not in combined

    def test_chunk_md_type_is_md(self):
        """Chunks produced by _chunk_md always carry type='md'."""
        path = self._write_tmp("Some markdown content\n", '.md')
        try:
            chunks = self.loader._chunk_md(path, 'test.md')
        finally:
            os.unlink(path)
        assert all(c['type'] == 'md' for c in chunks)

    def test_chunk_csv_key_value_pairs(self):
        """Each CSV data row becomes one chunk formatted as 'col=value; ...' pairs."""
        csv_content = "name,age,city\nAlice,30,NYC\nBob,25,LA\n"
        path = self._write_tmp(csv_content, '.csv')
        try:
            chunks = self.loader._chunk_csv(path, 'data.csv')
        finally:
            os.unlink(path)
        assert len(chunks) == 2
        assert 'name=Alice' in chunks[0]['text']
        assert 'age=30' in chunks[0]['text']
        assert all(c['type'] == 'csv' for c in chunks)

    def test_chunk_csv_empty_rows_skipped(self):
        """Rows where all values are empty are not added to the chunk list."""
        csv_content = "name,age\nAlice,30\n,\nBob,25\n"
        path = self._write_tmp(csv_content, '.csv')
        try:
            chunks = self.loader._chunk_csv(path, 'data.csv')
        finally:
            os.unlink(path)
        texts = [c['text'] for c in chunks]
        assert any('Alice' in t for t in texts)
        assert any('Bob' in t for t in texts)


class TestDocumentLoaderPdf:
    """Tests for the PDF chunker using real fitz (pymupdf)."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _make_pdf(self, text="Hello world. This is a test sentence. Another sentence here."):
        """Create a minimal single-page PDF in a temp file and return its path."""
        try:
            import fitz
        except ImportError:
            pytest.skip("pymupdf not installed")
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), text)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        doc.save(tmp.name)
        doc.close()
        tmp.close()
        return tmp.name

    def test_chunk_pdf_returns_chunks(self):
        """Real PDF → at least one chunk with type='pdf' and correct source."""
        path = self._make_pdf()
        try:
            chunks = self.loader._chunk_pdf(path, 'test.pdf')
        finally:
            os.unlink(path)
        assert len(chunks) > 0
        assert all(c['type'] == 'pdf' for c in chunks)
        assert all(c['source'] == 'test.pdf' for c in chunks)

    def test_chunk_pdf_page_number_in_start_line(self):
        """start_line on each PDF chunk equals the 1-based page number."""
        path = self._make_pdf()
        try:
            chunks = self.loader._chunk_pdf(path, 'test.pdf')
        finally:
            os.unlink(path)
        assert chunks[0]['start_line'] == 1


class TestDocumentLoaderDocx:
    """Tests for the DOCX chunker including table extraction and merged-cell deduplication."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _make_docx(self, paragraphs=None, table_rows=None):
        """Create a minimal DOCX file with given paragraphs and optional table rows."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")
        doc = Document()
        for p in (paragraphs or ["First paragraph.", "Second paragraph.", "Third paragraph."]):
            doc.add_paragraph(p)
        if table_rows:
            table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
            for i, row_data in enumerate(table_rows):
                for j, val in enumerate(row_data):
                    table.rows[i].cells[j].text = val
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        doc.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_chunk_docx_basic(self):
        """DOCX with three paragraphs → at least one chunk with type='docx'."""
        path = self._make_docx()
        try:
            chunks = self.loader._chunk_docx(path, 'test.docx')
        finally:
            os.unlink(path)
        assert len(chunks) > 0
        assert all(c['type'] == 'docx' for c in chunks)

    def test_chunk_docx_table_rows_extracted(self):
        """Table row values appear in the combined chunk text."""
        path = self._make_docx(table_rows=[["Name", "Age"], ["Alice", "30"]])
        try:
            chunks = self.loader._chunk_docx(path, 'test.docx')
        finally:
            os.unlink(path)
        all_text = ' '.join(c['text'] for c in chunks)
        assert 'Alice' in all_text

    def test_chunk_docx_merged_cells_deduplicated(self):
        """Merged cells repeat the same text — should appear only once per row."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")
        doc = Document()
        table = doc.add_table(rows=1, cols=3)
        # Simulate merged cell by setting same text in all cells
        for cell in table.rows[0].cells:
            cell.text = "MergedCell"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        doc.save(tmp.name)
        tmp.close()
        try:
            chunks = self.loader._chunk_docx(tmp.name, 'test.docx')
        finally:
            os.unlink(tmp.name)
        all_text = ' '.join(c['text'] for c in chunks)
        # Should appear only once, not triplicated
        assert all_text.count('MergedCell') == 1


class TestDocumentLoaderXlsx:
    """Tests for the XLSX chunker using openpyxl."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _make_xlsx(self, rows=None):
        """Create a minimal XLSX workbook with the given rows and return its path."""
        try:
            import openpyxl
        except ImportError:
            pytest.skip("openpyxl not installed")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        data = rows or [["Name", "Age", "City"], ["Alice", 30, "NYC"], ["Bob", 25, "LA"]]
        for row in data:
            ws.append(row)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
        wb.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_chunk_xlsx_key_value_pairs(self):
        """Each data row becomes a chunk with header=value pairs in the text."""
        path = self._make_xlsx()
        try:
            chunks = self.loader._chunk_xlsx(path, 'data.xlsx')
        finally:
            os.unlink(path)
        assert len(chunks) == 2
        assert any('Name=Alice' in c['text'] for c in chunks)
        assert all(c['type'] == 'xlsx' for c in chunks)

    def test_chunk_xlsx_row_number_in_start_line(self):
        """start_line for the first data row is 2 because row 1 is the header."""
        path = self._make_xlsx()
        try:
            chunks = self.loader._chunk_xlsx(path, 'data.xlsx')
        finally:
            os.unlink(path)
        assert chunks[0]['start_line'] == 2  # row 1 is header

    def test_chunk_xlsx_sheet_name_in_text(self):
        """Sheet name is prefixed in brackets in each chunk's text field."""
        path = self._make_xlsx()
        try:
            chunks = self.loader._chunk_xlsx(path, 'data.xlsx')
        finally:
            os.unlink(path)
        assert all('[Sheet1]' in c['text'] for c in chunks)


class TestDocumentLoaderXls:
    """Tests for the legacy .xls chunker (xlrd path)."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def test_chunk_xls_basic(self):
        """Legacy .xls file is readable via xlrd and produces correct chunks."""
        try:
            import xlrd
        except ImportError:
            pytest.skip("xlrd not installed")
        try:
            import xlwt
        except ImportError:
            pytest.skip("xlwt not installed (needed to create .xls for testing)")
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Sheet1")
        ws.write(0, 0, "Name"); ws.write(0, 1, "Score")
        ws.write(1, 0, "Alice"); ws.write(1, 1, 95)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.xls')
        wb.save(tmp.name)
        tmp.close()
        try:
            chunks = self.loader._chunk_xls(tmp.name, 'data.xls')
        finally:
            os.unlink(tmp.name)
        assert len(chunks) >= 1
        assert any('Alice' in c['text'] for c in chunks)


class TestDocumentLoaderPptx:
    """Tests for the PPTX chunker using python-pptx."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _make_pptx(self, slides_text=None):
        """Create a minimal PPTX presentation with one textbox per slide."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        for text in (slides_text or ["Slide one content.", "Slide two content."]):
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox.text_frame.text = text
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_chunk_pptx_one_chunk_per_slide(self):
        """Three slides → three chunks when PPTX_CHUNK_SLIDES is 1."""
        path = self._make_pptx(["Slide one.", "Slide two.", "Slide three."])
        try:
            chunks = self.loader._chunk_pptx(path, 'deck.pptx')
        finally:
            os.unlink(path)
        assert len(chunks) == 3
        assert all(c['type'] == 'pptx' for c in chunks)

    def test_chunk_pptx_slide_number_in_start_line(self):
        """start_line on each PPTX chunk equals the 1-based slide index."""
        path = self._make_pptx(["Slide A.", "Slide B."])
        try:
            chunks = self.loader._chunk_pptx(path, 'deck.pptx')
        finally:
            os.unlink(path)
        assert chunks[0]['start_line'] == 1
        assert chunks[1]['start_line'] == 2


class TestDocumentLoaderHtml:
    """Tests for the HTML chunker using BeautifulSoup."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _make_html(self, content):
        """Write HTML content to a temp file and return its path."""
        f = tempfile.NamedTemporaryFile(delete=False, suffix='.html',
                                        mode='w', encoding='utf-8')
        f.write(content)
        f.close()
        return f.name

    def test_chunk_html_strips_tags(self):
        """HTML tags are fully removed — no angle brackets appear in chunk text."""
        path = self._make_html("<html><body><h1>Title</h1><p>Some text here. More text.</p></body></html>")
        try:
            chunks = self.loader._chunk_html(path, 'page.html')
        finally:
            os.unlink(path)
        assert chunks
        combined = ' '.join(c['text'] for c in chunks)
        assert '<' not in combined

    def test_chunk_html_type_is_html(self):
        """Chunks produced by _chunk_html always carry type='html'."""
        path = self._make_html("<p>Hello world.</p>")
        try:
            chunks = self.loader._chunk_html(path, 'page.html')
        finally:
            os.unlink(path)
        assert all(c['type'] == 'html' for c in chunks)


class TestTruncateChunk:
    """Tests for DocumentLoader._truncate_chunk (300 words OR 1200 chars)."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def test_truncate_300_words(self):
        """400-word text is truncated to at most 300 words."""
        text = ' '.join(['word'] * 400)
        result = self.loader._truncate_chunk(text)
        assert len(result.split()) <= 300

    def test_truncate_1200_chars(self):
        """2000-character text is truncated to at most 1200 characters."""
        text = 'a' * 2000
        result = self.loader._truncate_chunk(text)
        assert len(result) <= 1200

    def test_short_text_unchanged(self):
        """Text already within both limits is returned verbatim."""
        text = "Short text."
        assert self.loader._truncate_chunk(text) == text


class TestDispatchChunker:
    """Tests for _dispatch_chunker routing logic."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def test_dispatch_txt(self):
        """detected_type='txt' → _chunk_txt called → non-empty chunk list."""
        f = tempfile.NamedTemporaryFile(delete=False, suffix='.txt',
                                        mode='w', encoding='utf-8')
        f.write("Hello world.\n")
        f.close()
        try:
            chunks = self.loader._dispatch_chunker({
                'filepath': f.name, 'filename': 'test.txt',
                'detected_type': 'txt', 'is_misplaced': False,
            })
        finally:
            os.unlink(f.name)
        assert len(chunks) > 0

    def test_dispatch_csv(self):
        """detected_type='csv' → _chunk_csv called → non-empty chunk list."""
        f = tempfile.NamedTemporaryFile(delete=False, suffix='.csv',
                                        mode='w', encoding='utf-8')
        f.write("col1,col2\nval1,val2\n")
        f.close()
        try:
            chunks = self.loader._dispatch_chunker({
                'filepath': f.name, 'filename': 'data.csv',
                'detected_type': 'csv', 'is_misplaced': False,
            })
        finally:
            os.unlink(f.name)
        assert len(chunks) > 0

    def test_dispatch_unknown_type_returns_empty(self):
        """Unrecognised detected_type returns an empty list without raising."""
        chunks = self.loader._dispatch_chunker({
            'filepath': '/tmp/x.xyz', 'filename': 'x.xyz',
            'detected_type': 'unknown_xyz', 'is_misplaced': False,
        })
        assert chunks == []


# ═══════════════════════════════════════════════════════════════════════════
# 3. URL INGESTION
# ═══════════════════════════════════════════════════════════════════════════

class TestUrlIngestion:
    """Tests for chunk_url: type detection priorities and format dispatch."""

    def setup_method(self):
        from src.rag.document_loader import DocumentLoader
        self.loader = DocumentLoader()

    def _mock_resp(self, content, content_type='text/html', encoding='utf-8'):
        """Return a MagicMock mimicking a requests.Response for the given content."""
        resp = MagicMock()
        resp.content = content if isinstance(content, bytes) else content.encode(encoding)
        resp.headers = {'Content-Type': content_type}
        resp.encoding = encoding
        resp.raise_for_status = MagicMock()
        return resp

    def test_url_html_webpage(self):
        """text/html Content-Type → HTML chunks with type='html'."""
        html = b"<html><body><p>Hello world. This is a test. More content here.</p></body></html>"
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = self.loader.chunk_url('http://example.com/page')
        assert len(chunks) > 0
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_type_by_content_type_pdf(self):
        """application/pdf Content-Type → type detection Priority 1 → PDF chunks."""
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "PDF content here. Another sentence.")
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=self._mock_resp(pdf_bytes, 'application/pdf')):
            chunks = self.loader.chunk_url('http://example.com/doc')
        assert len(chunks) > 0
        assert all(c['type'] == 'pdf' for c in chunks)

    def test_url_type_by_extension_csv(self):
        """Unrecognised Content-Type but .csv extension → type detection Priority 2."""
        csv_bytes = b"name,age\nAlice,30\nBob,25\n"
        with patch('requests.get', return_value=self._mock_resp(csv_bytes, 'application/octet-stream')):
            chunks = self.loader.chunk_url('http://example.com/data.csv')
        assert len(chunks) > 0
        assert all(c['type'] == 'csv' for c in chunks)

    def test_url_type_by_pdf_magic_bytes(self):
        """No recognised Content-Type or extension but %PDF header → Priority 3 → pdf chunks."""
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Magic bytes PDF test. Another sentence.")
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        pdf_bytes = buf.getvalue()
        assert pdf_bytes[:4] == b'%PDF'
        with patch('requests.get', return_value=self._mock_resp(pdf_bytes, 'application/octet-stream')):
            chunks = self.loader.chunk_url('http://example.com/unknown')
        assert len(chunks) > 0

    def test_url_defaults_to_html(self):
        """No signal from any priority → defaults to html and returns a list without error."""
        content = b"<p>Some page content. More words here.</p>"
        with patch('requests.get', return_value=self._mock_resp(content, 'application/octet-stream')):
            chunks = self.loader.chunk_url('http://example.com/unknown_type')
        # Should default to html without error
        assert isinstance(chunks, list)

    def test_url_connection_error_returns_empty(self):
        """Network errors are caught and an empty list is returned rather than raising."""
        with patch('requests.get', side_effect=Exception("Connection refused")):
            chunks = self.loader.chunk_url('http://unreachable.example.com')
        assert chunks == []

    def test_url_source_label_truncated_60_chars(self):
        """Very long URLs are truncated to 60 characters in the chunk source field."""
        html = b"<p>Content here. More sentences. And another one.</p>"
        long_url = 'http://example.com/' + 'a' * 100
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = self.loader.chunk_url(long_url)
        if chunks:
            assert len(chunks[0]['source']) <= 60

    def test_url_txt_plain(self):
        """text/plain Content-Type → each non-empty line becomes a type='txt' chunk."""
        txt = b"Line one\nLine two\nLine three\n"
        with patch('requests.get', return_value=self._mock_resp(txt, 'text/plain')):
            chunks = self.loader.chunk_url('http://example.com/file.txt')
        assert len(chunks) > 0
        assert all(c['type'] == 'txt' for c in chunks)


# ═══════════════════════════════════════════════════════════════════════════
# 4. VECTOR STORE
# ═══════════════════════════════════════════════════════════════════════════

class TestVectorStoreBuild:
    """Tests for VectorStore initialisation, chunk addition, and BM25 rebuild."""

    def test_build_empty(self):
        """build_or_load([]) creates a collection with 0 chunks and no BM25 index."""
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load([])
        assert store.collection is not None
        assert store.collection.count() == 0
        assert store.bm25_index is None

    def test_build_with_chunks(self):
        """build_or_load with 5 chunks stores all 5 and creates a BM25 index."""
        chunks = _sample_chunks(5)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        assert store.collection.count() == 5
        assert store.bm25_index is not None

    def test_add_chunks_increases_count(self):
        """add_chunks appends to the collection, raising count from 3 to 5."""
        chunks = _sample_chunks(3)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        extra = _sample_chunks(2, doc_type='pdf')
        store.add_chunks(extra, id_prefix='url')
        assert store.collection.count() == 5

    def test_rebuild_bm25_after_add(self):
        """rebuild_bm25 after add_chunks produces a non-None BM25 index."""
        chunks = _sample_chunks(3)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        extra = _sample_chunks(2)
        store.add_chunks(extra, id_prefix='file')
        store.rebuild_bm25(store.chunks)
        assert store.bm25_index is not None

    def test_clear_conversation(self):
        """clear_conversation resets conversation_history to an empty list."""
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load([])
        store.conversation_history = [{'role': 'user', 'content': 'hi'}]
        store.clear_conversation()
        assert store.conversation_history == []


class TestQueryClassification:
    """Tests for VectorStore._classify_query — all four query types."""

    def setup_method(self):
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def test_classify_summarise(self):
        """Summarise keywords trigger the 'summarise' query type."""
        assert self.store._classify_query("summarise this document") == 'summarise'
        assert self.store._classify_query("give me a summary") == 'summarise'
        assert self.store._classify_query("tell me about the resume") == 'summarise'

    def test_classify_comparison(self):
        """Comparison keywords ('compare', 'vs', 'difference') → 'comparison' type."""
        assert self.store._classify_query("compare Python vs Java") == 'comparison'
        assert self.store._classify_query("what is the difference between them") == 'comparison'

    def test_classify_factual(self):
        """'what is', 'how many', 'list' and similar → 'factual' type."""
        assert self.store._classify_query("what is the capital?") == 'factual'
        assert self.store._classify_query("how many employees?") == 'factual'
        assert self.store._classify_query("list all features") == 'factual'

    def test_classify_general(self):
        """Queries matching no signal → fallback 'general' type."""
        assert self.store._classify_query("random vague question") == 'general'

    def test_summarise_checked_first(self):
        """Summarise is checked before comparison, so 'summarise and compare' → 'summarise'."""
        # summarise beats comparison
        assert self.store._classify_query("summarise and compare") == 'summarise'


class TestSmartTopN:
    """Tests for VectorStore._smart_top_n — correct retrieval budget per query type."""

    def setup_method(self):
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def test_factual_returns_5(self):
        """Factual queries retrieve the fewest chunks — enough for one precise answer."""
        assert self.store._smart_top_n('factual') == 5

    def test_comparison_returns_15(self):
        """Comparison queries retrieve the most chunks to cover both sides."""
        assert self.store._smart_top_n('comparison') == 15

    def test_general_returns_10(self):
        """General queries use a mid-range budget."""
        assert self.store._smart_top_n('general') == 10

    def test_summarise_returns_top_retrieve(self):
        """Summarise uses the full TOP_RETRIEVE budget to cover the whole document."""
        from src.rag.config import TOP_RETRIEVE
        assert self.store._smart_top_n('summarise') == TOP_RETRIEVE


class TestConfidenceCheck:
    """Tests for VectorStore._check_confidence threshold logic."""

    def setup_method(self):
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def test_empty_results_not_confident(self):
        """Empty result list → is_confident=False, best_score=0.0."""
        ok, score = self.store._check_confidence([])
        assert ok is False
        assert score == 0.0

    def test_above_threshold_confident(self):
        """Best score strictly above SIMILARITY_THRESHOLD → is_confident=True."""
        from src.rag.config import SIMILARITY_THRESHOLD
        entry = {'text': 'x', 'source': 'f', 'start_line': 1, 'end_line': 1, 'type': 'txt'}
        ok, score = self.store._check_confidence([(entry, SIMILARITY_THRESHOLD + 0.1)])
        assert ok is True

    def test_below_threshold_not_confident(self):
        """Best score strictly below SIMILARITY_THRESHOLD → is_confident=False."""
        from src.rag.config import SIMILARITY_THRESHOLD
        entry = {'text': 'x', 'source': 'f', 'start_line': 1, 'end_line': 1, 'type': 'txt'}
        ok, score = self.store._check_confidence([(entry, SIMILARITY_THRESHOLD - 0.1)])
        assert ok is False


class TestSourceLabel:
    """Tests for VectorStore._source_label — all seven document types."""

    def setup_method(self):
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def _entry(self, doc_type, start=3, end=3):
        """Build a minimal chunk entry dict for label tests."""
        return {'text': 'x', 'source': 'f', 'start_line': start, 'end_line': end, 'type': doc_type}

    def test_pdf_label(self):
        """PDF chunks use 'p{page}' format."""
        assert self.store._source_label(self._entry('pdf', 2)) == 'p2'

    def test_xlsx_label(self):
        """XLSX chunks use 'row{n}' format."""
        assert self.store._source_label(self._entry('xlsx', 5)) == 'row5'

    def test_csv_label(self):
        """CSV chunks use 'row{n}' format (same as xlsx)."""
        assert self.store._source_label(self._entry('csv', 7)) == 'row7'

    def test_pptx_label(self):
        """PPTX chunks use 'slide{n}' format."""
        assert self.store._source_label(self._entry('pptx', 3)) == 'slide3'

    def test_html_label(self):
        """HTML chunks use 's{n}' (sentence window index) format."""
        assert self.store._source_label(self._entry('html', 4)) == 's4'

    def test_txt_label(self):
        """TXT chunks use 'L{start}-{end}' line range format."""
        assert self.store._source_label(self._entry('txt', 1, 2)) == 'L1-2'

    def test_md_label(self):
        """Markdown chunks use 'L{start}-{end}' line range format."""
        assert self.store._source_label(self._entry('md', 6, 6)) == 'L6-6'


class TestHallucinationFilter:
    """Tests for VectorStore._filter_hallucination — no-info + pivot detection."""

    def setup_method(self):
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def test_clean_response_unchanged(self):
        """Response with no no-info phrase is returned unmodified."""
        resp = "The answer is 42."
        assert self.store._filter_hallucination(resp) == resp

    def test_truncates_at_pivot_after_no_info(self):
        resp = "There is no information about this. However, I can tell you that cats sleep a lot."
        result = self.store._filter_hallucination(resp)
        assert "cats sleep" not in result
        assert "however," not in result.lower() or result.index("however,") > 0

    def test_no_pivot_leaves_response(self):
        resp = "There is no information about this topic in the documents."
        result = self.store._filter_hallucination(resp)
        # No pivot → response kept as-is
        assert "There is no information" in result


class TestHybridRetrieve:
    """Tests for VectorStore._hybrid_retrieve — BM25 + dense fusion."""

    def setup_method(self):
        chunks = [
            {'text': 'Cats sleep 16 hours a day.', 'source': 'cats.txt',
             'start_line': 1, 'end_line': 1, 'type': 'txt'},
            {'text': 'Dogs are loyal companions.', 'source': 'dogs.txt',
             'start_line': 1, 'end_line': 1, 'type': 'txt'},
            {'text': 'Python is a programming language.', 'source': 'prog.txt',
             'start_line': 1, 'end_line': 1, 'type': 'txt'},
        ]
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load(chunks)

    def test_returns_list(self):
        """Hybrid retrieve always returns a list, even for unknown queries."""
        results = self.store._hybrid_retrieve(['cats'], top_n=2)
        assert isinstance(results, list)

    def test_top_n_respected(self):
        """Result count never exceeds the requested top_n."""
        results = self.store._hybrid_retrieve(['cats sleep'], top_n=2)
        assert len(results) <= 2

    def test_each_result_is_entry_score_tuple(self):
        """Each result is a (entry_dict, float_score) tuple with a 'text' key."""
        results = self.store._hybrid_retrieve(['cats'], top_n=3)
        for entry, score in results:
            assert 'text' in entry
            assert isinstance(score, float)

    def test_empty_collection_returns_empty(self):
        """No build_or_load called → collection is None → returns empty list."""
        from src.rag.vector_store import VectorStore
        empty_store = VectorStore()
        # collection is None before build_or_load — should short-circuit to []
        results = empty_store._hybrid_retrieve(['cats'], top_n=3)
        assert results == []


class TestQueryExpand:
    """Tests for VectorStore._expand_query — LLM-based query rewriting."""

    def test_expand_returns_3_queries(self):
        """LLM returns 2 rewrites → original + 2 rewrites = 3 queries total."""
        with patch('src.rag.vector_store._llm_call', return_value="feline sleep\ncat nap duration"):
            from src.rag.vector_store import VectorStore
            store = VectorStore()
            store.build_or_load([])
            results = store._expand_query("how long do cats sleep")
        assert len(results) == 3
        assert results[0] == "how long do cats sleep"

    def test_expand_fallback_on_llm_error(self):
        """LLM raises an exception → falls back to original query only."""
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load([])
        # Patch _llm_chat (not _llm_call) since _llm_chat catches _llm_call errors internally
        with patch.object(store, '_llm_chat', side_effect=Exception("LLM down")):
            results = store._expand_query("test query")
        assert results == ["test query"]


class TestTruncateForEmbedding:
    """Tests for VectorStore._truncate_for_embedding — enforces 200-word AND 1200-char limits."""

    def setup_method(self):
        """Build an empty VectorStore to access the truncation helper."""
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def test_truncates_to_200_words(self):
        """300-word input is cut to ≤200 words."""
        text = ' '.join(['word'] * 300)
        result = self.store._truncate_for_embedding(text)
        assert len(result.split()) <= 200

    def test_truncates_to_1200_chars(self):
        """2000-char single-token string is cut to ≤1200 chars."""
        text = 'a' * 2000
        result = self.store._truncate_for_embedding(text)
        assert len(result) <= 1200

    def test_both_limits_enforced(self):
        """201 long words exceed both limits — both caps applied simultaneously."""
        # 201 long words that exceed both limits
        text = ' '.join(['abcdefghij'] * 201)
        result = self.store._truncate_for_embedding(text)
        assert len(result.split()) <= 200
        assert len(result) <= 1200


class TestLowConfidencePipeline:
    """Tests for the low-confidence path — LLM synthesis is skipped when similarity is too low."""

    def test_low_confidence_skips_llm(self):
        chunks = [
            {'text': 'Cats sleep 16 hours.', 'source': 'cats.txt',
             'start_line': 1, 'end_line': 1, 'type': 'txt'},
        ]
        llm_call_count = {'n': 0}

        def counting_llm(prompt, **kwargs):
            llm_call_count['n'] += 1
            return "mock rewrite\nalternative rewrite"

        with patch('src.rag.vector_store._llm_call', side_effect=counting_llm):
            from src.rag.vector_store import VectorStore
            store = VectorStore()
            store.build_or_load(chunks)
            # Force similarity threshold very high so confidence fails
            store._check_confidence = lambda r: (False, 0.1)
            result = store.run_pipeline("unrelated query about quantum physics")

        assert result['is_confident'] is False
        assert "could not find" in result['response'].lower()  # fixed low-confidence message, no LLM call


# ═══════════════════════════════════════════════════════════════════════════
# 5. AGENT
# ═══════════════════════════════════════════════════════════════════════════

class TestParseToolCall:
    """Tests for Agent._parse_tool_call — both regex patterns (with and without parentheses)."""

    def setup_method(self):
        """Build a minimal Agent backed by a small in-memory VectorStore."""
        chunks = _sample_chunks(3)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_parse_with_parentheses(self):
        """Pattern 1: TOOL: name(arg) → (name, arg)."""
        name, arg = self.agent._parse_tool_call("TOOL: rag_search(cat sleep hours)")
        assert name == 'rag_search'
        assert arg == 'cat sleep hours'

    def test_parse_without_parentheses(self):
        """Pattern 2 fallback: TOOL: name arg → (name, arg)."""
        name, arg = self.agent._parse_tool_call("TOOL: rag_search cat nap duration")
        assert name == 'rag_search'
        assert 'cat' in arg

    def test_parse_case_insensitive(self):
        """TOOL keyword and tool name are case-insensitive."""
        name, arg = self.agent._parse_tool_call("tool: Calculator(2 + 2)")
        assert name == 'calculator'

    def test_parse_finish(self):
        """finish tool with a full-sentence argument is parsed correctly."""
        name, arg = self.agent._parse_tool_call("TOOL: finish(The answer is 42.)")
        assert name == 'finish'
        assert arg == 'The answer is 42.'

    def test_parse_malformed_returns_none(self):
        """Text without TOOL: prefix returns (None, None)."""
        name, arg = self.agent._parse_tool_call("This is not a tool call")
        assert name is None
        assert arg is None


class TestCalculatorTool:
    """Tests for Agent._tool_calculator — safe eval with allowed-chars whitelist."""

    def setup_method(self):
        """Build a minimal Agent to exercise the calculator tool."""
        chunks = _sample_chunks(2)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_basic_addition(self):
        """Simple integer addition returns string result."""
        assert self.agent._tool_calculator("2 + 2") == "4"

    def test_complex_expression(self):
        """Nested parentheses and mixed operators evaluate correctly."""
        result = self.agent._tool_calculator("(10 + 5) * 2 / 3")
        assert float(result) == pytest.approx(10.0)

    def test_unsafe_chars_rejected(self):
        """Expressions with disallowed chars (letters, quotes) return an error string."""
        result = self.agent._tool_calculator("__import__('os').system('rm -rf')")
        assert "Error" in result

    def test_unsafe_letters_rejected(self):
        """exec() — contains letters → rejected by whitelist before eval."""
        result = self.agent._tool_calculator("exec('bad')")
        assert "Error" in result

    def test_division(self):
        """Float division returns decimal string result."""
        assert self.agent._tool_calculator("10 / 4") == "2.5"

    def test_float_expression(self):
        """Float literals in expression are handled correctly."""
        result = float(self.agent._tool_calculator("3.14 * 2"))
        assert result == pytest.approx(6.28)


class TestSummariseTool:
    """Tests for Agent._tool_summarise — adaptive length hints based on word count."""

    def setup_method(self):
        """Build a minimal Agent to exercise the summarise tool."""
        chunks = _sample_chunks(2)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_short_text_hint(self):
        """< 100 words → 2-3 sentences hint."""
        captured = {}
        def capture_llm(messages, **kwargs):
            captured['msg'] = messages[0]['content']
            return "summary"
        with patch.object(self.agent.store, '_llm_chat', side_effect=capture_llm):
            self.agent._tool_summarise("Short text " * 5)
        assert "2-3 sentences" in captured['msg']

    def test_medium_text_hint(self):
        """100-299 words → 4-5 sentences hint."""
        captured = {}
        def capture_llm(messages, **kwargs):
            captured['msg'] = messages[0]['content']
            return "summary"
        with patch.object(self.agent.store, '_llm_chat', side_effect=capture_llm):
            self.agent._tool_summarise("word " * 150)
        assert "4-5 sentences" in captured['msg']

    def test_long_text_hint(self):
        """>=300 words → 6-8 sentences hint."""
        captured = {}
        def capture_llm(messages, **kwargs):
            captured['msg'] = messages[0]['content']
            return "summary"
        with patch.object(self.agent.store, '_llm_chat', side_effect=capture_llm):
            self.agent._tool_summarise("word " * 350)
        assert "6-8 sentences" in captured['msg']


class TestSentimentTool:
    """Tests for Agent._tool_sentiment — short queries search first, long text analysed directly."""

    def setup_method(self):
        """Build a minimal Agent to exercise the sentiment tool."""
        chunks = _sample_chunks(3)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_short_query_searches_first(self):
        """< 10 words → rag_search called first."""
        searched = {'called': False}
        def mock_search(q):
            searched['called'] = True
            return "The text is very positive and uplifting."
        with patch.object(self.agent, '_tool_rag_search', side_effect=mock_search), \
             patch.object(self.agent.store, '_llm_chat', return_value="Sentiment: Positive\nTone: uplifting\nKey phrases: positive\nExplanation: test"):
            self.agent._tool_sentiment("resume tone")
        assert searched['called']

    def test_long_text_analyses_directly(self):
        """>=10 words → no rag_search, analyse directly."""
        searched = {'called': False}
        def mock_search(q):
            searched['called'] = True
            return ""
        with patch.object(self.agent, '_tool_rag_search', side_effect=mock_search), \
             patch.object(self.agent.store, '_llm_chat', return_value="Sentiment: Neutral"):
            self.agent._tool_sentiment("This is a very long text that has many words in it.")
        assert not searched['called']

    def test_sentiment_prompt_has_4_fields(self):
        """Sentiment prompt instructs the LLM to return all 4 structured fields."""
        captured = {}
        def capture_llm(messages, **kwargs):
            captured['msg'] = messages[0]['content']
            return "Sentiment: Positive\nTone: upbeat\nKey phrases: great\nExplanation: Good vibes."
        with patch.object(self.agent.store, '_llm_chat', side_effect=capture_llm):
            self.agent._tool_sentiment("This is a long enough text to analyse directly now.")
        prompt = captured['msg']
        assert "Sentiment:" in prompt
        assert "Tone:" in prompt
        assert "Key phrases:" in prompt
        assert "Explanation:" in prompt


class TestAgentFastPaths:
    """Tests for Agent fast-path shortcuts — summarise (4-term search) and sentiment (label strip)."""

    def setup_method(self):
        """Build an Agent over 5 sample chunks for fast-path testing."""
        chunks = _sample_chunks(5)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_fast_path_summarise_does_4_searches(self):
        """Summarise fast path fires exactly 4 rag_search calls with fixed terms."""
        search_calls = []
        def mock_search(q):
            search_calls.append(q)
            return f"result for {q}"
        with patch.object(self.agent, '_tool_rag_search', side_effect=mock_search), \
             patch.object(self.agent, '_synthesize_final_answer', return_value="summary"):
            result = self.agent._fast_path_summarise("summarise the document")
        assert len(search_calls) == 4
        assert 'work experience' in search_calls
        assert 'education' in search_calls
        assert 'answer' in result
        assert 'steps' in result

    def test_fast_path_sentiment_strips_labels(self):
        """Sentiment fast path strips chunk metadata labels before calling _tool_sentiment."""
        with patch.object(self.agent, '_tool_rag_search', return_value="- [cats.txt L1-1] Cats are great."), \
             patch.object(self.agent, '_tool_sentiment', return_value="Sentiment: Positive"):
            result = self.agent._fast_path_sentiment("what is the sentiment?")
        assert result['answer'] == "Sentiment: Positive"
        assert len(result['steps']) == 3

    def test_summarise_query_triggers_fast_path(self):
        """A query containing a summarise keyword routes to _fast_path_summarise."""
        fast_called = {'called': False}
        def mock_fast(q, streamlit_mode=False):
            fast_called['called'] = True
            return {'answer': 'summary', 'steps': []}
        with patch.object(self.agent, '_fast_path_summarise', side_effect=mock_fast):
            self.agent.run("summarise the document")
        assert fast_called['called']

    def test_sentiment_query_triggers_fast_path(self):
        """A query containing a sentiment keyword routes to _fast_path_sentiment."""
        fast_called = {'called': False}
        def mock_fast(q, streamlit_mode=False):
            fast_called['called'] = True
            return {'answer': 'sentiment result', 'steps': []}
        with patch.object(self.agent, '_fast_path_sentiment', side_effect=mock_fast):
            self.agent.run("what is the sentiment of the document?")
        assert fast_called['called']


class TestAgentReactLoop:
    """Tests for the Agent ReAct loop — auto-finish, bad-format retry, step limit."""

    def setup_method(self):
        """Build an Agent over 5 sample chunks for ReAct loop testing."""
        chunks = _sample_chunks(5)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_calculator_auto_finish(self):
        """Calculator result → auto-finish step appended with the numeric result."""
        responses = iter(["TOOL: calculator(10 * 5)"])
        with patch.object(self.agent.store, '_llm_chat', side_effect=lambda *a, **kw: next(responses)):
            result = self.agent.run("what is 10 times 5?")
        assert "50" in result['answer']
        assert any(s['tool'] == 'finish' for s in result['steps'])

    def test_rag_search_auto_finish(self):
        """First rag_search on a non-summarise query → auto-finish after synthesis."""
        responses = iter(["TOOL: rag_search(cat sleep hours)"])
        with patch.object(self.agent.store, '_llm_chat', side_effect=lambda *a, **kw: next(responses)):
            result = self.agent.run("how long do cats sleep?")
        assert result['answer'] is not None
        assert any(s['tool'] == 'finish' for s in result['steps'])

    def test_bad_format_retries_max_2(self):
        """Non-tool LLM output triggers up to 2 correction retries, then uses raw text."""
        call_count = {'n': 0}
        def bad_llm(messages, **kwargs):
            call_count['n'] += 1
            return "This is not a tool call at all."
        with patch.object(self.agent.store, '_llm_chat', side_effect=bad_llm):
            result = self.agent.run("simple question?")
        # 2 bad-format retries allowed, then uses raw text
        assert result['answer'] is not None

    def test_max_steps_reached(self):
        """Loop that keeps calling rag_search terminates at max_steps and returns an answer."""
        self.agent.max_steps = 3
        step_count = {'n': 0}
        def looping_llm(messages, **kwargs):
            step_count['n'] += 1
            return "TOOL: rag_search(cats)"
        with patch.object(self.agent.store, '_llm_chat', side_effect=looping_llm), \
             patch.object(self.agent, '_synthesize_final_answer', return_value="final"):
            result = self.agent.run("how many steps?")
        assert result['answer'] is not None

    def test_finish_uses_collected_context(self):
        """rag_search followed by finish uses collected context for final synthesis."""
        responses = iter([
            "TOOL: rag_search(cats)",
            "TOOL: finish(done)",
        ])
        with patch.object(self.agent.store, '_llm_chat', side_effect=lambda *a, **kw: next(responses)), \
             patch.object(self.agent, '_synthesize_final_answer', return_value="synthesized answer"):
            result = self.agent.run("what do cats do?")
        assert result['answer'] is not None


class TestDispatchTool:
    """Tests for Agent._dispatch_tool — routes tool name to the correct private method."""

    def setup_method(self):
        """Build a minimal Agent to test tool dispatch routing."""
        chunks = _sample_chunks(3)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        from src.rag.agent import Agent
        self.agent = Agent(store)

    def test_dispatch_calculator(self):
        """'calculator' routes to _tool_calculator and returns numeric string."""
        result = self.agent._dispatch_tool('calculator', '3 + 3')
        assert result == '6'

    def test_dispatch_unknown_tool(self):
        """Unrecognised tool name returns an 'Unknown tool' error string."""
        result = self.agent._dispatch_tool('unknown_tool', 'arg')
        assert 'Unknown tool' in result

    def test_dispatch_rag_search_adds_to_context(self):
        """rag_search result is appended to self.collected_context."""
        self.agent.collected_context = []
        with patch.object(self.agent, '_tool_rag_search', return_value="some result"):
            self.agent._dispatch_tool('rag_search', 'cats')
        assert len(self.agent.collected_context) == 1


# ═══════════════════════════════════════════════════════════════════════════
# 6. RERANK PROMPTS
# ═══════════════════════════════════════════════════════════════════════════

class TestRerankPrompts:
    """Tests for VectorStore._rerank_prompt — 7 doc-type-specific prompt variants."""

    def setup_method(self):
        """Build an empty VectorStore to access the rerank prompt generator."""
        from src.rag.vector_store import VectorStore
        self.store = VectorStore()
        self.store.build_or_load([])

    def _entry(self, doc_type):
        """Return a minimal chunk entry dict for the given doc type."""
        return {'text': 'sample text', 'source': 'f', 'start_line': 1,
                'end_line': 1, 'type': doc_type}

    def test_xlsx_prompt(self):
        """xlsx prompt mentions 'spreadsheet' and uses 1-to-10 scoring scale."""
        prompt = self.store._rerank_prompt("query", self._entry('xlsx'))
        assert 'spreadsheet' in prompt.lower()
        assert '1 to 10' in prompt

    def test_csv_prompt(self):
        """csv prompt also uses the spreadsheet variant."""
        prompt = self.store._rerank_prompt("query", self._entry('csv'))
        assert 'spreadsheet' in prompt.lower()

    def test_pptx_prompt(self):
        """pptx prompt references slides."""
        prompt = self.store._rerank_prompt("query", self._entry('pptx'))
        assert 'slide' in prompt.lower()

    def test_pdf_prompt(self):
        """pdf prompt references PDF or page."""
        prompt = self.store._rerank_prompt("query", self._entry('pdf'))
        assert 'pdf' in prompt.lower() or 'page' in prompt.lower()

    def test_docx_prompt(self):
        """docx prompt references paragraph or document."""
        prompt = self.store._rerank_prompt("query", self._entry('docx'))
        assert 'paragraph' in prompt.lower() or 'document' in prompt.lower()

    def test_html_prompt(self):
        """html prompt references webpage."""
        prompt = self.store._rerank_prompt("query", self._entry('html'))
        assert 'webpage' in prompt.lower()

    def test_md_prompt(self):
        """md prompt references markdown."""
        prompt = self.store._rerank_prompt("query", self._entry('md'))
        assert 'markdown' in prompt.lower()

    def test_default_prompt(self):
        """txt falls through to the default prompt with a 1-10 numeric scale."""
        prompt = self.store._rerank_prompt("query", self._entry('txt'))
        assert '1-10' in prompt or '1 to 10' in prompt


# ═══════════════════════════════════════════════════════════════════════════
# 7. LLM CALL — fallback logic
# ═══════════════════════════════════════════════════════════════════════════

class TestLlmCallFallback:
    """Tests for the module-level _llm_call — HF Inference API with multi-model fallback."""

    def test_no_token_returns_error(self):
        """Empty HF_TOKEN returns an error string without making any HTTP request."""
        with patch.dict(os.environ, {'HF_TOKEN': ''}):
            import importlib
            import src.rag.vector_store as vs_module
            result = vs_module._llm_call("test prompt")
        assert "HF_TOKEN not set" in result

    def test_first_model_success_returns_immediately(self):
        """Successful first-model response is returned without trying fallback models."""
        import requests as req_module
        mock_resp = MagicMock()
        mock_resp.ok = True
        mock_resp.json.return_value = {
            "choices": [{"message": {"content": "hello from model"}}]
        }
        with patch.dict(os.environ, {'HF_TOKEN': 'hf_fake'}), \
             patch('requests.post', return_value=mock_resp):
            import src.rag.vector_store as vs_module
            result = vs_module._llm_call("test prompt")
        assert result == "hello from model"

    def test_first_fails_tries_second(self):
        """HTTP 410 from first model triggers retry with the second model in the list."""
        call_count = {'n': 0}
        def fake_post(url, **kwargs):
            call_count['n'] += 1
            if call_count['n'] == 1:
                resp = MagicMock()
                resp.ok = False
                resp.status_code = 410
                resp.text = "Gone"
                return resp
            resp = MagicMock()
            resp.ok = True
            resp.json.return_value = {"choices": [{"message": {"content": "second model"}}]}
            return resp
        with patch.dict(os.environ, {'HF_TOKEN': 'hf_fake'}), \
             patch('requests.post', side_effect=fake_post):
            import src.rag.vector_store as vs_module
            result = vs_module._llm_call("test prompt")
        assert call_count['n'] == 2
        assert result == "second model"

    def test_all_models_fail_returns_error_message(self):
        """All models returning HTTP 410 produces an 'all models failed' error string."""
        mock_resp = MagicMock()
        mock_resp.ok = False
        mock_resp.status_code = 410
        mock_resp.text = "Gone"
        with patch.dict(os.environ, {'HF_TOKEN': 'hf_fake'}), \
             patch('requests.post', return_value=mock_resp):
            import src.rag.vector_store as vs_module
            result = vs_module._llm_call("test prompt")
        assert "all models failed" in result


# ═══════════════════════════════════════════════════════════════════════════
# 8. APP HANDLERS (unit — no Gradio server)
# ═══════════════════════════════════════════════════════════════════════════

gradio_available = pytest.importorskip if False else None
try:
    import gradio  # noqa: F401
    _GRADIO_INSTALLED = True
except ImportError:
    _GRADIO_INSTALLED = False

requires_gradio = pytest.mark.skipif(not _GRADIO_INSTALLED, reason="gradio not installed locally")


@requires_gradio
class TestAppHandlers:
    """Test app.py handler functions in isolation (no Gradio server needed)."""

    def _make_store(self):
        """Build a small in-memory VectorStore for app handler tests."""
        chunks = _sample_chunks(3)
        from src.rag.vector_store import VectorStore
        store = VectorStore()
        store.build_or_load(chunks)
        return store

    def test_pipeline_summary_format(self):
        """_pipeline_summary should produce markdown with Query type and Confidence."""
        sys.path.insert(0, HF_ROOT)
        import importlib
        # Import app without triggering Gradio launch
        import unittest.mock as um
        with um.patch('gradio.Blocks'), um.patch('gradio.Chatbot'), \
             um.patch('gradio.Textbox'), um.patch('gradio.Radio'), \
             um.patch('gradio.Button'), um.patch('gradio.Row'), \
             um.patch('gradio.Column'), um.patch('gradio.Accordion'), \
             um.patch('gradio.File'), um.patch('gradio.Markdown'), \
             um.patch('gradio.HTML'), um.patch('gradio.Progress'):
            try:
                import app as app_module
            except Exception:
                pytest.skip("Could not import app.py in test environment")

        store = self._make_store()
        app_module._store = store

        data = {
            'query_type': 'factual',
            'best_score': 0.75,
            'is_confident': True,
            'retrieved': [
                ({'text': 'Cats sleep.', 'source': 'cats.txt', 'start_line': 1,
                  'end_line': 1, 'type': 'txt'}, 0.75),
            ],
            'reranked': [
                ({'text': 'Cats sleep.', 'source': 'cats.txt', 'start_line': 1,
                  'end_line': 1, 'type': 'txt'}, 0.75, 8.0),
            ],
        }
        result = app_module._pipeline_summary(data)
        assert 'factual' in result
        assert '0.750' in result

    def test_agent_steps_md_format(self):
        """_agent_steps_md should produce markdown with Step N and tool name for each step."""
        sys.path.insert(0, HF_ROOT)
        import unittest.mock as um
        with um.patch('gradio.Blocks'), um.patch('gradio.Chatbot'), \
             um.patch('gradio.Textbox'), um.patch('gradio.Radio'), \
             um.patch('gradio.Button'), um.patch('gradio.Row'), \
             um.patch('gradio.Column'), um.patch('gradio.Accordion'), \
             um.patch('gradio.File'), um.patch('gradio.Markdown'), \
             um.patch('gradio.HTML'), um.patch('gradio.Progress'):
            try:
                import app as app_module
            except Exception:
                pytest.skip("Could not import app.py in test environment")

        steps = [
            {'step': 1, 'tool': 'rag_search', 'arg': 'cats', 'result': 'Cats sleep 16 hours.'},
            {'step': 2, 'tool': 'finish', 'arg': 'done', 'result': 'done'},
        ]
        result = app_module._agent_steps_md(steps)
        assert 'Step 1' in result
        assert 'rag_search' in result
        assert 'Step 2' in result

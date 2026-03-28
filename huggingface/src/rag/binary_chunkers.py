"""binary_chunkers.py — Chunkers for binary document formats.

Handles the five formats that require binary parsing libraries:
PDF (fitz/pymupdf), DOCX (python-docx), XLSX (openpyxl),
XLS legacy (xlrd), and PPTX (python-pptx).

Each function takes a file path and filename, reads the file using the
appropriate library, and returns a list of chunk dicts. All third-party
imports are done inside each function so a missing library only skips
that one format instead of crashing the whole module on import.
"""

import re
from typing import List

from src.rag.config import (
    PDF_CHUNK_SENTENCES,
    DOCX_CHUNK_PARAS,
    PPTX_CHUNK_SLIDES,
)

__all__ = [
    'chunk_pdf',
    'chunk_docx',
    'chunk_xlsx',
    'chunk_xls',
    'chunk_pptx',
]


def chunk_pdf(
    filepath: str,
    filename: str,
    sentences_per_chunk: int = PDF_CHUNK_SENTENCES,
) -> List[dict]:
    """Chunk a PDF using sentence-based windows, one page at a time.

    Page-level isolation means start_line == end_line == page number,
    so source citations are exact (e.g. '[report.pdf p3]').

    Args:
        filepath:            Absolute path to the .pdf file.
        filename:            Display name stored in each chunk's 'source' field.
        sentences_per_chunk: Number of sentences per sliding window.

    Returns:
        List of chunk dicts with type='pdf'. Empty list on parse error.
    """
    try:
        import fitz
    except ImportError:
        print("  [WARNING] pymupdf not installed — skipping PDF. Install: pip install pymupdf")
        return []

    try:
        doc = fitz.open(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    chunks = []
    for page_num, page in enumerate(doc, start=1):
        raw = page.get_text()
        if not raw.strip():
            continue
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', raw) if s.strip()]
        for i in range(0, len(sentences), sentences_per_chunk):
            window = sentences[i: i + sentences_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': page_num,
                'end_line':   page_num,
                'type':       'pdf',
            })
    doc.close()
    return chunks


def chunk_docx(
    filepath: str,
    filename: str,
    paras_per_chunk: int = DOCX_CHUNK_PARAS,
) -> List[dict]:
    """Chunk a DOCX file using paragraph groups plus table rows.

    Table cells are included because resumes often use tables for layout.
    Merged cells (which python-docx repeats) are deduplicated per row.

    Args:
        filepath:        Absolute path to the .docx file.
        filename:        Display name stored in each chunk's 'source' field.
        paras_per_chunk: Number of paragraphs (or table rows) per chunk.

    Returns:
        List of chunk dicts with type='docx'. Empty list on parse error.
    """
    try:
        from docx import Document
    except ImportError:
        print("  [WARNING] python-docx not installed — skipping DOCX. Install: pip install python-docx")
        return []

    try:
        doc = Document(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    # Collect body paragraphs
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Also extract text from tables (resumes often use tables for layout)
    for table in doc.tables:
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            # Deduplicate merged cells (python-docx repeats merged cell text)
            seen = []
            for cell in row_cells:
                if cell not in seen:
                    seen.append(cell)
            if seen:
                paras.append(' | '.join(seen))

    chunks = []
    for i in range(0, len(paras), paras_per_chunk):
        window = paras[i: i + paras_per_chunk]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'docx',
        })
    return chunks


def chunk_xlsx(filepath: str, filename: str) -> List[dict]:
    """Chunk an XLSX file — each non-empty row becomes one chunk.

    Each row is formatted as 'header=value; header=value; ...' pairs
    so the embedding model can understand what each value means.
    Sheet name is prefixed in brackets for multi-sheet workbooks.

    Args:
        filepath: Absolute path to the .xlsx file.
        filename: Display name stored in each chunk's 'source' field.

    Returns:
        List of chunk dicts with type='xlsx'. Empty list on parse error.
    """
    try:
        import openpyxl
    except ImportError:
        print("  [WARNING] openpyxl not installed — skipping XLSX. Install: pip install openpyxl")
        return []

    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    chunks = []
    for sheet in wb.sheetnames:
        ws   = wb[sheet]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else f"col{i}" for i, c in enumerate(rows[0])]
        for row_idx, row in enumerate(rows[1:], start=2):
            cells = [str(v) for v in row if v is not None and str(v).strip()]
            if not cells:
                continue
            pairs = '; '.join(
                f"{header[i] if i < len(header) else f'col{i}'}={str(row[i])}"
                for i in range(len(row)) if row[i] is not None and str(row[i]).strip()
            )
            chunks.append({
                'text':       f"[{sheet}] {pairs}",
                'source':     filename,
                'start_line': row_idx,
                'end_line':   row_idx,
                'type':       'xlsx',
            })
    wb.close()
    return chunks


def chunk_xls(filepath: str, filename: str) -> List[dict]:
    """Chunk a legacy .xls file using xlrd.

    openpyxl cannot read the old .xls binary format, so xlrd is used
    as a fallback. Each data row becomes one key=value pair chunk.

    Args:
        filepath: Absolute path to the .xls file.
        filename: Display name stored in each chunk's 'source' field.

    Returns:
        List of chunk dicts with type='xlsx'. Empty list on parse error.
    """
    try:
        import xlrd
    except ImportError:
        print("  [WARNING] xlrd not installed — skipping XLS. Install: pip install xlrd")
        return []

    try:
        wb = xlrd.open_workbook(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    chunks = []
    for sheet in wb.sheets():
        if sheet.nrows < 2:
            continue
        header = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
        for row_idx in range(1, sheet.nrows):
            pairs = '; '.join(
                f"{header[c] if c < len(header) else f'col{c}'}={sheet.cell_value(row_idx, c)}"
                for c in range(sheet.ncols)
                if str(sheet.cell_value(row_idx, c)).strip()
            )
            if not pairs:
                continue
            chunks.append({
                'text':       f"[{sheet.name}] {pairs}",
                'source':     filename,
                'start_line': row_idx + 1,
                'end_line':   row_idx + 1,
                'type':       'xlsx',
            })
    return chunks


def chunk_pptx(
    filepath: str,
    filename: str,
    slides_per_chunk: int = PPTX_CHUNK_SLIDES,
) -> List[dict]:
    """Chunk a PPTX file — text shapes from each slide, grouped by slide count.

    All text shapes on a slide are joined into one string, then slides
    are grouped into windows of size slides_per_chunk.

    Args:
        filepath:         Absolute path to the .pptx file.
        filename:         Display name stored in each chunk's 'source' field.
        slides_per_chunk: Number of slides per chunk window.

    Returns:
        List of chunk dicts with type='pptx'. Empty list on parse error.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("  [WARNING] python-pptx not installed — skipping PPTX. Install: pip install python-pptx")
        return []

    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    slide_texts = []
    for slide in prs.slides:
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, 'text') and shape.text.strip()
        ]
        if texts:
            slide_texts.append(' '.join(texts))

    chunks = []
    for i in range(0, len(slide_texts), slides_per_chunk):
        window = slide_texts[i: i + slides_per_chunk]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'pptx',
        })
    return chunks

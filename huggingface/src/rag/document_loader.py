"""document_loader.py — DocumentLoader class: owns all ingestion. HF Space version."""

import csv as _csv
import os
import re
import tempfile
from typing import List
from urllib.parse import urlparse

from src.rag.config import (
    EXT_TO_TYPE,
    TXT_CHUNK_SIZE, TXT_CHUNK_OVERLAP,
    PDF_CHUNK_SENTENCES, DOCX_CHUNK_PARAS,
    PPTX_CHUNK_SLIDES, HTML_CHUNK_SENTENCES,
)

__all__ = ['DocumentLoader']


class DocumentLoader:
    """Owns all document ingestion — chunkers, URL fetching.

    State:
        ext_to_type (dict): Maps file extension strings to canonical type keys.
        chunk_sizes (dict): All chunk-size constants keyed by config name.

    Public API:
        chunk_url(url)          -- Fetch a URL and route to the correct chunker.
        _dispatch_chunker(info) -- Pick and call the right private chunker.
    """

    def __init__(self) -> None:
        """Load config constants as instance vars so they're easy to override in tests."""
        self.ext_to_type = EXT_TO_TYPE
        self.chunk_sizes = {
            'txt_chunk_size':        TXT_CHUNK_SIZE,
            'txt_chunk_overlap':     TXT_CHUNK_OVERLAP,
            'pdf_chunk_sentences':   PDF_CHUNK_SENTENCES,
            'docx_chunk_paras':      DOCX_CHUNK_PARAS,
            'pptx_chunk_slides':     PPTX_CHUNK_SLIDES,
            'html_chunk_sentences':  HTML_CHUNK_SENTENCES,
        }

    # ------------------------------------------------------------------ Public

    def chunk_url(self, url: str) -> List[dict]:
        """Fetch a URL and produce chunks using the appropriate format handler.

        Type detection uses four priorities in strict order:
            1. Content-Type response header (most reliable).
            2. File extension in the URL path (strips query strings first).
            3. PDF magic bytes sniff (``content[:4] == b'%PDF'``).
            4. Default to 'html' when no other signal is present.

        Args:
            url: Public HTTP/HTTPS URL to fetch.

        Returns:
            List of chunk dicts, each with keys: text, source, start_line,
            end_line, type.  Returns an empty list on network error.
        """
        try:
            import requests
        except ImportError:
            print("  [WARNING] requests not installed.")
            return []

        url = url.strip()
        try:
            resp = requests.get(url, timeout=60, headers={
                'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) '
                              'AppleWebKit/537.36 (KHTML, like Gecko) '
                              'Chrome/120.0.0.0 Safari/537.36'
            }, allow_redirects=True, stream=True)
            resp.raise_for_status()
            content = resp.content
            text    = None
        except Exception as e:
            print(f"  [ERROR] Could not fetch URL: {e}")
            return []

        # ── Priority 1: exact Content-Type header match ──────────────
        # Strip charset and boundary params before comparing
        content_type = resp.headers.get('Content-Type', '').lower().split(';')[0].strip()
        dtype = None

        CONTENT_TYPE_MAP = {
            'application/pdf':                                          'pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
            'application/msword':                                       'docx',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
            'application/vnd.ms-excel':                                 'xlsx',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
            'application/vnd.ms-powerpoint':                            'pptx',
            'text/csv':                                                 'csv',
            'text/plain':                                               'txt',
            'text/markdown':                                            'md',
            'text/html':                                                'html',
            'application/xhtml+xml':                                    'html',
        }
        dtype = CONTENT_TYPE_MAP.get(content_type)

        # Fuzzy fallback for unusual MIME subtypes (e.g. "application/x-pdf")
        if dtype is None:
            if 'pdf'          in content_type: dtype = 'pdf'
            elif 'word'       in content_type: dtype = 'docx'
            elif 'excel'      in content_type or 'spreadsheet' in content_type: dtype = 'xlsx'
            elif 'powerpoint' in content_type or 'presentation' in content_type: dtype = 'pptx'
            elif 'csv'        in content_type: dtype = 'csv'

        # ── Priority 2: file extension in URL path ───────────────────
        # Strip query and fragment before extracting extension
        if dtype is None:
            parsed_path = urlparse(url).path.lower()
            clean_path  = re.sub(r'[?#].*$', '', parsed_path)
            ext         = os.path.splitext(clean_path)[1]
            dtype       = self.ext_to_type.get(ext)

        # ── Priority 3: PDF magic bytes ───────────────────────────────
        # Some servers send application/octet-stream for PDFs
        if dtype is None and content[:4] == b'%PDF':
            dtype = 'pdf'

        # ── Priority 4: default ───────────────────────────────────────
        if dtype is None:
            dtype = 'html'

        # ── Source label ─────────────────────────────────────────────
        parsed      = urlparse(url)
        source_name = (parsed.netloc + parsed.path).rstrip('/')
        if parsed.query:
            source_name += '?' + parsed.query[:20] + ('...' if len(parsed.query) > 20 else '')
        if len(source_name) > 60:
            source_name = source_name[:57] + '...'
        if not source_name:
            source_name = url[:60]

        # ── Dispatch to chunker ───────────────────────────────────────
        if dtype in ('pdf', 'docx', 'xlsx', 'pptx', 'xls'):
            suffix = {'pdf':'.pdf','docx':'.docx','xlsx':'.xlsx','pptx':'.pptx','xls':'.xls'}[dtype]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                file_info = {
                    'filepath':      tmp_path,
                    'filename':      source_name,
                    'detected_type': dtype,
                    'is_misplaced':  False,
                }
                chunks = self._dispatch_chunker(file_info)
            finally:
                try: os.unlink(tmp_path)
                except: pass
            return chunks

        try:
            text = content.decode(resp.encoding or 'utf-8', errors='replace')
        except Exception:
            text = content.decode('utf-8', errors='replace')

        if dtype == 'csv':
            with tempfile.NamedTemporaryFile(delete=False, suffix='.csv', mode='w',
                                             encoding='utf-8', errors='replace') as tmp:
                tmp.write(text)
                tmp_path = tmp.name
            try:
                chunks = self._chunk_csv(tmp_path, source_name)
            finally:
                try: os.unlink(tmp_path)
                except: pass
            return chunks

        if dtype == 'txt':
            lines = [l.strip() for l in text.splitlines() if l.strip()]
            return [{'text': l, 'source': source_name, 'start_line': i+1,
                     'end_line': i+1, 'type': 'txt'} for i, l in enumerate(lines)]

        if dtype == 'md':
            with tempfile.NamedTemporaryFile(delete=False, suffix='.md', mode='w',
                                              encoding='utf-8', errors='replace') as tmp:
                tmp.write(text)
                tmp_path = tmp.name
            try:
                chunks = self._chunk_md(tmp_path, source_name)
            finally:
                try: os.unlink(tmp_path)
                except: pass
            return chunks

        # HTML / webpage
        try:
            from bs4 import BeautifulSoup
            text = BeautifulSoup(text, 'html.parser').get_text(separator=' ', strip=True)
        except ImportError:
            text = re.sub(r'<[^>]+>', ' ', text)

        sents  = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
        chunks = []
        for i in range(0, len(sents), HTML_CHUNK_SENTENCES):
            window = sents[i: i + HTML_CHUNK_SENTENCES]
            if window:
                chunks.append({'text': ' '.join(window), 'source': source_name,
                               'start_line': i+1, 'end_line': i+len(window), 'type': 'html'})
        return chunks

    # ----------------------------------------------------------------- Private

    def _dispatch_chunker(self, file_info: dict) -> List[dict]:
        """Route a file_info dict to the correct private chunker method.

        Args:
            file_info: Dict with keys filepath, filename, detected_type, is_misplaced.

        Returns:
            List of chunk dicts.  Empty list when the type is unrecognised.
        """
        fp  = file_info['filepath']
        fn  = file_info['filename']
        ext = os.path.splitext(fn)[1].lower()
        t   = file_info['detected_type']

        if t == 'txt':
            return self._chunk_txt(fp, fn)
        elif t == 'md':
            return self._chunk_md(fp, fn)
        elif t == 'pdf':
            return self._chunk_pdf(fp, fn)
        elif t == 'docx':
            return self._chunk_docx(fp, fn)
        elif t == 'xlsx':
            # .xls uses xlrd (legacy binary format); .csv despite 'xlsx' type
            # only arises from misplaced-file detection — route correctly
            if ext == '.xls':
                return self._chunk_xls(fp, fn)
            elif ext == '.csv':
                return self._chunk_csv(fp, fn)
            else:
                return self._chunk_xlsx(fp, fn)
        elif t == 'csv':
            return self._chunk_csv(fp, fn)
        elif t == 'pptx':
            return self._chunk_pptx(fp, fn)
        elif t == 'html':
            return self._chunk_html(fp, fn)
        else:
            print(f"  [SKIP] No chunker for type '{t}' ({fn})")
            return []

    def _chunk_txt(self, filepath: str, filename: str) -> List[dict]:
        """Chunk a plain-text file — one logical line per chunk, no overlap by default."""
        chunk_size = self.chunk_sizes['txt_chunk_size']
        overlap    = self.chunk_sizes['txt_chunk_overlap']
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            lines = [l.strip() for l in f.readlines() if l.strip()]
        step   = max(1, chunk_size - overlap)
        chunks = []
        for i in range(0, len(lines), step):
            window = lines[i: i + chunk_size]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'txt',
            })
        return chunks

    def _chunk_md(self, filepath: str, filename: str) -> List[dict]:
        """Strips markdown syntax and chunks like plain text."""
        chunk_size = self.chunk_sizes['txt_chunk_size']
        overlap    = self.chunk_sizes['txt_chunk_overlap']
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()
        clean = re.sub(r'#{1,6}\s*', '', raw)
        clean = re.sub(r'[*_`~]{1,3}', '', clean)
        clean = re.sub(r'!\[.*?\]\(.*?\)', '', clean)
        clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)
        lines = [l.strip() for l in clean.splitlines() if l.strip()]
        step  = max(1, chunk_size - overlap)
        chunks = []
        for i in range(0, len(lines), step):
            window = lines[i: i + chunk_size]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'md',
            })
        return chunks

    def _chunk_pdf(self, filepath: str, filename: str) -> List[dict]:
        """Chunk a PDF file — sliding sentence windows per page, page number as start_line."""
        sentences_per_chunk = self.chunk_sizes['pdf_chunk_sentences']
        try:
            import fitz
        except ImportError:
            print("  [WARNING] pymupdf not installed — skipping PDF.")
            return []

        chunks   = []
        try:
            doc = fitz.open(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        for page_num, page in enumerate(doc, start=1):
            raw = page.get_text()
            if not raw.strip():
                continue
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', raw) if s.strip()]
            for i in range(0, len(sentences), sentences_per_chunk):
                window = sentences[i: i + sentences_per_chunk]
                if not window:
                    continue
                chunks.append({
                    'text':       ' '.join(window),
                    'source':     filename,
                    'start_line': page_num,
                    'end_line':   page_num,
                    'type':       'pdf',
                })
        doc.close()
        return chunks

    def _chunk_docx(self, filepath: str, filename: str) -> List[dict]:
        """Chunk a DOCX file — paragraphs plus table rows; merged cells are deduplicated."""
        paras_per_chunk = self.chunk_sizes['docx_chunk_paras']
        try:
            from docx import Document
        except ImportError:
            print("  [WARNING] python-docx not installed — skipping DOCX.")
            return []

        chunks = []
        try:
            doc = Document(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        for table in doc.tables:
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                # python-docx repeats the same cell object for merged cells;
                # deduplicate by preserving insertion order
                seen = []
                for cell in row_cells:
                    if cell not in seen:
                        seen.append(cell)
                if seen:
                    paras.append(' | '.join(seen))

        for i in range(0, len(paras), paras_per_chunk):
            window = paras[i: i + paras_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'docx',
            })
        return chunks

    def _chunk_xlsx(self, filepath: str, filename: str) -> List[dict]:
        """Each non-empty row becomes a chunk (joined as key=value pairs)."""
        try:
            import openpyxl
        except ImportError:
            print("  [WARNING] openpyxl not installed — skipping XLSX.")
            return []

        chunks = []
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        for sheet in wb.sheetnames:
            ws     = wb[sheet]
            rows   = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            header = [str(c) if c is not None else f"col{i}" for i, c in enumerate(rows[0])]
            for row_idx, row in enumerate(rows[1:], start=2):
                cells = [str(v) for v in row if v is not None and str(v).strip()]
                if not cells:
                    continue
                pairs = '; '.join(
                    f"{header[i] if i < len(header) else f'col{i}'}={str(row[i])}"
                    for i in range(len(row)) if row[i] is not None and str(row[i]).strip()
                )
                chunks.append({
                    'text':       f"[{sheet}] {pairs}",
                    'source':     filename,
                    'start_line': row_idx,
                    'end_line':   row_idx,
                    'type':       'xlsx',
                })
        wb.close()
        return chunks

    def _chunk_xls(self, filepath: str, filename: str) -> List[dict]:
        """Fallback chunker for legacy .xls files using xlrd (openpyxl cannot read them)."""
        try:
            import xlrd
        except ImportError:
            print("  [WARNING] xlrd not installed — skipping XLS.")
            return []

        chunks = []
        try:
            wb = xlrd.open_workbook(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        for sheet in wb.sheets():
            if sheet.nrows < 2:
                continue
            header = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
            for row_idx in range(1, sheet.nrows):
                pairs = '; '.join(
                    f"{header[c] if c < len(header) else f'col{c}'}={sheet.cell_value(row_idx, c)}"
                    for c in range(sheet.ncols)
                    if str(sheet.cell_value(row_idx, c)).strip()
                )
                if not pairs:
                    continue
                chunks.append({
                    'text':       f"[{sheet.name}] {pairs}",
                    'source':     filename,
                    'start_line': row_idx + 1,
                    'end_line':   row_idx + 1,
                    'type':       'xlsx',
                })
        return chunks

    def _chunk_csv(self, filepath: str, filename: str) -> List[dict]:
        """Each CSV row becomes a chunk."""
        chunks = []
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                reader = _csv.DictReader(f)
                for row_idx, row in enumerate(reader, start=2):
                    pairs = '; '.join(f"{k}={v}" for k, v in row.items() if v and str(v).strip())
                    if not pairs:
                        continue
                    chunks.append({
                        'text':       pairs,
                        'source':     filename,
                        'start_line': row_idx,
                        'end_line':   row_idx,
                        'type':       'csv',
                    })
        except Exception as e:
            print(f"  [ERROR] Could not read '{filename}': {e}")
        return chunks

    def _chunk_pptx(self, filepath: str, filename: str) -> List[dict]:
        """Chunk a PPTX file — text shapes from each slide, one chunk per slide."""
        slides_per_chunk = self.chunk_sizes['pptx_chunk_slides']
        try:
            from pptx import Presentation
        except ImportError:
            print("  [WARNING] python-pptx not installed — skipping PPTX.")
            return []

        chunks = []
        try:
            prs = Presentation(filepath)
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        slide_texts = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(' '.join(texts))

        for i in range(0, len(slide_texts), slides_per_chunk):
            window = slide_texts[i: i + slides_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'pptx',
            })
        return chunks

    def _chunk_html(self, filepath: str, filename: str) -> List[dict]:
        """Chunk an HTML file — strip navigation/boilerplate tags, then sentence windows."""
        sentences_per_chunk = self.chunk_sizes['html_chunk_sentences']
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            print("  [WARNING] beautifulsoup4 not installed — skipping HTML.")
            return []

        chunks = []
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
        except Exception as e:
            print(f"  [ERROR] Could not open '{filename}': {e}")
            return []

        # Remove navigation boilerplate before extracting text
        for tag in soup.find_all(['nav', 'header', 'footer', 'aside',
                                  'script', 'style', 'noscript']):
            tag.decompose()
        for tag in soup.find_all(attrs={'role': ['navigation', 'banner', 'complementary']}):
            tag.decompose()
        for tag in soup.find_all(attrs={'id': ['mw-navigation', 'mw-head', 'mw-panel',
                                               'p-navigation', 'p-tb', 'footer',
                                               'toc', 'catlinks', 'mw-page-base']}):
            tag.decompose()
        for tag in soup.find_all(attrs={'class': ['navbox', 'sidebar', 'mw-editsection',
                                                   'reflist', 'reference', 'noprint']}):
            tag.decompose()
        text = soup.get_text(separator=' ', strip=True)
        # 40-char threshold filters single-word nav items and short menu labels
        # while keeping substantive content lines
        lines = [ln.strip() for ln in text.splitlines() if len(ln.strip()) > 40]
        text  = ' '.join(lines)
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
        for i in range(0, len(sentences), sentences_per_chunk):
            window = sentences[i: i + sentences_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': i + 1,
                'end_line':   i + len(window),
                'type':       'html',
            })
        return chunks

    def _truncate_chunk(self, text: str, max_words: int = 300, max_chars: int = 1200) -> str:
        """Truncate to 300 words OR 1200 chars, whichever is shorter."""
        words     = text.split()
        truncated = ' '.join(words[:max_words]) if len(words) > max_words else text
        return truncated[:max_chars] if len(truncated) > max_chars else truncated

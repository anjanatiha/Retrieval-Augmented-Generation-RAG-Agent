"""test_url_ingestion.py — URL ingestion tests for all 8 file types in chat and
agent mode.

Mock strategy:
  Always mock: ollama.embed, ollama.chat, chromadb → EphemeralClient, requests.get
  Never mock:  fitz, python-docx, openpyxl, python-pptx, beautifulsoup4, BM25Okapi

Reason for split: max 500 lines per file per CLAUDE.md.
Pipeline feature tests (query type, confidence, reranked, source labels,
source cited for URL-loaded docs) are in test_url_pipeline.py.
"""

import io
import os
import uuid
from unittest.mock import MagicMock, patch

import chromadb
import pytest
from rank_bm25 import BM25Okapi

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _fake_embed(**kwargs):
    """Return a fixed 4-dim embedding vector to satisfy ollama.embed calls."""
    return {'embeddings': [[0.1, 0.2, 0.3, 0.4]]}


def _fake_chat(**kwargs):
    """Return a canned chat response (stream or single) to satisfy ollama.chat calls."""
    if kwargs.get('stream'):
        return [{'message': {'content': 'mock response'}}]
    return {'message': {'content': 'mock response'}}


def _make_store(chunks):
    """Build an isolated VectorStore backed by an EphemeralClient for the given chunks."""
    from src.rag.vector_store import VectorStore
    vs         = VectorStore()
    client     = chromadb.EphemeralClient()
    collection = client.get_or_create_collection(
        f'test_{uuid.uuid4().hex}', metadata={'hnsw:space': 'cosine'}
    )
    if chunks:
        ids    = [f'c{i}' for i in range(len(chunks))]
        texts  = [c['text'] for c in chunks]
        metas  = [{'source': c['source'], 'start_line': c['start_line'],
                   'end_line': c['end_line'], 'type': c['type']} for c in chunks]
        embeds = [[0.1, 0.2, 0.3, 0.4]] * len(chunks)
        collection.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
    vs.collection           = collection
    vs.chunks               = chunks
    vs.bm25_index           = BM25Okapi([c['text'].lower().split() for c in chunks]) if chunks else None
    vs.conversation_history = []
    return vs


PATCHES = [
    patch('ollama.embed', side_effect=_fake_embed),
    patch('ollama.chat',  side_effect=_fake_chat),
]


def _apply_patches(func):
    """Decorator that applies both ollama.embed and ollama.chat patches to a test method."""
    for p in reversed(PATCHES):
        func = p(func)
    return func


def _mock_url_response(content, content_type='text/html', encoding='utf-8'):
    """Build a MagicMock simulating a requests.Response with the given content and headers."""
    resp = MagicMock()
    resp.content  = content
    resp.headers  = {'Content-Type': content_type}
    resp.encoding = encoding
    resp.raise_for_status = MagicMock()
    return resp


# ---------------------------------------------------------------------------
# 1. URL ingestion — chat mode (all 8 file types)
# ---------------------------------------------------------------------------

class TestUrlIngestionChatMode:
    """DocumentLoader.chunk_url → VectorStore.run_pipeline in chat mode for all 8 URL types."""

    @_apply_patches
    def test_url_html_chat_mode(self, mock_chat, mock_embed):
        """HTML URL → chunk_url produces chunks; run_pipeline returns a response key."""
        html = (b'<html><body><p>Dogs have an extraordinary sense of smell. '
                b'They can detect scents from miles away. '
                b'Their noses have 300 million receptors.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        assert len(chunks) >= 1
        result = _make_store(chunks).run_pipeline('What is special about dog smell?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_url_pdf_chat_mode(self, mock_chat, mock_embed):
        """PDF URL (application/pdf) → chunk_url; run_pipeline returns a response key."""
        import fitz
        doc  = fitz.open(); page = doc.new_page()
        page.insert_text((50, 50),
            'Nutrition facts: Bananas contain 89 calories per 100 grams and are rich in potassium.')
        buf = io.BytesIO(); doc.save(buf); pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_url_response(pdf_bytes, 'application/pdf')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/nutrition.pdf')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('How many calories in bananas?', streamlit_mode=True)
            assert 'response' in result

    @_apply_patches
    def test_url_docx_chat_mode(self, mock_chat, mock_embed):
        """DOCX URL (wordprocessingml content-type) → chunk_url; run_pipeline returns a response key."""
        from docx import Document
        doc = Document()
        doc.add_paragraph('Python is a programming language created by Guido van Rossum in 1991.')
        doc.add_paragraph('It emphasizes code readability and simplicity.')
        buf = io.BytesIO(); doc.save(buf); docx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        with patch('requests.get', return_value=_mock_url_response(docx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/languages.docx')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
            assert 'response' in result

    @_apply_patches
    def test_url_xlsx_chat_mode(self, mock_chat, mock_embed):
        """XLSX URL (spreadsheetml content-type) → chunk_url; run_pipeline returns a response key."""
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Country', 'Capital', 'Population'])
        ws.append(['Brazil', 'Brasilia', '215 million'])
        ws.append(['Japan',  'Tokyo',    '125 million'])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        with patch('requests.get', return_value=_mock_url_response(xlsx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/countries.xlsx')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('What is the capital of Brazil?', streamlit_mode=True)
            assert 'response' in result

    @_apply_patches
    def test_url_csv_chat_mode(self, mock_chat, mock_embed):
        """CSV URL (text/csv content-type) → chunk_url; run_pipeline returns a response key."""
        csv_content = b'language,year,creator\nPython,1991,Guido van Rossum\nRust,2010,Graydon Hoare\n'
        with patch('requests.get', return_value=_mock_url_response(csv_content, 'text/csv')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/languages.csv')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
            assert 'response' in result

    @_apply_patches
    def test_url_pptx_chat_mode(self, mock_chat, mock_embed):
        """PPTX URL (presentationml content-type) → chunk_url; run_pipeline returns a response key."""
        from pptx import Presentation
        from pptx.util import Inches
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'The Berlin Wall fell on November 9 1989 ending the Cold War division.'
        buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        with patch('requests.get', return_value=_mock_url_response(pptx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/history.pptx')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('When did the Berlin Wall fall?', streamlit_mode=True)
            assert 'response' in result

    @_apply_patches
    def test_url_md_chat_mode(self, mock_chat, mock_embed):
        """Markdown URL (text/markdown content-type) → chunk_url; run_pipeline returns a response key."""
        md_content = (b'# Coffee Facts\n\n'
                      b'Coffee was first discovered in Ethiopia around 850 AD.\n\n'
                      b"Brazil is the world's largest coffee producer.")
        with patch('requests.get', return_value=_mock_url_response(md_content, 'text/markdown')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/coffee.md')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('Where was coffee discovered?', streamlit_mode=True)
            assert 'response' in result

    @_apply_patches
    def test_url_txt_chat_mode(self, mock_chat, mock_embed):
        """Plain-text URL (text/plain content-type) → chunk_url; run_pipeline returns a response key."""
        txt_content = (b'Dogs were domesticated from wolves approximately 15000 years ago.\n'
                       b'There are over 340 recognized dog breeds in the world.')
        with patch('requests.get', return_value=_mock_url_response(txt_content, 'text/plain')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs.txt')
        assert isinstance(chunks, list)
        if chunks:
            result = _make_store(chunks).run_pipeline('How were dogs domesticated?', streamlit_mode=True)
            assert 'response' in result


# ---------------------------------------------------------------------------
# 2. URL ingestion — agent mode (all 8 file types)
# ---------------------------------------------------------------------------

class TestUrlIngestionAgentMode:
    """DocumentLoader.chunk_url → Agent.run in agent mode for all 8 URL types."""

    @_apply_patches
    def test_url_html_agent_mode(self, mock_chat, mock_embed):
        """HTML URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(dog smell sense)'}}
        html = (b'<html><body><p>Dogs have an extraordinary sense of smell. '
                b'They can detect scents from miles away.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        assert len(chunks) >= 1
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('What is special about dog smell?', streamlit_mode=True)
        assert 'answer' in result
        assert any(s['tool'] == 'rag_search' for s in result['steps'])

    @_apply_patches
    def test_url_pdf_agent_mode(self, mock_chat, mock_embed):
        """PDF URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(banana calories)'}}
        import fitz
        doc  = fitz.open(); page = doc.new_page()
        page.insert_text((50, 50), 'Bananas contain 89 calories per 100 grams.')
        buf = io.BytesIO(); doc.save(buf); pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_url_response(pdf_bytes, 'application/pdf')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/nutrition.pdf')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('How many calories in bananas?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_docx_agent_mode(self, mock_chat, mock_embed):
        """DOCX URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Python creator)'}}
        from docx import Document
        doc = Document()
        doc.add_paragraph('Python is a programming language created by Guido van Rossum in 1991.')
        buf = io.BytesIO(); doc.save(buf); docx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        with patch('requests.get', return_value=_mock_url_response(docx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/languages.docx')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('Who created Python?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_xlsx_agent_mode(self, mock_chat, mock_embed):
        """XLSX URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Brazil capital)'}}
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Country', 'Capital']); ws.append(['Brazil', 'Brasilia'])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        with patch('requests.get', return_value=_mock_url_response(xlsx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/countries.xlsx')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('What is the capital of Brazil?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_csv_agent_mode(self, mock_chat, mock_embed):
        """CSV URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Python creator)'}}
        csv_content = b'language,year,creator\nPython,1991,Guido van Rossum\nRust,2010,Graydon Hoare\n'
        with patch('requests.get', return_value=_mock_url_response(csv_content, 'text/csv')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/langs.csv')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('Who created Python?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_pptx_agent_mode(self, mock_chat, mock_embed):
        """PPTX URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Berlin Wall)'}}
        from pptx import Presentation
        from pptx.util import Inches
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'The Berlin Wall fell on November 9 1989.'
        buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        with patch('requests.get', return_value=_mock_url_response(pptx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/history.pptx')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('When did the Berlin Wall fall?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_md_agent_mode(self, mock_chat, mock_embed):
        """Markdown URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(coffee Ethiopia)'}}
        md_content = b'# Coffee Facts\n\nCoffee was first discovered in Ethiopia around 850 AD.'
        with patch('requests.get', return_value=_mock_url_response(md_content, 'text/markdown')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/coffee.md')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('Where was coffee discovered?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_txt_agent_mode(self, mock_chat, mock_embed):
        """Plain-text URL in agent mode: rag_search tool is used and answer key is present."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(dog breeds)'}}
        txt_content = (b'Dogs were domesticated from wolves approximately 15000 years ago.\n'
                       b'There are over 340 recognized dog breeds in the world.')
        with patch('requests.get', return_value=_mock_url_response(txt_content, 'text/plain')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs.txt')
        if chunks:
            from src.rag.agent import Agent
            result = Agent(_make_store(chunks)).run('How many dog breeds are there?', streamlit_mode=True)
            assert 'answer' in result

    @_apply_patches
    def test_url_chunks_indexed_and_searchable(self, mock_chat, mock_embed):
        """Chunks from URL should be retrievable via rag_search."""
        html = (b'<html><body><p>Finland consumes the most coffee per capita in the world '
                b'at 12 kg per person per year.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/coffee-facts')
        assert len(chunks) >= 1
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_rag_search('Finland coffee consumption')
        assert len(result) > 0

    @_apply_patches
    def test_url_source_stored_in_chunks(self, mock_chat, mock_embed):
        """URL hostname should be stored as source in each chunk."""
        html = (b'<html><body><p>Some interesting content about dogs and their behavior '
                b'that is long enough to pass the filter.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/test-page')
        if chunks:
            assert all('example.com' in c['source'] for c in chunks)

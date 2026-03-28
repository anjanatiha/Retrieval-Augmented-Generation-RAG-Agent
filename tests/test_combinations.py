"""test_combinations.py — Combination matrix tests: doc types × pipeline modes + URL types.

Tests every (doc_type × pipeline_mode) and (URL content_type × url_type) combination
to ensure coverage across the full input space rather than just the happy path.

Scope (this file): chat/agent mode × doc types, URL fetch × content types, chunker contract.
Truncation, classification, and source-label combinations are in test_combinations_analysis.py.

Mock strategy (per CLAUDE.md):
  Always mock: ollama.embed, ollama.chat, requests.get
  Never mock:  fitz, python-docx, openpyxl, python-pptx, beautifulsoup4, BM25Okapi

Uses real temp files for binary formats (pdf, docx, xlsx, pptx) so the parsers
run against genuine file bytes — not synthetic in-memory mocks.
"""

import io
import csv
import os
import tempfile
from unittest.mock import MagicMock, patch

import pytest
import chromadb
from rank_bm25 import BM25Okapi

from src.rag.chunkers import (
    chunk_txt, chunk_md, chunk_csv, chunk_html,
)
from src.rag.binary_chunkers import (
    chunk_pdf, chunk_docx, chunk_xlsx, chunk_pptx,
)

# Map doc type string → chunker function (used instead of getattr on loader)
_CHUNKER_MAP = {
    'txt':  chunk_txt,
    'md':   chunk_md,
    'csv':  chunk_csv,
    'html': chunk_html,
    'pdf':  chunk_pdf,
    'docx': chunk_docx,
    'xlsx': chunk_xlsx,
    'pptx': chunk_pptx,
}


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _make_store(chunks):
    """Return a VectorStore wired to an in-memory EphemeralClient.

    Args:
        chunks: List of chunk dicts to seed into the store.

    Returns:
        Fully-wired VectorStore with BM25 and ChromaDB populated.
    """
    from src.rag.vector_store import VectorStore
    vs     = VectorStore()
    client = chromadb.EphemeralClient()
    col    = client.get_or_create_collection(
        'comb_test', metadata={'hnsw:space': 'cosine'}
    )
    if chunks:
        ids    = [f'c{i}' for i in range(len(chunks))]
        texts  = [c['text'] for c in chunks]
        metas  = [
            {'source': c['source'], 'start_line': c['start_line'],
             'end_line': c['end_line'], 'type': c['type']}
            for c in chunks
        ]
        embeds = [[0.1, 0.2, 0.3, 0.4]] * len(chunks)
        col.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
    vs.collection           = col
    vs.chunks               = list(chunks)
    vs.bm25_index           = BM25Okapi([c['text'].lower().split() for c in chunks]) if chunks else None
    vs.conversation_history = []
    return vs


def _chat_mock(*args, **kwargs):
    """Smart mock: returns stream list when stream=True, dict otherwise."""
    if kwargs.get('stream'):
        return [{'message': {'content': 'Cats sleep about 16 hours a day.'}}]
    return {'message': {'content': '8'}}


# ---------------------------------------------------------------------------
# Binary file fixture helpers (real parsers)
# ---------------------------------------------------------------------------

def _make_pdf_file(tmp_path, filename='test.pdf'):
    """Write a real single-page PDF to tmp_path and return the file path."""
    import fitz
    doc  = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Cats sleep sixteen hours a day. They are nocturnal hunters.")
    path = str(tmp_path / filename)
    doc.save(path); doc.close()
    return path


def _make_docx_file(tmp_path, filename='test.docx'):
    """Write a real DOCX file to tmp_path and return the file path."""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Cats sleep sixteen hours a day.")
    doc.add_paragraph("They are excellent nocturnal hunters.")
    path = str(tmp_path / filename)
    doc.save(path)
    return path


def _make_xlsx_file(tmp_path, filename='test.xlsx'):
    """Write a real XLSX file to tmp_path and return the file path."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(['animal', 'fact'])
    ws.append(['cat', 'sleeps 16 hours a day'])
    path = str(tmp_path / filename)
    wb.save(path)
    return path


def _make_pptx_file(tmp_path, filename='test.pptx'):
    """Write a real PPTX file to tmp_path and return the file path."""
    from pptx import Presentation
    prs   = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text               = "Cat Facts"
    slide.placeholders[1].text_frame.text = "Cats sleep 16 hours a day."
    path  = str(tmp_path / filename)
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# 1. Chat mode × all 8 doc types
# ---------------------------------------------------------------------------

class TestChatModeAllDocTypes:
    """run_pipeline returns a non-empty response for every supported doc type.

    Each parametrize case creates a real file, calls the matching chunker,
    seeds a VectorStore, and runs the pipeline with ollama fully mocked.
    """

    @pytest.mark.parametrize("doc_type,text_content,filename", [
        (
            "txt",
            "Cats sleep 16 hours a day.\nDogs are loyal companions.",
            "test.txt",
        ),
        (
            "md",
            "# Animals\n\nCats sleep 16 hours.\nDogs are loyal companions.",
            "test.md",
        ),
        (
            "csv",
            "animal,fact\ncat,sleeps 16 hours\ndog,loyal companion",
            "test.csv",
        ),
        (
            "html",
            "<html><body><p>Cats sleep 16 hours a day. Dogs are loyal companions.</p></body></html>",
            "test.html",
        ),
    ])
    def test_text_doc_type_pipeline(self, doc_type, text_content, filename, tmp_path):
        """Pipeline returns a non-empty string response for text-based doc type.

        Args:
            doc_type:     One of txt, md, csv, html.
            text_content: Raw string written to the temp file.
            filename:     Filename passed to the chunker.
        """
        f = tmp_path / filename
        f.write_text(text_content)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(str(f), filename)
        assert len(chunks) >= 1, f"No chunks produced for {doc_type}"

        store  = _make_store(chunks)
        with patch('ollama.embed', return_value={'embeddings': [[0.1, 0.2, 0.3, 0.4]]}), \
             patch('ollama.chat', side_effect=_chat_mock):
            result = store.run_pipeline('How long do cats sleep?', streamlit_mode=True)

        assert isinstance(result['response'], str)
        assert len(result['response']) > 0, f"Empty response for doc_type={doc_type}"

    @pytest.mark.parametrize("doc_type,make_fn", [
        ("pdf",  _make_pdf_file),
        ("docx", _make_docx_file),
        ("xlsx", _make_xlsx_file),
        ("pptx", _make_pptx_file),
    ])
    def test_binary_doc_type_pipeline(self, doc_type, make_fn, tmp_path):
        """Pipeline returns a non-empty string response for binary doc type.

        Args:
            doc_type: One of pdf, docx, xlsx, pptx.
            make_fn:  Fixture helper that creates a real binary temp file.
        """
        path     = make_fn(tmp_path)
        filename = os.path.basename(path)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(path, filename)
        assert len(chunks) >= 1, f"No chunks produced for {doc_type}"

        store  = _make_store(chunks)
        with patch('ollama.embed', return_value={'embeddings': [[0.1, 0.2, 0.3, 0.4]]}), \
             patch('ollama.chat', side_effect=_chat_mock):
            result = store.run_pipeline('How long do cats sleep?', streamlit_mode=True)

        assert isinstance(result['response'], str)
        assert len(result['response']) > 0, f"Empty response for doc_type={doc_type}"


# ---------------------------------------------------------------------------
# 2. Agent mode × all 8 doc types
# ---------------------------------------------------------------------------

class TestAgentModeAllDocTypes:
    """agent.run() returns a non-empty answer for every supported doc type."""

    @pytest.mark.parametrize("doc_type,text_content,filename", [
        ("txt",  "Cats sleep 16 hours a day.\nDogs are loyal.",    "test.txt"),
        ("md",   "# Animals\nCats sleep 16 hours.\nDogs loyal.",   "test.md"),
        ("csv",  "animal,fact\ncat,sleeps 16 hours\ndog,loyal",    "test.csv"),
        ("html", "<html><body><p>Cats sleep 16 hours.</p></body></html>", "test.html"),
    ])
    def test_text_doc_type_agent(self, doc_type, text_content, filename, tmp_path):
        """Agent returns a non-empty answer for text-based doc types.

        Args:
            doc_type:     One of txt, md, csv, html.
            text_content: Raw string written to the temp file.
            filename:     Filename passed to the chunker.
        """
        from src.rag.agent import Agent

        f = tmp_path / filename
        f.write_text(text_content)
        chunker = _CHUNKER_MAP[doc_type]
        chunks  = chunker(str(f), filename)
        assert len(chunks) >= 1

        store = _make_store(chunks)
        # Wire agent tools to use the real store
        store._expand_query       = lambda q: [q]
        store._hybrid_retrieve    = lambda queries, top_n, alpha=0.5: [(chunks[0], 0.8)]
        store._rerank             = lambda q, cands, top_n: [(chunks[0], 0.8, 0.8)]
        store._source_label       = lambda e: 'L1-1'
        store._synthesize         = MagicMock(return_value='Cats sleep 16 hours.')

        agent = Agent(store)

        def _chat(*args, **kwargs):
            return {'message': {'content': 'TOOL: finish(Cats sleep 16 hours a day.)'}}

        with patch('ollama.chat', side_effect=_chat):
            result = agent.run('How long do cats sleep?', streamlit_mode=True)

        assert isinstance(result['answer'], str)
        assert len(result['answer']) > 0, f"Empty agent answer for doc_type={doc_type}"

    @pytest.mark.parametrize("doc_type,make_fn", [
        ("pdf",  _make_pdf_file),
        ("docx", _make_docx_file),
        ("xlsx", _make_xlsx_file),
        ("pptx", _make_pptx_file),
    ])
    def test_binary_doc_type_agent(self, doc_type, make_fn, tmp_path):
        """Agent returns a non-empty answer for binary doc types.

        Args:
            doc_type: One of pdf, docx, xlsx, pptx.
            make_fn:  Fixture helper that creates a real binary temp file.
        """
        from src.rag.agent import Agent

        path     = make_fn(tmp_path)
        filename = os.path.basename(path)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(path, filename)
        assert len(chunks) >= 1

        store = _make_store(chunks)
        store._expand_query    = lambda q: [q]
        store._hybrid_retrieve = lambda queries, top_n, alpha=0.5: [(chunks[0], 0.8)]
        store._rerank          = lambda q, cands, top_n: [(chunks[0], 0.8, 0.8)]
        store._source_label    = lambda e: 'L1-1'
        store._synthesize      = MagicMock(return_value='Cats sleep 16 hours.')

        agent = Agent(store)

        def _chat(*args, **kwargs):
            return {'message': {'content': 'TOOL: finish(Cats sleep 16 hours a day.)'}}

        with patch('ollama.chat', side_effect=_chat):
            result = agent.run('How long do cats sleep?', streamlit_mode=True)

        assert isinstance(result['answer'], str)
        assert len(result['answer']) > 0, f"Empty agent answer for doc_type={doc_type}"

"""test_file_upload.py — File upload tests for all 8 document types using real
temp files through DocumentLoader._dispatch_chunker, then chat mode, agent mode,
and all 5 agent tools.

Mock strategy:
  Always mock: ollama.embed, ollama.chat, chromadb → EphemeralClient
  Never mock:  fitz, python-docx, openpyxl, python-pptx, BM25Okapi, file I/O
"""

import io
import uuid
import os
import pytest
import chromadb
from rank_bm25 import BM25Okapi
from unittest.mock import patch, MagicMock


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _fake_embed(**kwargs):
    return {'embeddings': [[0.1, 0.2, 0.3, 0.4]]}


def _fake_chat(**kwargs):
    if kwargs.get('stream'):
        return [{'message': {'content': 'mock response'}}]
    return {'message': {'content': 'mock response'}}


def _make_store(chunks):
    """Build an isolated VectorStore with EphemeralClient."""
    from src.rag.vector_store import VectorStore
    vs         = VectorStore()
    client     = chromadb.EphemeralClient()
    collection = client.get_or_create_collection(
        f'test_{uuid.uuid4().hex}', metadata={'hnsw:space': 'cosine'}
    )
    if chunks:
        ids    = [f'c{i}' for i in range(len(chunks))]
        texts  = [c['text'] for c in chunks]
        metas  = [{'source': c['source'], 'start_line': c['start_line'],
                   'end_line': c['end_line'], 'type': c['type']} for c in chunks]
        embeds = [[0.1, 0.2, 0.3, 0.4]] * len(chunks)
        collection.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
    vs.collection           = collection
    vs.chunks               = chunks
    vs.bm25_index           = BM25Okapi([c['text'].lower().split() for c in chunks]) if chunks else None
    vs.conversation_history = []
    return vs


PATCHES = [
    patch('ollama.embed', side_effect=_fake_embed),
    patch('ollama.chat',  side_effect=_fake_chat),
]


def _apply_patches(func):
    for p in reversed(PATCHES):
        func = p(func)
    return func


def _upload_chunks(tmp_path, filename, content_fn):
    """Create a real temp file and run it through DocumentLoader._dispatch_chunker."""
    from src.rag.document_loader import DocumentLoader
    loader   = DocumentLoader()
    ext      = os.path.splitext(filename)[1].lower()
    dtype    = loader.ext_to_type.get(ext, 'txt')
    filepath = str(tmp_path / filename)
    content_fn(filepath)
    file_info = {
        'filepath':      filepath,
        'filename':      filename,
        'detected_type': dtype,
        'is_misplaced':  False,
    }
    return loader._dispatch_chunker(file_info)


# ---------------------------------------------------------------------------
# Per-type file builders — write real files using the actual libraries
# ---------------------------------------------------------------------------

def _write_pdf(path):
    """Single-page PDF with nutrition facts."""
    import fitz
    doc  = fitz.open(); page = doc.new_page()
    page.insert_text((50, 50),
        'Nutrition facts: Bananas contain 89 calories per 100 grams. '
        'They are rich in potassium and vitamin B6. '
        'Avocados contain 160 calories per 100 grams and healthy fats.')
    doc.save(path); doc.close()


def _write_txt(path):
    """Plain-text dog facts, one fact per line."""
    with open(path, 'w') as f:
        f.write('Dogs were domesticated from wolves approximately 15000 years ago.\n')
        f.write('There are over 340 recognized dog breeds in the world.\n')
        f.write('A dog can understand about 165 words and gestures.\n')


def _write_docx(path):
    """DOCX with three programming language paragraphs."""
    from docx import Document
    doc = Document()
    doc.add_paragraph('Python is a programming language created by Guido van Rossum in 1991.')
    doc.add_paragraph('Java was created by James Gosling at Sun Microsystems in 1995.')
    doc.add_paragraph('Rust was developed at Mozilla Research and released in 2010.')
    doc.save(path)


def _write_xlsx(path):
    """XLSX with country/capital/population data."""
    import openpyxl
    wb = openpyxl.Workbook(); ws = wb.active
    ws.append(['Country', 'Capital', 'Population_Millions'])
    ws.append(['Brazil',  'Brasilia', 215])
    ws.append(['Japan',   'Tokyo',    125])
    ws.append(['Germany', 'Berlin',   84])
    wb.save(path)


def _write_csv(path):
    """CSV with programming language creator and year."""
    with open(path, 'w') as f:
        f.write('language,year,creator\n')
        f.write('Python,1991,Guido van Rossum\n')
        f.write('Rust,2010,Graydon Hoare\n')
        f.write('Go,2009,Robert Griesemer\n')


def _write_pptx(path):
    """PPTX with three history fact slides."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for title, body in [
        ('World War Two', 'World War Two ended in 1945 with Allied victory.'),
        ('Berlin Wall',   'The Berlin Wall fell on November 9 1989.'),
        ('Moon Landing',  'Apollo 11 landed on the Moon on July 20 1969.'),
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tb    = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        tb.text_frame.text = f'{title}: {body}'
    prs.save(path)


def _write_md(path):
    """Markdown file with coffee facts."""
    with open(path, 'w') as f:
        f.write('# Coffee Facts\n\n')
        f.write("Coffee was first discovered in Ethiopia around 850 AD by a goat herder named Kaldi.\n\n")
        f.write("Brazil is the world's largest coffee producer, accounting for 40 percent of global supply.\n\n")
        f.write('A standard cup of coffee contains between 80 and 100 milligrams of caffeine.\n')


def _write_html(path):
    """HTML file with dog biology facts."""
    with open(path, 'w') as f:
        f.write('<html><body><h1>Dog Facts</h1>')
        f.write("<p>A dog's sense of smell is up to 100000 times more powerful than a human's.</p>")
        f.write('<p>Border Collies are considered the most intelligent dog breed.</p>')
        f.write('<p>Greyhounds can run at speeds up to 72 kilometers per hour.</p>')
        f.write('</body></html>')


# ---------------------------------------------------------------------------
# 1. Chat mode — full pipeline per file type
# ---------------------------------------------------------------------------

class TestFileUploadChatModeTxt:
    """TXT file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        assert len(chunks) >= 1 and all(c['type'] == 'txt' for c in chunks)
        result = _make_store(chunks).run_pipeline('How many dog breeds are there?', streamlit_mode=True)
        assert 'response' in result and len(result['response']) > 0

    @_apply_patches
    def test_query_type(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        result = _make_store(chunks).run_pipeline('How many dog breeds are there?', streamlit_mode=True)
        assert result.get('query_type') in ('factual', 'general', 'comparison', 'summarise')

    @_apply_patches
    def test_source_label(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        assert 'L' in _make_store(chunks)._source_label(chunks[0])   # txt → L{s}-{e}

    @_apply_patches
    def test_source_cited(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        result = _make_store(chunks).run_pipeline('How many dog breeds are there?', streamlit_mode=True)
        assert any('dogs.txt' in e['source'] for e, _ in result['retrieved'])

    @_apply_patches
    def test_confidence_present(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        result = _make_store(chunks).run_pipeline('How many dog breeds are there?', streamlit_mode=True)
        assert 'is_confident' in result and result['is_confident'] in (True, False)


class TestFileUploadChatModePdf:
    """PDF file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'nutrition.pdf', _write_pdf)
        assert len(chunks) >= 1 and all(c['type'] == 'pdf' for c in chunks)
        result = _make_store(chunks).run_pipeline('How many calories in bananas?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_source_label(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'nutrition.pdf', _write_pdf)
        label  = _make_store(chunks)._source_label(chunks[0])
        assert 'p' in label or 'L' in label   # pdf → p{n}

    @_apply_patches
    def test_reranked(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'nutrition.pdf', _write_pdf)
        result = _make_store(chunks).run_pipeline('How many calories in bananas?', streamlit_mode=True)
        assert 'reranked' in result

    @_apply_patches
    def test_source_cited(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'nutrition.pdf', _write_pdf)
        result = _make_store(chunks).run_pipeline('How many calories in bananas?', streamlit_mode=True)
        assert any('nutrition.pdf' in e['source'] for e, _ in result['retrieved'])


class TestFileUploadChatModeDocx:
    """DOCX file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'languages.docx', _write_docx)
        assert len(chunks) >= 1 and all(c['type'] == 'docx' for c in chunks)
        result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_reranked(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'languages.docx', _write_docx)
        result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
        assert 'reranked' in result and len(result['reranked']) >= 1

    @_apply_patches
    def test_confidence(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'languages.docx', _write_docx)
        result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
        assert 'is_confident' in result and result['is_confident'] in (True, False)


class TestFileUploadChatModeXlsx:
    """XLSX file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'countries.xlsx', _write_xlsx)
        assert len(chunks) >= 1 and all(c['type'] == 'xlsx' for c in chunks)
        result = _make_store(chunks).run_pipeline('What is the capital of Brazil?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_source_label_row(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'countries.xlsx', _write_xlsx)
        label  = _make_store(chunks)._source_label(chunks[0])
        assert 'row' in label or 'L' in label   # xlsx → row{n}

    @_apply_patches
    def test_source_cited(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'countries.xlsx', _write_xlsx)
        result = _make_store(chunks).run_pipeline('What is the capital of Brazil?', streamlit_mode=True)
        assert any('countries.xlsx' in e['source'] for e, _ in result['retrieved'])


class TestFileUploadChatModeCsv:
    """CSV file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'languages.csv', _write_csv)
        assert len(chunks) >= 1 and all(c['type'] == 'csv' for c in chunks)
        result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_source_label_row(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'languages.csv', _write_csv)
        label  = _make_store(chunks)._source_label(chunks[0])
        assert 'row' in label or 'L' in label   # csv → row{n}

    @_apply_patches
    def test_query_type(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'languages.csv', _write_csv)
        result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
        assert result.get('query_type') in ('factual', 'general', 'comparison', 'summarise')


class TestFileUploadChatModePptx:
    """PPTX file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'history.pptx', _write_pptx)
        assert len(chunks) >= 1 and all(c['type'] == 'pptx' for c in chunks)
        result = _make_store(chunks).run_pipeline('When did the Berlin Wall fall?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_source_label_slide(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'history.pptx', _write_pptx)
        label  = _make_store(chunks)._source_label(chunks[0])
        assert 'slide' in label or 'L' in label   # pptx → slide{n}

    @_apply_patches
    def test_reranked(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'history.pptx', _write_pptx)
        result = _make_store(chunks).run_pipeline('When did the Berlin Wall fall?', streamlit_mode=True)
        assert 'reranked' in result


class TestFileUploadChatModeMd:
    """Markdown file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'coffee.md', _write_md)
        assert len(chunks) >= 1 and all(c['type'] == 'md' for c in chunks)
        result = _make_store(chunks).run_pipeline('Where was coffee discovered?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_source_label(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'coffee.md', _write_md)
        assert 'L' in _make_store(chunks)._source_label(chunks[0])   # md → L{s}-{e}

    @_apply_patches
    def test_confidence(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'coffee.md', _write_md)
        result = _make_store(chunks).run_pipeline('Where was coffee discovered?', streamlit_mode=True)
        assert 'is_confident' in result and result['is_confident'] in (True, False)


class TestFileUploadChatModeHtml:
    """HTML file upload → chat pipeline."""

    @_apply_patches
    def test_factual(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.html', _write_html)
        assert len(chunks) >= 1 and all(c['type'] == 'html' for c in chunks)
        result = _make_store(chunks).run_pipeline('How fast can a greyhound run?', streamlit_mode=True)
        assert 'response' in result

    @_apply_patches
    def test_source_label(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.html', _write_html)
        label  = _make_store(chunks)._source_label(chunks[0])
        assert 's' in label or 'L' in label   # html → s{n}

    @_apply_patches
    def test_source_cited(self, mock_chat, mock_embed, tmp_path):
        chunks = _upload_chunks(tmp_path, 'dogs.html', _write_html)
        result = _make_store(chunks).run_pipeline('How fast can a greyhound run?', streamlit_mode=True)
        assert any('dogs.html' in e['source'] for e, _ in result['retrieved'])


# ---------------------------------------------------------------------------
# 2. Agent mode — rag_search for each file type
# ---------------------------------------------------------------------------

class TestFileUploadAgentMode:
    """All 8 file types loaded via upload, then queried through Agent.run()."""

    @_apply_patches
    def test_txt_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(dog breeds)'}}
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('How many dog breeds are there?', streamlit_mode=True)
        assert 'answer' in result
        assert any(s['tool'] == 'rag_search' for s in result['steps'])

    @_apply_patches
    def test_pdf_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(banana calories)'}}
        chunks = _upload_chunks(tmp_path, 'nutrition.pdf', _write_pdf)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('How many calories in bananas?', streamlit_mode=True)
        assert 'answer' in result
        assert any(s['tool'] == 'rag_search' for s in result['steps'])

    @_apply_patches
    def test_docx_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Python creator)'}}
        chunks = _upload_chunks(tmp_path, 'languages.docx', _write_docx)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('Who created Python?', streamlit_mode=True)
        assert 'answer' in result

    @_apply_patches
    def test_xlsx_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Brazil capital)'}}
        chunks = _upload_chunks(tmp_path, 'countries.xlsx', _write_xlsx)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('What is the capital of Brazil?', streamlit_mode=True)
        assert 'answer' in result

    @_apply_patches
    def test_csv_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Python year)'}}
        chunks = _upload_chunks(tmp_path, 'languages.csv', _write_csv)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('When was Python created?', streamlit_mode=True)
        assert 'answer' in result

    @_apply_patches
    def test_pptx_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(Berlin Wall 1989)'}}
        chunks = _upload_chunks(tmp_path, 'history.pptx', _write_pptx)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('When did the Berlin Wall fall?', streamlit_mode=True)
        assert 'answer' in result

    @_apply_patches
    def test_md_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(coffee Ethiopia)'}}
        chunks = _upload_chunks(tmp_path, 'coffee.md', _write_md)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('Where was coffee discovered?', streamlit_mode=True)
        assert 'answer' in result

    @_apply_patches
    def test_html_agent_rag_search(self, mock_chat, mock_embed, tmp_path):
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(dog smell)'}}
        chunks = _upload_chunks(tmp_path, 'dogs.html', _write_html)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('What is special about dog smell?', streamlit_mode=True)
        assert 'answer' in result


# ---------------------------------------------------------------------------
# 3. All 5 agent tools with real uploaded file content
# ---------------------------------------------------------------------------

class TestFileUploadAllTools:
    """All 5 tools (rag_search, calculator, summarise, sentiment, finish) exercised
    with content loaded from real temp files via DocumentLoader._dispatch_chunker."""

    @_apply_patches
    def test_rag_search_on_txt(self, mock_chat, mock_embed, tmp_path):
        """rag_search returns chunks from a real uploaded txt file."""
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_rag_search('dog breeds')
        assert len(result) > 0 and 'dogs.txt' in result

    @_apply_patches
    def test_rag_search_on_csv(self, mock_chat, mock_embed, tmp_path):
        """rag_search returns chunks from a real uploaded csv file."""
        chunks = _upload_chunks(tmp_path, 'languages.csv', _write_csv)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_rag_search('Python creator')
        assert len(result) > 0

    @_apply_patches
    def test_calculator_on_xlsx_values(self, mock_chat, mock_embed, tmp_path):
        """calculator evaluates numbers that might come from an xlsx file."""
        chunks = _upload_chunks(tmp_path, 'countries.xlsx', _write_xlsx)
        from src.rag.agent import Agent
        assert Agent(_make_store(chunks))._tool_calculator('215 + 125 + 84') == '424'

    @_apply_patches
    def test_calculator_percentage(self, mock_chat, mock_embed, tmp_path):
        """calculator handles percentage expressions."""
        from src.rag.agent import Agent
        result = Agent(_make_store([]))._tool_calculator('15% of 89')
        assert float(result) == pytest.approx(13.35)

    @_apply_patches
    def test_summarise_on_md(self, mock_chat, mock_embed, tmp_path):
        """summarise tool works on text from a real md file."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'Coffee originated in Ethiopia.'}}
        chunks = _upload_chunks(tmp_path, 'coffee.md', _write_md)
        text   = ' '.join(c['text'] for c in chunks)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_summarise(text)
        assert isinstance(result, str) and len(result) > 0

    @_apply_patches
    def test_summarise_on_docx(self, mock_chat, mock_embed, tmp_path):
        """summarise tool works on text from a real docx file."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'Python was created by Guido.'}}
        chunks = _upload_chunks(tmp_path, 'languages.docx', _write_docx)
        text   = ' '.join(c['text'] for c in chunks)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_summarise(text)
        assert isinstance(result, str) and len(result) > 0

    @_apply_patches
    def test_sentiment_on_html(self, mock_chat, mock_embed, tmp_path):
        """sentiment tool returns all 4 required fields for html file content."""
        resp = ('Sentiment: Positive\nTone: informative\n'
                'Key phrases: powerful, intelligent\nExplanation: Factual and positive.')
        mock_chat.side_effect = lambda **kw: {'message': {'content': resp}}
        chunks = _upload_chunks(tmp_path, 'dogs.html', _write_html)
        text   = ' '.join(c['text'] for c in chunks)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_sentiment(text)
        assert 'Sentiment' in result and 'Tone' in result
        assert 'Key phrases' in result and 'Explanation' in result

    @_apply_patches
    def test_sentiment_on_txt(self, mock_chat, mock_embed, tmp_path):
        """sentiment tool works on text from a real txt file."""
        resp = ('Sentiment: Neutral\nTone: factual\n'
                'Key phrases: domesticated, breeds\nExplanation: Objective facts.')
        mock_chat.side_effect = lambda **kw: {'message': {'content': resp}}
        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        text   = ' '.join(c['text'] for c in chunks)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks))._tool_sentiment(text)
        assert 'Sentiment' in result

    @_apply_patches
    def test_finish_after_rag_search(self, mock_chat, mock_embed, tmp_path):
        """finish step is reached after rag_search on a real uploaded file."""
        responses = iter([
            {'message': {'content': 'TOOL: rag_search(coffee Ethiopia)'}},
            {'message': {'content': 'TOOL: finish(Coffee was discovered in Ethiopia around 850 AD.)'}},
        ])
        mock_chat.side_effect = lambda **kw: next(responses)
        chunks = _upload_chunks(tmp_path, 'coffee.md', _write_md)
        from src.rag.agent import Agent
        result = Agent(_make_store(chunks)).run('Where was coffee discovered?', streamlit_mode=True)
        assert any(s['tool'] == 'finish' for s in result['steps'])
        assert result['answer'] is not None

    @_apply_patches
    def test_finish_calc_auto(self, mock_chat, mock_embed, tmp_path):
        """calculator auto-finish: tool result goes directly to finish."""
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: calculator(215 + 125)'}}
        from src.rag.agent import Agent
        result = Agent(_make_store([])).run('What is 215 plus 125?', streamlit_mode=True)
        assert any(s['tool'] == 'calculator' for s in result['steps'])
        assert any(s['tool'] == 'finish' for s in result['steps'])
        assert '340' in result['answer']

    @_apply_patches
    def test_all_5_tools_smoke(self, mock_chat, mock_embed, tmp_path):
        """Smoke test: all 5 tools are callable in sequence with real file content."""
        from src.rag.agent import Agent

        chunks = _upload_chunks(tmp_path, 'dogs.txt', _write_txt)
        store  = _make_store(chunks)
        agent  = Agent(store)

        # 1. rag_search
        assert len(agent._tool_rag_search('dog breeds')) > 0

        # 2. calculator
        assert agent._tool_calculator('340 * 2') == '680'

        # 3. summarise
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'Dogs were domesticated long ago.'}}
        assert len(agent._tool_summarise(' '.join(c['text'] for c in chunks))) > 0

        # 4. sentiment
        mock_chat.side_effect = lambda **kw: {
            'message': {'content': 'Sentiment: Neutral\nTone: factual\nKey phrases: dogs\nExplanation: OK.'}
        }
        assert 'Sentiment' in agent._tool_sentiment(' '.join(c['text'] for c in chunks))

        # 5. finish — via auto-finish after rag_search
        mock_chat.side_effect = lambda **kw: {'message': {'content': 'TOOL: rag_search(dog breeds)'}}
        result = Agent(store).run('How many dog breeds are there?', streamlit_mode=True)
        assert any(s['tool'] == 'finish' for s in result['steps'])

"""test_url_pipeline.py — Pipeline feature tests for URL-loaded documents.

Covers: query type detection, confidence flag, reranked list, source labels,
        and source-cited checks for all 8 URL types.

Mock strategy:
  Always mock: ollama.embed, ollama.chat, chromadb → EphemeralClient, requests.get
  Never mock:  fitz, python-docx, openpyxl, python-pptx, beautifulsoup4, BM25Okapi

Reason for split: max 500 lines per file per CLAUDE.md.
Chat and agent mode URL ingestion tests are in test_url_ingestion.py.
"""

import io
import uuid
import pytest
import chromadb
from rank_bm25 import BM25Okapi
from unittest.mock import patch, MagicMock


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _fake_embed(**kwargs):
    """Return a fixed 4-dim embedding vector to satisfy ollama.embed calls."""
    return {'embeddings': [[0.1, 0.2, 0.3, 0.4]]}


def _fake_chat(**kwargs):
    """Return a canned chat response (stream or single) to satisfy ollama.chat calls."""
    if kwargs.get('stream'):
        return [{'message': {'content': 'mock response'}}]
    return {'message': {'content': 'mock response'}}


def _make_store(chunks):
    """Build an isolated VectorStore backed by an EphemeralClient for the given chunks."""
    from src.rag.vector_store import VectorStore
    vs         = VectorStore()
    client     = chromadb.EphemeralClient()
    collection = client.get_or_create_collection(
        f'test_{uuid.uuid4().hex}', metadata={'hnsw:space': 'cosine'}
    )
    if chunks:
        ids    = [f'c{i}' for i in range(len(chunks))]
        texts  = [c['text'] for c in chunks]
        metas  = [{'source': c['source'], 'start_line': c['start_line'],
                   'end_line': c['end_line'], 'type': c['type']} for c in chunks]
        embeds = [[0.1, 0.2, 0.3, 0.4]] * len(chunks)
        collection.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
    vs.collection           = collection
    vs.chunks               = chunks
    vs.bm25_index           = BM25Okapi([c['text'].lower().split() for c in chunks]) if chunks else None
    vs.conversation_history = []
    return vs


PATCHES = [
    patch('ollama.embed', side_effect=_fake_embed),
    patch('ollama.chat',  side_effect=_fake_chat),
]


def _apply_patches(func):
    """Decorator that applies both ollama.embed and ollama.chat patches to a test method."""
    for p in reversed(PATCHES):
        func = p(func)
    return func


def _mock_url_response(content, content_type='text/html', encoding='utf-8'):
    """Build a MagicMock simulating a requests.Response with the given content and headers."""
    resp = MagicMock()
    resp.content  = content
    resp.headers  = {'Content-Type': content_type}
    resp.encoding = encoding
    resp.raise_for_status = MagicMock()
    return resp


# ---------------------------------------------------------------------------
# 3a. URL — query type detection
# ---------------------------------------------------------------------------

class TestUrlQueryTypeDetected:
    """run_pipeline returns a valid query_type string for URL-loaded documents."""

    @_apply_patches
    def test_url_html_query_type(self, mock_chat, mock_embed):
        """HTML URL: query_type in result is one of the 4 valid classifier labels."""
        html = (b'<html><body><p>Dogs have an extraordinary sense of smell. '
                b'They can detect scents from miles away.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        assert len(chunks) >= 1
        result = _make_store(chunks).run_pipeline('What is special about dog smell?', streamlit_mode=True)
        assert result.get('query_type') in ('factual', 'general', 'comparison', 'summarise')

    @_apply_patches
    def test_url_csv_query_type(self, mock_chat, mock_embed):
        """CSV URL: query_type in result is one of the 4 valid classifier labels."""
        csv_content = b'language,year,creator\nPython,1991,Guido van Rossum\n'
        with patch('requests.get', return_value=_mock_url_response(csv_content, 'text/csv')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/languages.csv')
        if chunks:
            result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
            assert result.get('query_type') in ('factual', 'general', 'comparison', 'summarise')

    @_apply_patches
    def test_url_txt_query_type(self, mock_chat, mock_embed):
        """Plain-text URL: query_type in result is one of the 4 valid classifier labels."""
        txt_content = (b'Dogs were domesticated from wolves approximately 15000 years ago.\n'
                       b'There are over 340 recognized dog breeds in the world.')
        with patch('requests.get', return_value=_mock_url_response(txt_content, 'text/plain')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs.txt')
        if chunks:
            result = _make_store(chunks).run_pipeline('How were dogs domesticated?', streamlit_mode=True)
            assert result.get('query_type') in ('factual', 'general', 'comparison', 'summarise')


# ---------------------------------------------------------------------------
# 3b. URL — confidence flag
# ---------------------------------------------------------------------------

class TestUrlConfidenceFlag:
    """run_pipeline result contains a boolean is_confident key for URL-loaded documents."""

    @_apply_patches
    def test_url_html_confidence(self, mock_chat, mock_embed):
        """HTML URL: is_confident key is present and is a boolean in the pipeline result."""
        html = (b'<html><body><p>Dogs have an extraordinary sense of smell. '
                b'They can detect scents from miles away.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        assert len(chunks) >= 1
        result = _make_store(chunks).run_pipeline('What is special about dog smell?', streamlit_mode=True)
        assert 'is_confident' in result and result['is_confident'] in (True, False)

    @_apply_patches
    def test_url_csv_confidence(self, mock_chat, mock_embed):
        """CSV URL: is_confident key is present and is a boolean in the pipeline result."""
        csv_content = b'language,year,creator\nPython,1991,Guido van Rossum\n'
        with patch('requests.get', return_value=_mock_url_response(csv_content, 'text/csv')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/langs.csv')
        if chunks:
            result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
            assert 'is_confident' in result and result['is_confident'] in (True, False)

    @_apply_patches
    def test_url_md_confidence(self, mock_chat, mock_embed):
        """Markdown URL: is_confident key is present and is a boolean in the pipeline result."""
        md_content = b'# Coffee\n\nCoffee was first discovered in Ethiopia around 850 AD.'
        with patch('requests.get', return_value=_mock_url_response(md_content, 'text/markdown')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/coffee.md')
        if chunks:
            result = _make_store(chunks).run_pipeline('Where was coffee discovered?', streamlit_mode=True)
            assert 'is_confident' in result and result['is_confident'] in (True, False)


# ---------------------------------------------------------------------------
# 3c. URL — reranked list
# ---------------------------------------------------------------------------

class TestUrlReranked:
    """run_pipeline result contains a non-empty reranked list for URL-loaded documents."""

    @_apply_patches
    def test_url_html_reranked_present(self, mock_chat, mock_embed):
        """HTML URL with multiple paragraphs: reranked list is present and has at least one entry."""
        html = (b'<html><body>'
                b'<p>Dogs have an extraordinary sense of smell and can detect scents from miles away.</p>'
                b'<p>Dogs were domesticated from wolves approximately 15000 years ago in East Asia.</p>'
                b'</body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        if chunks:
            result = _make_store(chunks).run_pipeline('What is special about dog smell?', streamlit_mode=True)
            assert 'reranked' in result and len(result['reranked']) >= 1

    @_apply_patches
    def test_url_xlsx_reranked_present(self, mock_chat, mock_embed):
        """XLSX URL with multiple rows: reranked key is present in the pipeline result."""
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Country', 'Capital', 'Population'])
        ws.append(['Brazil', 'Brasilia', '215 million'])
        ws.append(['Japan',  'Tokyo',    '125 million'])
        ws.append(['Germany','Berlin',   '84 million'])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        with patch('requests.get', return_value=_mock_url_response(xlsx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/countries.xlsx')
        if chunks:
            result = _make_store(chunks).run_pipeline('What is the capital of Brazil?', streamlit_mode=True)
            assert 'reranked' in result


# ---------------------------------------------------------------------------
# 3d. URL — source labels
# ---------------------------------------------------------------------------

class TestUrlSourceLabels:
    """_source_label returns the correct format string for each URL document type."""

    @_apply_patches
    def test_url_html_source_label_section(self, mock_chat, mock_embed):
        """HTML URL chunk: _source_label contains 's' (section) or 'L' (line) marker."""
        html = (b'<html><body><p>Dogs have an extraordinary sense of smell. '
                b'They can detect scents from miles away.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        assert len(chunks) >= 1
        label = _make_store(chunks)._source_label(chunks[0])
        assert 's' in label or 'L' in label   # html → s{n}

    @_apply_patches
    def test_url_csv_source_label_row(self, mock_chat, mock_embed):
        """CSV URL chunk: _source_label contains 'row' or 'L' marker."""
        csv_content = b'language,year,creator\nPython,1991,Guido van Rossum\n'
        with patch('requests.get', return_value=_mock_url_response(csv_content, 'text/csv')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/langs.csv')
        if chunks:
            label = _make_store(chunks)._source_label(chunks[0])
            assert 'row' in label or 'L' in label   # csv → row{n}

    @_apply_patches
    def test_url_xlsx_source_label_row(self, mock_chat, mock_embed):
        """XLSX URL chunk: _source_label contains 'row' or 'L' marker."""
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Country', 'Capital']); ws.append(['Brazil', 'Brasilia'])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        with patch('requests.get', return_value=_mock_url_response(xlsx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/countries.xlsx')
        if chunks:
            label = _make_store(chunks)._source_label(chunks[0])
            assert 'row' in label or 'L' in label   # xlsx → row{n}

    @_apply_patches
    def test_url_pptx_source_label_slide(self, mock_chat, mock_embed):
        """PPTX URL chunk: _source_label contains 'slide' or 'L' marker."""
        from pptx import Presentation
        from pptx.util import Inches
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'The Berlin Wall fell on November 9 1989.'
        buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        with patch('requests.get', return_value=_mock_url_response(pptx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/history.pptx')
        if chunks:
            label = _make_store(chunks)._source_label(chunks[0])
            assert 'slide' in label or 'L' in label   # pptx → slide{n}

    @_apply_patches
    def test_url_pdf_source_label_page(self, mock_chat, mock_embed):
        """PDF URL chunk: _source_label contains 'p' (page) or 'L' marker."""
        import fitz
        doc  = fitz.open(); page = doc.new_page()
        page.insert_text((50, 50), 'Bananas contain 89 calories per 100 grams.')
        buf = io.BytesIO(); doc.save(buf); pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_url_response(pdf_bytes, 'application/pdf')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/nutrition.pdf')
        if chunks:
            label = _make_store(chunks)._source_label(chunks[0])
            assert 'p' in label or 'L' in label   # pdf → p{n}

    @_apply_patches
    def test_url_txt_source_label(self, mock_chat, mock_embed):
        """Plain-text URL chunk: _source_label contains 'L' (line range) marker."""
        txt_content = (b'Dogs were domesticated from wolves approximately 15000 years ago.\n'
                       b'There are over 340 recognized dog breeds in the world.')
        with patch('requests.get', return_value=_mock_url_response(txt_content, 'text/plain')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs.txt')
        if chunks:
            label = _make_store(chunks)._source_label(chunks[0])
            assert 'L' in label   # txt → L{s}-{e}

    @_apply_patches
    def test_url_md_source_label(self, mock_chat, mock_embed):
        """Markdown URL chunk: _source_label contains 'L' (line range) marker."""
        md_content = b'# Coffee\n\nCoffee was first discovered in Ethiopia around 850 AD.'
        with patch('requests.get', return_value=_mock_url_response(md_content, 'text/markdown')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/coffee.md')
        if chunks:
            label = _make_store(chunks)._source_label(chunks[0])
            assert 'L' in label   # md → L{s}-{e}


# ---------------------------------------------------------------------------
# 3e. URL — source cited in retrieved list
# ---------------------------------------------------------------------------

class TestUrlSourceCited:
    """run_pipeline retrieved list includes the URL hostname as the source for each document type."""

    @_apply_patches
    def test_url_html_source_cited(self, mock_chat, mock_embed):
        """HTML URL: retrieved chunks contain 'example.com' in their source field."""
        html = (b'<html><body><p>Dogs have an extraordinary sense of smell. '
                b'They can detect scents from miles away.</p></body></html>')
        with patch('requests.get', return_value=_mock_url_response(html, 'text/html')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/dogs')
        assert len(chunks) >= 1
        result = _make_store(chunks).run_pipeline('What is special about dog smell?', streamlit_mode=True)
        assert 'retrieved' in result
        assert any('example.com' in e['source'] for e, _ in result['retrieved'])

    @_apply_patches
    def test_url_csv_source_cited(self, mock_chat, mock_embed):
        """CSV URL: retrieved chunks contain 'example.com' in their source field."""
        csv_content = b'language,year,creator\nPython,1991,Guido van Rossum\n'
        with patch('requests.get', return_value=_mock_url_response(csv_content, 'text/csv')):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/langs.csv')
        if chunks:
            result = _make_store(chunks).run_pipeline('Who created Python?', streamlit_mode=True)
            assert any('example.com' in e['source'] for e, _ in result['retrieved'])

    @_apply_patches
    def test_url_xlsx_source_cited(self, mock_chat, mock_embed):
        """XLSX URL: retrieved chunks contain 'example.com' in their source field."""
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Country', 'Capital']); ws.append(['Brazil', 'Brasilia'])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        with patch('requests.get', return_value=_mock_url_response(xlsx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/countries.xlsx')
        if chunks:
            result = _make_store(chunks).run_pipeline('What is the capital of Brazil?', streamlit_mode=True)
            assert any('example.com' in e['source'] for e, _ in result['retrieved'])

    @_apply_patches
    def test_url_pptx_source_cited(self, mock_chat, mock_embed):
        """PPTX URL: retrieved chunks contain 'example.com' in their source field."""
        from pptx import Presentation
        from pptx.util import Inches
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'The Berlin Wall fell on November 9 1989.'
        buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()
        ct = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        with patch('requests.get', return_value=_mock_url_response(pptx_bytes, ct)):
            from src.rag.document_loader import DocumentLoader
            chunks = DocumentLoader().chunk_url('https://example.com/history.pptx')
        if chunks:
            result = _make_store(chunks).run_pipeline('When did the Berlin Wall fall?', streamlit_mode=True)
            assert any('example.com' in e['source'] for e, _ in result['retrieved'])

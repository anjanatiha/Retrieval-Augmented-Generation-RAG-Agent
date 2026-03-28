"""Integration tests — document loading: all 8 file types, xls, and URL ingestion.

Covers sections 1 and 2 of the original integration test suite:
  1. All 8 file types + xls (real libraries, tmp files)
  2. URL ingestion (mock requests.get only)

Mock strategy (per CLAUDE.md):
  Always mock:   ollama.embed, ollama.chat, chromadb → EphemeralClient, requests.get
  Never mock:    fitz, python-docx, openpyxl, xlrd, python-pptx, beautifulsoup4,
                 BM25Okapi, chunk truncation, misplaced detection, calculator eval

Reason for split: max 500 lines per file per CLAUDE.md.
"""

import io
import os
import tempfile
from unittest.mock import MagicMock, patch

import pytest

from src.rag.binary_chunkers import chunk_docx, chunk_pdf, chunk_pptx, chunk_xls, chunk_xlsx

# ── Chunker functions are now standalone module-level functions ───────────────
from src.rag.chunkers import chunk_csv, chunk_html, chunk_md, chunk_txt, truncate_chunk

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _mock_resp(content, content_type='text/html', encoding='utf-8'):
    """Build a MagicMock simulating requests.Response with the given content and headers."""
    resp = MagicMock()
    resp.content = content if isinstance(content, bytes) else content.encode(encoding)
    resp.headers = {'Content-Type': content_type}
    resp.encoding = encoding
    resp.raise_for_status = MagicMock()
    return resp


# ============================================================
# 1. All 8 file types + xls (real libraries, tmp files)
# ============================================================

class TestLoadPdf:
    """Integration tests for PDF chunk extraction using real fitz."""

    def test_load_pdf(self, tmp_path):
        """PDF with one page of text: at least one chunk of type 'pdf' is returned."""
        import fitz

        from src.rag.document_loader import DocumentLoader
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Cats sleep sixteen hours a day. They are nocturnal hunters.")
        path = str(tmp_path / 'test.pdf')
        doc.save(path); doc.close()
        chunks = chunk_pdf(path, 'test.pdf')
        assert len(chunks) >= 1
        assert all(c['type'] == 'pdf' for c in chunks)
        assert any('cats' in c['text'].lower() for c in chunks)


class TestLoadTxt:
    """Integration tests for plain-text chunk extraction."""

    def test_load_txt(self, tmp_path):
        """Two-line TXT file: exactly 2 chunks of type 'txt' are returned."""
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.txt'
        f.write_text('Cats sleep sixteen hours a day.\nThey are nocturnal hunters.\n')
        chunks = chunk_txt(str(f), 'test.txt')
        assert len(chunks) == 2
        assert all(c['type'] == 'txt' for c in chunks)


class TestLoadDocx:
    """Integration tests for DOCX paragraph and table extraction using real python-docx."""

    def test_load_docx(self, tmp_path):
        """DOCX with two paragraphs: at least one chunk of type 'docx' is returned."""
        from docx import Document

        from src.rag.document_loader import DocumentLoader
        doc = Document()
        doc.add_paragraph('Cats sleep sixteen hours a day.')
        doc.add_paragraph('They are excellent hunters.')
        path = str(tmp_path / 'test.docx')
        doc.save(path)
        chunks = chunk_docx(path, 'test.docx')
        assert len(chunks) >= 1
        assert all(c['type'] == 'docx' for c in chunks)

    def test_docx_table_rows_extracted(self, tmp_path):
        """DOCX with a 2×2 table: table cell text appears in the combined chunk output."""
        from docx import Document

        from src.rag.document_loader import DocumentLoader
        doc   = Document()
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = 'Name'
        table.rows[0].cells[1].text = 'Alice'
        table.rows[1].cells[0].text = 'Job'
        table.rows[1].cells[1].text = 'Engineer'
        path = str(tmp_path / 'table.docx')
        doc.save(path)
        chunks = chunk_docx(path, 'table.docx')
        full   = ' '.join(c['text'] for c in chunks)
        assert 'Alice' in full or 'Name' in full

    def test_docx_merged_cells_deduplicated(self, tmp_path):
        """DOCX with a simulated merged cell: duplicate cell text appears only once."""
        from docx import Document

        from src.rag.document_loader import DocumentLoader
        doc   = Document()
        table = doc.add_table(rows=1, cols=3)
        table.rows[0].cells[0].text = 'Merged'
        table.rows[0].cells[1].text = 'Merged'   # simulate merged cell
        table.rows[0].cells[2].text = 'Different'
        path = str(tmp_path / 'merged.docx')
        doc.save(path)
        chunks = chunk_docx(path, 'merged.docx')
        full   = ' '.join(c['text'] for c in chunks)
        assert full.count('Merged') == 1


class TestLoadXlsx:
    """Integration tests for XLSX row extraction using real openpyxl."""

    def test_load_xlsx(self, tmp_path):
        """XLSX with 2 data rows (plus header): exactly 2 chunks of type 'xlsx' are returned."""
        import openpyxl

        from src.rag.document_loader import DocumentLoader
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(['Name', 'Age'])
        ws.append(['Alice', 30])
        ws.append(['Bob', 25])
        path = str(tmp_path / 'test.xlsx')
        wb.save(path)
        chunks = chunk_xlsx(path, 'test.xlsx')
        assert len(chunks) == 2
        assert all(c['type'] == 'xlsx' for c in chunks)
        assert 'Alice' in chunks[0]['text']


class TestLoadXls:
    """Integration tests for legacy XLS extraction using real xlrd."""

    def test_load_xls(self, tmp_path):
        """XLS with 1 data row: exactly 1 chunk of type 'xls' is returned with 'Alice' in text."""
        import xlwt

        from src.rag.document_loader import DocumentLoader
        wb = xlwt.Workbook()
        ws = wb.add_sheet('Sheet1')
        ws.write(0, 0, 'Name'); ws.write(0, 1, 'Score')
        ws.write(1, 0, 'Alice'); ws.write(1, 1, 95)
        path = str(tmp_path / 'test.xls')
        wb.save(path)
        chunks = chunk_xls(path, 'test.xls')
        assert len(chunks) == 1
        assert 'Alice' in chunks[0]['text']


class TestLoadPptx:
    """Integration tests for PPTX slide extraction using real python-pptx."""

    def test_load_pptx(self, tmp_path):
        """PPTX with one text box: at least one chunk of type 'pptx' containing 'Cats' is returned."""
        from pptx import Presentation
        from pptx.util import Inches

        from src.rag.document_loader import DocumentLoader
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'Cats are wonderful nocturnal hunters.'
        path  = str(tmp_path / 'test.pptx')
        prs.save(path)
        chunks = chunk_pptx(path, 'test.pptx')
        assert len(chunks) >= 1
        assert all(c['type'] == 'pptx' for c in chunks)
        assert 'Cats' in chunks[0]['text']


class TestLoadCsv:
    """Integration tests for CSV row extraction using the real csv.DictReader."""

    def test_load_csv(self, tmp_path):
        """CSV with 2 data rows (plus header): exactly 2 chunks of type 'csv' are returned."""
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.csv'
        f.write_text('name,age\nAlice,30\nBob,25\n')
        chunks = chunk_csv(str(f), 'test.csv')
        assert len(chunks) == 2
        assert all(c['type'] == 'csv' for c in chunks)


class TestLoadMd:
    """Integration tests for Markdown chunk extraction including heading/bold/italic stripping."""

    def test_load_md(self, tmp_path):
        """Markdown with heading and bold: at least one chunk of type 'md' with markers stripped."""
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.md'
        f.write_text('# Title\n**bold** cats sleep a lot\n')
        chunks = chunk_md(str(f), 'test.md')
        assert len(chunks) >= 1
        assert all(c['type'] == 'md' for c in chunks)
        assert all('**' not in c['text'] for c in chunks)
        assert all('# ' not in c['text'] for c in chunks)


class TestLoadHtml:
    """Integration tests for HTML chunk extraction using real BeautifulSoup."""

    def test_load_html(self, tmp_path):
        """HTML with one paragraph: at least one chunk of type 'html' with all tags stripped."""
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.html'
        f.write_text('<html><body><p>Cats sleep. They hunt at night. Amazing animals.</p></body></html>')
        chunks = chunk_html(str(f), 'test.html')
        assert len(chunks) >= 1
        assert all(c['type'] == 'html' for c in chunks)
        assert all('<' not in c['text'] for c in chunks)


class TestMisplacedFile:
    """Integration tests for misplaced file detection in scan_all_files."""

    def _fake_folders(self, tmp_path):
        """Return a doc_folders dict pointing all type dirs into tmp_path."""
        return {
            'pdf':  str(tmp_path / 'pdfs'),
            'txt':  str(tmp_path / 'txts'),
            'docx': str(tmp_path / 'docx'),
            'xlsx': str(tmp_path / 'xlsx'),
            'pptx': str(tmp_path / 'pptx'),
            'csv':  str(tmp_path / 'csv'),
            'md':   str(tmp_path / 'md'),
            'html': str(tmp_path / 'html'),
        }

    def test_misplaced_file_detected_and_chunked(self, tmp_path):
        """PDF file placed in the txts folder: detected as misplaced with correct detected_type."""
        from src.rag.document_loader import DocumentLoader
        txts_dir = tmp_path / 'txts'
        txts_dir.mkdir()
        (txts_dir / 'report.pdf').write_bytes(b'%PDF fake')
        loader = DocumentLoader()
        loader.docs_root   = str(tmp_path)
        loader.doc_folders = self._fake_folders(tmp_path)
        files = loader.scan_all_files()
        assert len(files) == 1
        assert files[0]['is_misplaced'] is True
        assert files[0]['detected_type'] == 'pdf'

    def test_recursive_scan_finds_deeply_nested_file(self, tmp_path):
        """File three levels deep under docs_root is found and has a correct relative path."""
        from src.rag.document_loader import DocumentLoader
        deep = tmp_path / 'project' / 'reports' / '2025'
        deep.mkdir(parents=True)
        (deep / 'annual.pdf').write_bytes(b'%PDF fake')
        loader = DocumentLoader()
        loader.docs_root   = str(tmp_path)
        loader.doc_folders = self._fake_folders(tmp_path)
        files = loader.scan_all_files()
        assert len(files) == 1
        assert files[0]['filename'] == os.path.join('project', 'reports', '2025', 'annual.pdf')
        assert files[0]['is_misplaced'] is True  # not in canonical 'pdfs/' folder


class TestTruncateChunk:
    """Integration tests for _truncate_chunk enforcing 300-word and 1200-char limits."""

    def test_truncate_chunk_300_words(self):
        """400-word input: truncate_chunk returns at most 300 words."""
        text   = ' '.join(['word'] * 400)
        result = truncate_chunk(text)
        assert len(result.split()) <= 300

    def test_truncate_chunk_1200_chars(self):
        """2000-char input: truncate_chunk returns at most 1200 characters."""
        text   = 'a' * 2000
        result = truncate_chunk(text)
        assert len(result) <= 1200


# ============================================================
# 2. URL ingestion (mock requests.get only)
# ============================================================

class TestUrlIngestion:
    """Integration tests for chunk_url — all URL type detection paths and error handling."""

    def test_url_html_webpage(self):
        """HTML content-type URL: chunk_url returns at least one chunk of type 'html'."""
        from src.rag.document_loader import DocumentLoader
        html = b'<html><body><p>Hello world. Second sentence. Third one here.</p></body></html>'
        with patch('requests.get', return_value=_mock_resp(html, 'text/html')):
            chunks = DocumentLoader().chunk_url('https://example.com')
        assert len(chunks) >= 1
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_remote_pdf(self):
        """PDF URL (application/pdf content-type): chunk_url returns a list (possibly empty)."""
        import fitz

        from src.rag.document_loader import DocumentLoader
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF content. Another sentence.")
        buf = io.BytesIO(); doc.save(buf); pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_resp(pdf_bytes, 'application/pdf')):
            chunks = DocumentLoader().chunk_url('https://example.com/file.pdf')
        assert isinstance(chunks, list)

    def test_url_remote_docx(self):
        """DOCX URL (wordprocessingml content-type): chunk_url returns a list."""
        from docx import Document

        from src.rag.document_loader import DocumentLoader
        doc = Document(); doc.add_paragraph('Remote docx content here.')
        buf = io.BytesIO(); doc.save(buf); docx_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_resp(
            docx_bytes,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )):
            chunks = DocumentLoader().chunk_url('https://example.com/file.docx')
        assert isinstance(chunks, list)

    def test_url_remote_xlsx(self):
        """XLSX URL (spreadsheetml content-type): chunk_url returns a list."""
        import openpyxl

        from src.rag.document_loader import DocumentLoader
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Name', 'Value']); ws.append(['Alice', 42])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_resp(
            xlsx_bytes,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )):
            chunks = DocumentLoader().chunk_url('https://example.com/data.xlsx')
        assert isinstance(chunks, list)

    def test_url_remote_csv(self):
        """CSV URL (text/csv content-type): chunk_url returns a list of type 'csv' chunks."""
        from src.rag.document_loader import DocumentLoader
        csv_content = b'name,age\nAlice,30\nBob,25\n'
        with patch('requests.get', return_value=_mock_resp(csv_content, 'text/csv')):
            chunks = DocumentLoader().chunk_url('https://example.com/data.csv')
        assert isinstance(chunks, list)
        assert all(c['type'] == 'csv' for c in chunks)

    def test_url_remote_pptx(self):
        """PPTX URL (presentationml content-type): chunk_url returns a list."""
        from pptx import Presentation
        from pptx.util import Inches

        from src.rag.document_loader import DocumentLoader
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'Remote slide content here.'
        buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()
        with patch('requests.get', return_value=_mock_resp(
            pptx_bytes,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )):
            chunks = DocumentLoader().chunk_url('https://example.com/slides.pptx')
        assert isinstance(chunks, list)

    def test_url_type_by_content_type(self):
        """text/html content-type: all chunks have type 'html' (content-type priority 1)."""
        from src.rag.document_loader import DocumentLoader
        html = b'<html><body><p>Hello world. More content here now.</p></body></html>'
        with patch('requests.get', return_value=_mock_resp(html, 'text/html')):
            chunks = DocumentLoader().chunk_url('https://example.com')
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_type_by_extension(self):
        """octet-stream with .html extension in URL path: chunks have type 'html' (priority 2)."""
        from src.rag.document_loader import DocumentLoader
        html = b'<html><body><p>Content. More content. Third sentence here.</p></body></html>'
        with patch('requests.get', return_value=_mock_resp(html, 'application/octet-stream')):
            chunks = DocumentLoader().chunk_url('https://example.com/page.html')
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_type_by_pdf_magic_bytes(self):
        """octet-stream response with %PDF magic bytes: chunk_url returns a list (priority 3)."""
        from src.rag.document_loader import DocumentLoader
        fake_pdf = b'%PDF-1.4 fake content that is not real'
        with patch('requests.get', return_value=_mock_resp(fake_pdf, 'application/octet-stream')):
            chunks = DocumentLoader().chunk_url('https://example.com/nodoc')
        assert isinstance(chunks, list)

    def test_url_defaults_to_html(self):
        """octet-stream with no extension and no magic bytes: chunk_url defaults to html (priority 4)."""
        from src.rag.document_loader import DocumentLoader
        content = b'Hello world. This is a sentence. Third sentence here.'
        with patch('requests.get', return_value=_mock_resp(content, 'application/octet-stream')):
            chunks = DocumentLoader().chunk_url('https://example.com/nodoc')
        assert isinstance(chunks, list)

    def test_url_connection_error_returns_empty(self):
        """Network exception during requests.get: chunk_url returns an empty list."""
        from src.rag.document_loader import DocumentLoader
        with patch('requests.get', side_effect=Exception('timeout')):
            chunks = DocumentLoader().chunk_url('https://example.com')
        assert chunks == []

    def test_url_source_label_truncated_60_chars(self):
        """URL longer than 60 chars: source field in each chunk is at most 60 characters."""
        from src.rag.document_loader import DocumentLoader
        long_url = 'https://example.com/' + 'a' * 100
        html = b'<html><body><p>Text. More text. Here we go.</p></body></html>'
        with patch('requests.get', return_value=_mock_resp(html, 'text/html')):
            chunks = DocumentLoader().chunk_url(long_url)
        if chunks:
            assert len(chunks[0]['source']) <= 60


# ---------------------------------------------------------------------------
# chunk_directory
# ---------------------------------------------------------------------------

class TestChunkDirectory:
    """Integration tests for DocumentLoader.chunk_directory().

    Uses real temporary directories and real files so no mocking is needed.
    The method must handle: supported files, unsupported files, missing dirs.
    """

    def test_missing_directory_returns_empty_list(self):
        """Non-existent directory path: returns an empty list without raising."""
        from src.rag.document_loader import DocumentLoader
        chunks = DocumentLoader().chunk_directory('/tmp/this_path_does_not_exist_rag')
        assert chunks == []

    def test_empty_directory_returns_empty_list(self):
        """Directory with no files: returns an empty list."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            chunks = DocumentLoader().chunk_directory(tmpdir)
        assert chunks == []

    def test_txt_file_is_chunked(self):
        """A .txt file in the directory: returns at least one chunk."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            txt_path = os.path.join(tmpdir, 'test.txt')
            with open(txt_path, 'w') as f:
                f.write('Line one.\nLine two.\nLine three.\n')
            chunks = DocumentLoader().chunk_directory(tmpdir)
        assert len(chunks) > 0

    def test_csv_file_is_chunked(self):
        """A .csv file in the directory: returns at least one chunk."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            csv_path = os.path.join(tmpdir, 'data.csv')
            with open(csv_path, 'w', newline='') as f:
                writer = __import__('csv').DictWriter(f, fieldnames=['Name', 'Role'])
                writer.writeheader()
                writer.writerow({'Name': 'Alice', 'Role': 'Engineer'})
            chunks = DocumentLoader().chunk_directory(tmpdir)
        assert len(chunks) > 0

    def test_md_file_is_chunked(self):
        """A .md file in the directory: returns at least one chunk."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            md_path = os.path.join(tmpdir, 'notes.md')
            with open(md_path, 'w') as f:
                f.write('# Heading\nSome content here.\nMore text.\n')
            chunks = DocumentLoader().chunk_directory(tmpdir)
        assert len(chunks) > 0

    def test_unsupported_extension_is_skipped(self):
        """A .xyz file: skipped silently, returns empty list for that file."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            bad_path = os.path.join(tmpdir, 'unknown.xyz')
            with open(bad_path, 'w') as f:
                f.write('irrelevant content')
            chunks = DocumentLoader().chunk_directory(tmpdir)
        assert chunks == []

    def test_mixed_files_processes_only_supported(self):
        """Directory with one supported (.txt) and one unsupported (.xyz): only .txt chunked."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            txt_path = os.path.join(tmpdir, 'good.txt')
            bad_path = os.path.join(tmpdir, 'skip.xyz')
            with open(txt_path, 'w') as f:
                f.write('Hello world.\n')
            with open(bad_path, 'w') as f:
                f.write('ignored')
            chunks = DocumentLoader().chunk_directory(tmpdir)
        assert len(chunks) > 0
        # All chunks must come from the supported .txt file
        for chunk in chunks:
            assert 'good.txt' in chunk['source']

    def test_chunks_have_required_keys(self):
        """Each chunk returned by chunk_directory has the standard required keys."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            txt_path = os.path.join(tmpdir, 'sample.txt')
            with open(txt_path, 'w') as f:
                f.write('First line.\nSecond line.\n')
            chunks = DocumentLoader().chunk_directory(tmpdir)
        required_keys = {'text', 'source', 'start_line', 'end_line', 'type'}
        for chunk in chunks:
            missing = required_keys - set(chunk.keys())
            assert not missing, f"Chunk missing keys: {missing}"

    def test_subdirectories_are_skipped(self):
        """Subdirectories inside the target folder are not recursed into."""
        from src.rag.document_loader import DocumentLoader
        with tempfile.TemporaryDirectory() as tmpdir:
            # Create a file in a subdirectory — should be ignored
            subdir = os.path.join(tmpdir, 'sub')
            os.makedirs(subdir)
            nested_txt = os.path.join(subdir, 'nested.txt')
            with open(nested_txt, 'w') as f:
                f.write('nested content')
            chunks = DocumentLoader().chunk_directory(tmpdir)
        # No files at the top level → empty result
        assert chunks == []

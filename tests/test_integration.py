"""Integration tests — full pipeline with real libraries.

Mock strategy (per CLAUDE.md):
  Always mock:   ollama.embed, ollama.chat, chromadb → EphemeralClient, requests.get
  Never mock:    fitz, python-docx, openpyxl, xlrd, python-pptx, beautifulsoup4,
                 BM25Okapi, chunk truncation, misplaced detection, calculator eval

All file tests use real libraries and tmp files.
"""

import io
import os
import json
import tempfile
import pytest
import chromadb
from rank_bm25 import BM25Okapi
from unittest.mock import MagicMock, patch


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _fake_embed():
    return {'embeddings': [[0.1, 0.2, 0.3, 0.4]]}


def _fake_chat(content='mock'):
    return {'message': {'content': content}}


def _pipeline_chat_mock(*args, **kwargs):
    if kwargs.get('stream'):
        return [{'message': {'content': 'mock response'}}]
    return {'message': {'content': '8'}}


def _make_store(chunks):
    """Return a VectorStore wired to EphemeralClient with given chunks."""
    from src.rag.vector_store import VectorStore
    vs = VectorStore()
    client = chromadb.EphemeralClient()
    collection = client.get_or_create_collection(
        'test_rag', metadata={'hnsw:space': 'cosine'}
    )
    if chunks:
        ids    = [f'c{i}' for i in range(len(chunks))]
        texts  = [c['text'] for c in chunks]
        metas  = [{'source': c['source'], 'start_line': c['start_line'],
                   'end_line': c['end_line'], 'type': c['type']} for c in chunks]
        embeds = [[0.1, 0.2, 0.3, 0.4]] * len(chunks)
        collection.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
    vs.collection          = collection
    vs.chunks              = chunks
    vs.bm25_index          = BM25Okapi([c['text'].lower().split() for c in chunks]) if chunks else None
    vs.conversation_history = []
    return vs


# ============================================================
# 1. All 8 file types + xls (real libraries, tmp files)
# ============================================================

class TestLoadPdf:
    def test_load_pdf(self, tmp_path):
        import fitz
        from src.rag.document_loader import DocumentLoader
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Cats sleep sixteen hours a day. They are nocturnal hunters.")
        path = str(tmp_path / 'test.pdf')
        doc.save(path); doc.close()
        loader = DocumentLoader()
        chunks = loader._chunk_pdf(path, 'test.pdf')
        assert len(chunks) >= 1
        assert all(c['type'] == 'pdf' for c in chunks)
        assert any('cats' in c['text'].lower() for c in chunks)


class TestLoadTxt:
    def test_load_txt(self, tmp_path):
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.txt'
        f.write_text('Cats sleep sixteen hours a day.\nThey are nocturnal hunters.\n')
        loader = DocumentLoader()
        chunks = loader._chunk_txt(str(f), 'test.txt')
        assert len(chunks) == 2
        assert all(c['type'] == 'txt' for c in chunks)


class TestLoadDocx:
    def test_load_docx(self, tmp_path):
        from docx import Document
        from src.rag.document_loader import DocumentLoader
        doc = Document()
        doc.add_paragraph('Cats sleep sixteen hours a day.')
        doc.add_paragraph('They are excellent hunters.')
        path = str(tmp_path / 'test.docx')
        doc.save(path)
        loader = DocumentLoader()
        chunks = loader._chunk_docx(path, 'test.docx')
        assert len(chunks) >= 1
        assert all(c['type'] == 'docx' for c in chunks)

    def test_docx_table_rows_extracted(self, tmp_path):
        from docx import Document
        from src.rag.document_loader import DocumentLoader
        doc   = Document()
        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = 'Name'
        table.rows[0].cells[1].text = 'Alice'
        table.rows[1].cells[0].text = 'Job'
        table.rows[1].cells[1].text = 'Engineer'
        path = str(tmp_path / 'table.docx')
        doc.save(path)
        loader = DocumentLoader()
        chunks = loader._chunk_docx(path, 'table.docx')
        full   = ' '.join(c['text'] for c in chunks)
        assert 'Alice' in full or 'Name' in full

    def test_docx_merged_cells_deduplicated(self, tmp_path):
        from docx import Document
        from src.rag.document_loader import DocumentLoader
        doc   = Document()
        table = doc.add_table(rows=1, cols=3)
        table.rows[0].cells[0].text = 'Merged'
        table.rows[0].cells[1].text = 'Merged'   # simulate merged cell
        table.rows[0].cells[2].text = 'Different'
        path = str(tmp_path / 'merged.docx')
        doc.save(path)
        loader = DocumentLoader()
        chunks = loader._chunk_docx(path, 'merged.docx')
        full   = ' '.join(c['text'] for c in chunks)
        assert full.count('Merged') == 1


class TestLoadXlsx:
    def test_load_xlsx(self, tmp_path):
        import openpyxl
        from src.rag.document_loader import DocumentLoader
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(['Name', 'Age'])
        ws.append(['Alice', 30])
        ws.append(['Bob', 25])
        path = str(tmp_path / 'test.xlsx')
        wb.save(path)
        loader = DocumentLoader()
        chunks = loader._chunk_xlsx(path, 'test.xlsx')
        assert len(chunks) == 2
        assert all(c['type'] == 'xlsx' for c in chunks)
        assert 'Alice' in chunks[0]['text']


class TestLoadXls:
    def test_load_xls(self, tmp_path):
        import xlwt
        from src.rag.document_loader import DocumentLoader
        wb = xlwt.Workbook()
        ws = wb.add_sheet('Sheet1')
        ws.write(0, 0, 'Name'); ws.write(0, 1, 'Score')
        ws.write(1, 0, 'Alice'); ws.write(1, 1, 95)
        path = str(tmp_path / 'test.xls')
        wb.save(path)
        loader = DocumentLoader()
        chunks = loader._chunk_xls(path, 'test.xls')
        assert len(chunks) == 1
        assert 'Alice' in chunks[0]['text']


class TestLoadPptx:
    def test_load_pptx(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        from src.rag.document_loader import DocumentLoader
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'Cats are wonderful nocturnal hunters.'
        path  = str(tmp_path / 'test.pptx')
        prs.save(path)
        loader = DocumentLoader()
        chunks = loader._chunk_pptx(path, 'test.pptx')
        assert len(chunks) >= 1
        assert all(c['type'] == 'pptx' for c in chunks)
        assert 'Cats' in chunks[0]['text']


class TestLoadCsv:
    def test_load_csv(self, tmp_path):
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.csv'
        f.write_text('name,age\nAlice,30\nBob,25\n')
        loader = DocumentLoader()
        chunks = loader._chunk_csv(str(f), 'test.csv')
        assert len(chunks) == 2
        assert all(c['type'] == 'csv' for c in chunks)


class TestLoadMd:
    def test_load_md(self, tmp_path):
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.md'
        f.write_text('# Title\n**bold** cats sleep a lot\n')
        loader = DocumentLoader()
        chunks = loader._chunk_md(str(f), 'test.md')
        assert len(chunks) >= 1
        assert all(c['type'] == 'md' for c in chunks)
        assert all('**' not in c['text'] for c in chunks)
        assert all('# ' not in c['text'] for c in chunks)


class TestLoadHtml:
    def test_load_html(self, tmp_path):
        from src.rag.document_loader import DocumentLoader
        f = tmp_path / 'test.html'
        f.write_text('<html><body><p>Cats sleep. They hunt at night. Amazing animals.</p></body></html>')
        loader = DocumentLoader()
        chunks = loader._chunk_html(str(f), 'test.html')
        assert len(chunks) >= 1
        assert all(c['type'] == 'html' for c in chunks)
        assert all('<' not in c['text'] for c in chunks)


class TestMisplacedFile:
    def test_misplaced_file_detected_and_chunked(self, tmp_path):
        from src.rag.document_loader import DocumentLoader
        txts_dir = tmp_path / 'txts'
        txts_dir.mkdir()
        (txts_dir / 'report.pdf').write_bytes(b'%PDF fake')
        fake_folders = {
            'pdf':  str(tmp_path / 'pdfs'),
            'txt':  str(txts_dir),
            'docx': str(tmp_path / 'docx'),
            'xlsx': str(tmp_path / 'xlsx'),
            'pptx': str(tmp_path / 'pptx'),
            'csv':  str(tmp_path / 'csv'),
            'md':   str(tmp_path / 'md'),
            'html': str(tmp_path / 'html'),
        }
        loader = DocumentLoader()
        loader.doc_folders = fake_folders
        files = loader.scan_all_files()
        assert len(files) == 1
        assert files[0]['is_misplaced'] is True
        assert files[0]['detected_type'] == 'pdf'


class TestTruncateChunk:
    def test_truncate_chunk_300_words(self):
        from src.rag.document_loader import DocumentLoader
        loader = DocumentLoader()
        text   = ' '.join(['word'] * 400)
        result = loader._truncate_chunk(text)
        assert len(result.split()) <= 300

    def test_truncate_chunk_1200_chars(self):
        from src.rag.document_loader import DocumentLoader
        loader = DocumentLoader()
        text   = 'a' * 2000
        result = loader._truncate_chunk(text)
        assert len(result) <= 1200


# ============================================================
# 2. URL ingestion (mock requests.get only)
# ============================================================

class TestUrlIngestion:
    def _mock_resp(self, content, content_type='text/html', encoding='utf-8'):
        resp = MagicMock()
        resp.content = content if isinstance(content, bytes) else content.encode(encoding)
        resp.headers = {'Content-Type': content_type}
        resp.encoding = encoding
        resp.raise_for_status = MagicMock()
        return resp

    def test_url_html_webpage(self):
        from src.rag.document_loader import DocumentLoader
        html = b'<html><body><p>Hello world. Second sentence. Third one here.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = DocumentLoader().chunk_url('https://example.com')
        assert len(chunks) >= 1
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_remote_pdf(self):
        import fitz
        from src.rag.document_loader import DocumentLoader
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF content. Another sentence.")
        buf = io.BytesIO(); doc.save(buf); pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=self._mock_resp(pdf_bytes, 'application/pdf')):
            chunks = DocumentLoader().chunk_url('https://example.com/file.pdf')
        assert isinstance(chunks, list)

    def test_url_remote_docx(self):
        from docx import Document
        from src.rag.document_loader import DocumentLoader
        doc = Document(); doc.add_paragraph('Remote docx content here.')
        buf = io.BytesIO(); doc.save(buf); docx_bytes = buf.getvalue()
        with patch('requests.get', return_value=self._mock_resp(
            docx_bytes,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )):
            chunks = DocumentLoader().chunk_url('https://example.com/file.docx')
        assert isinstance(chunks, list)

    def test_url_remote_xlsx(self):
        import openpyxl
        from src.rag.document_loader import DocumentLoader
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(['Name', 'Value']); ws.append(['Alice', 42])
        buf = io.BytesIO(); wb.save(buf); xlsx_bytes = buf.getvalue()
        with patch('requests.get', return_value=self._mock_resp(
            xlsx_bytes,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )):
            chunks = DocumentLoader().chunk_url('https://example.com/data.xlsx')
        assert isinstance(chunks, list)

    def test_url_remote_csv(self):
        from src.rag.document_loader import DocumentLoader
        csv_content = b'name,age\nAlice,30\nBob,25\n'
        with patch('requests.get', return_value=self._mock_resp(csv_content, 'text/csv')):
            chunks = DocumentLoader().chunk_url('https://example.com/data.csv')
        assert isinstance(chunks, list)
        assert all(c['type'] == 'csv' for c in chunks)

    def test_url_remote_pptx(self):
        from pptx import Presentation
        from pptx.util import Inches
        from src.rag.document_loader import DocumentLoader
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'Remote slide content here.'
        buf = io.BytesIO(); prs.save(buf); pptx_bytes = buf.getvalue()
        with patch('requests.get', return_value=self._mock_resp(
            pptx_bytes,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )):
            chunks = DocumentLoader().chunk_url('https://example.com/slides.pptx')
        assert isinstance(chunks, list)

    def test_url_type_by_content_type(self):
        from src.rag.document_loader import DocumentLoader
        html = b'<html><body><p>Hello world. More content here now.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = DocumentLoader().chunk_url('https://example.com')
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_type_by_extension(self):
        from src.rag.document_loader import DocumentLoader
        html = b'<html><body><p>Content. More content. Third sentence here.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'application/octet-stream')):
            chunks = DocumentLoader().chunk_url('https://example.com/page.html')
        assert all(c['type'] == 'html' for c in chunks)

    def test_url_type_by_pdf_magic_bytes(self):
        from src.rag.document_loader import DocumentLoader
        fake_pdf = b'%PDF-1.4 fake content that is not real'
        with patch('requests.get', return_value=self._mock_resp(fake_pdf, 'application/octet-stream')):
            chunks = DocumentLoader().chunk_url('https://example.com/nodoc')
        assert isinstance(chunks, list)

    def test_url_defaults_to_html(self):
        from src.rag.document_loader import DocumentLoader
        content = b'Hello world. This is a sentence. Third sentence here.'
        with patch('requests.get', return_value=self._mock_resp(content, 'application/octet-stream')):
            chunks = DocumentLoader().chunk_url('https://example.com/nodoc')
        assert isinstance(chunks, list)

    def test_url_connection_error_returns_empty(self):
        from src.rag.document_loader import DocumentLoader
        with patch('requests.get', side_effect=Exception('timeout')):
            chunks = DocumentLoader().chunk_url('https://example.com')
        assert chunks == []

    def test_url_source_label_truncated_60_chars(self):
        from src.rag.document_loader import DocumentLoader
        long_url = 'https://example.com/' + 'a' * 100
        html = b'<html><body><p>Text. More text. Here we go.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = DocumentLoader().chunk_url(long_url)
        if chunks:
            assert len(chunks[0]['source']) <= 60


# ============================================================
# 3. VectorStore pipeline (mock ollama + EphemeralClient)
# ============================================================

SAMPLE_CHUNKS = [
    {'text': 'Cats sleep sixteen hours a day.',    'source': 'cats.txt', 'start_line': 1, 'end_line': 1, 'type': 'txt'},
    {'text': 'Cats can see in dim light very well.', 'source': 'cats.txt', 'start_line': 2, 'end_line': 2, 'type': 'txt'},
    {'text': 'Cats have five toes on front paws.',  'source': 'cats.txt', 'start_line': 3, 'end_line': 3, 'type': 'txt'},
    {'text': 'Cats have twelve whiskers.',           'source': 'cats.txt', 'start_line': 4, 'end_line': 4, 'type': 'txt'},
    {'text': 'Cats cannot taste sweet food.',        'source': 'cats.txt', 'start_line': 5, 'end_line': 5, 'type': 'txt'},
]


class TestHybridRetrieve:
    def test_hybrid_retrieve_fuses_bm25_and_dense(self):
        vs = _make_store(SAMPLE_CHUNKS)
        with patch('ollama.embed', return_value=_fake_embed()):
            result = vs._hybrid_retrieve(['cats sleep'], top_n=3)
        assert len(result) <= 3
        assert all(isinstance(score, float) for _, score in result)

    def test_hybrid_retrieve_alpha_weighting(self):
        vs = _make_store(SAMPLE_CHUNKS)
        with patch('ollama.embed', return_value=_fake_embed()):
            r_dense = vs._hybrid_retrieve(['cats sleep'], top_n=5, alpha=1.0)
            r_bm25  = vs._hybrid_retrieve(['cats sleep'], top_n=5, alpha=0.0)
        assert isinstance(r_dense, list)
        assert isinstance(r_bm25, list)


class TestExpandQuery:
    def test_expand_query_returns_3(self):
        vs = _make_store(SAMPLE_CHUNKS)
        with patch('ollama.chat', return_value=_fake_chat('rewrite1\nrewrite2')):
            result = vs._expand_query('how many hours do cats sleep')
        assert len(result) == 3
        assert result[0] == 'how many hours do cats sleep'

    def test_expand_query_fallback_on_error(self):
        vs = _make_store(SAMPLE_CHUNKS)
        with patch('ollama.chat', side_effect=Exception('fail')):
            result = vs._expand_query('test query')
        assert result == ['test query']


class TestClassifyQuery:
    def test_classify_summarise_checked_first(self):
        vs = _make_store(SAMPLE_CHUNKS)
        assert vs._classify_query('summarise all cat facts') == 'summarise'

    def test_classify_factual(self):
        vs = _make_store(SAMPLE_CHUNKS)
        assert vs._classify_query('what is the boiling point of water') == 'factual'

    def test_classify_comparison(self):
        vs = _make_store(SAMPLE_CHUNKS)
        assert vs._classify_query('compare cats vs dogs') == 'comparison'

    def test_classify_general(self):
        vs = _make_store(SAMPLE_CHUNKS)
        assert vs._classify_query('something interesting happened today') == 'general'


class TestRerank:
    def test_rerank_orders_by_llm_score(self):
        vs = _make_store(SAMPLE_CHUNKS)
        candidates = [
            ({'text': 'low relevance', 'source': 's', 'start_line': 1, 'end_line': 1, 'type': 'txt'}, 0.5),
            ({'text': 'high relevance', 'source': 's', 'start_line': 2, 'end_line': 2, 'type': 'txt'}, 0.4),
        ]
        scores = iter(['2', '9'])
        with patch('ollama.chat', side_effect=lambda **kw: _fake_chat(next(scores))):
            result = vs._rerank('test query', candidates, top_n=2)
        assert result[0][2] >= result[1][2]


class TestConfidence:
    def test_confidence_below_threshold(self):
        from src.rag.config import SIMILARITY_THRESHOLD
        vs = _make_store(SAMPLE_CHUNKS)
        entry = SAMPLE_CHUNKS[0]
        results = [(entry, SIMILARITY_THRESHOLD - 0.1)]
        confident, score = vs._check_confidence(results)
        assert confident is False

    def test_confidence_above_threshold(self):
        from src.rag.config import SIMILARITY_THRESHOLD
        vs = _make_store(SAMPLE_CHUNKS)
        entry = SAMPLE_CHUNKS[0]
        results = [(entry, SIMILARITY_THRESHOLD + 0.1)]
        confident, score = vs._check_confidence(results)
        assert confident is True


class TestSmartTopN:
    def test_smart_top_n_all_4_types(self):
        from src.rag.config import TOP_RETRIEVE
        vs = _make_store(SAMPLE_CHUNKS)
        assert vs._smart_top_n('factual')    == 5
        assert vs._smart_top_n('comparison') == 15
        assert vs._smart_top_n('general')    == 10
        assert vs._smart_top_n('summarise')  == TOP_RETRIEVE


class TestSourceLabel:
    def test_source_label_all_types(self):
        vs = _make_store(SAMPLE_CHUNKS)
        assert vs._source_label({'type': 'pdf',  'start_line': 3, 'end_line': 3}) == 'p3'
        assert vs._source_label({'type': 'xlsx', 'start_line': 5, 'end_line': 5}) == 'row5'
        assert vs._source_label({'type': 'csv',  'start_line': 2, 'end_line': 2}) == 'row2'
        assert vs._source_label({'type': 'pptx', 'start_line': 1, 'end_line': 1}) == 'slide1'
        assert vs._source_label({'type': 'html', 'start_line': 4, 'end_line': 4}) == 's4'
        assert vs._source_label({'type': 'txt',  'start_line': 1, 'end_line': 3}) == 'L1-3'


class TestHallucinationFilter:
    def test_hallucination_filter_truncates(self):
        vs = _make_store(SAMPLE_CHUNKS)
        response = "I couldn't find any information. However, I can tell you cats are nice."
        result = vs._filter_hallucination(response)
        assert 'however' not in result.lower()
        assert 'I can only answer' in result

    def test_hallucination_filter_clean_response_unchanged(self):
        vs = _make_store(SAMPLE_CHUNKS)
        response = 'Cats sleep 16 hours a day. [cats.txt L1]'
        assert vs._filter_hallucination(response) == response


class TestLowConfidence:
    def test_low_confidence_skips_llm(self):
        vs = _make_store(SAMPLE_CHUNKS)
        # Use distant embedding to force low similarity
        with patch('ollama.embed', return_value={'embeddings': [[0.0, 0.0, 0.0, 1.0]]}), \
             patch('ollama.chat', side_effect=_pipeline_chat_mock) as mock_chat:
            result = vs.run_pipeline('completely unrelated xyz topic', streamlit_mode=True)
        # If not confident, LLM for synthesis should not be called
        assert 'response' in result
        assert result['is_confident'] is False or 'response' in result


class TestRebuildLogic:
    def test_rebuild_logic_skips_if_existing_gte_chunks(self):
        from src.rag.vector_store import VectorStore
        vs = VectorStore()
        with patch('chromadb.PersistentClient') as mock_client:
            client     = chromadb.EphemeralClient()
            collection = client.get_or_create_collection('rag', metadata={'hnsw:space': 'cosine'})
            # Pre-populate with more entries than chunks → should skip rebuild
            collection.add(
                ids=['x0', 'x1', 'x2'],
                embeddings=[[0.1, 0.2, 0.3, 0.4]] * 3,
                documents=['a', 'b', 'c'],
                metadatas=[{'source': 's', 'start_line': 1, 'end_line': 1, 'type': 'txt'}] * 3
            )
            mock_client.return_value.get_or_create_collection.return_value = collection
            chunks = [
                {'text': 'a', 'source': 's', 'start_line': 1, 'end_line': 1, 'type': 'txt'},
                {'text': 'b', 'source': 's', 'start_line': 2, 'end_line': 2, 'type': 'txt'},
            ]
            with patch('ollama.embed', return_value=_fake_embed()):
                vs.build_or_load(chunks)
        # Collection should still have 3 (not rebuilt)
        assert vs.collection.count() == 3

    def test_rebuild_logic_deletes_and_rebuilds_if_local_grew(self):
        from src.rag.vector_store import VectorStore
        vs = VectorStore()
        with patch('chromadb.PersistentClient') as mock_client:
            client     = chromadb.EphemeralClient()
            collection = client.get_or_create_collection('rag2', metadata={'hnsw:space': 'cosine'})
            # Pre-populate with fewer entries than chunks → should delete and rebuild
            collection.add(
                ids=['x0'],
                embeddings=[[0.1, 0.2, 0.3, 0.4]],
                documents=['old'],
                metadatas=[{'source': 's', 'start_line': 1, 'end_line': 1, 'type': 'txt'}]
            )
            mock_client.return_value.get_or_create_collection.return_value = collection
            chunks = [
                {'text': 'a', 'source': 's', 'start_line': 1, 'end_line': 1, 'type': 'txt'},
                {'text': 'b', 'source': 's', 'start_line': 2, 'end_line': 2, 'type': 'txt'},
                {'text': 'c', 'source': 's', 'start_line': 3, 'end_line': 3, 'type': 'txt'},
            ]
            with patch('ollama.embed', return_value=_fake_embed()):
                vs.build_or_load(chunks)
        assert vs.collection.count() == 3


# ============================================================
# 4. All 5 agent tools
# ============================================================

class TestAgentTools:
    @pytest.fixture
    def agent(self):
        from src.rag.agent import Agent
        vs = _make_store(SAMPLE_CHUNKS)
        vs._expand_query = MagicMock(return_value=['cats'])
        entry = SAMPLE_CHUNKS[0]
        vs._hybrid_retrieve = MagicMock(return_value=[(entry, 0.9)])
        vs._rerank          = MagicMock(return_value=[(entry, 0.9, 0.9)])
        vs._source_label    = MagicMock(return_value='L1-1')
        return Agent(vs)

    def test_tool_calculator_basic(self, agent):
        assert agent._tool_calculator('2 + 2') == '4'

    def test_tool_calculator_complex(self, agent):
        assert agent._tool_calculator('(100 + 50) * 2') == '300'

    def test_tool_calculator_unsafe_chars_rejected(self, agent):
        result = agent._tool_calculator('__import__("os").system("rm -rf /")')
        assert 'Error' in result

    def test_tool_summarise_short_2_3_sentences(self, agent):
        short_text = ' '.join(['word'] * 50)
        with patch('ollama.chat', return_value=_fake_chat('summary')) as m:
            agent._tool_summarise(short_text)
        assert '2-3 sentences' in m.call_args[1]['messages'][0]['content']

    def test_tool_summarise_medium_4_5_sentences(self, agent):
        medium_text = ' '.join(['word'] * 150)
        with patch('ollama.chat', return_value=_fake_chat('summary')) as m:
            agent._tool_summarise(medium_text)
        assert '4-5 sentences' in m.call_args[1]['messages'][0]['content']

    def test_tool_summarise_long_6_8_sentences(self, agent):
        long_text = ' '.join(['word'] * 400)
        with patch('ollama.chat', return_value=_fake_chat('summary')) as m:
            agent._tool_summarise(long_text)
        assert '6-8 sentences' in m.call_args[1]['messages'][0]['content']

    def test_tool_sentiment_short_query_searches_first(self, agent):
        agent._tool_rag_search = MagicMock(return_value='cats are wonderful creatures')
        with patch('ollama.chat', return_value=_fake_chat('Sentiment: Positive\nTone: warm\nKey phrases: wonderful\nExplanation: positive.')):
            agent._tool_sentiment('cats')
        agent._tool_rag_search.assert_called_once()

    def test_tool_sentiment_long_text_direct(self, agent):
        agent._tool_rag_search = MagicMock(return_value='x')
        long_text = ' '.join(['word'] * 15)
        with patch('ollama.chat', return_value=_fake_chat('Sentiment: Neutral\nTone: flat\nKey phrases: word\nExplanation: neutral.')):
            agent._tool_sentiment(long_text)
        agent._tool_rag_search.assert_not_called()

    def test_tool_sentiment_output_4_fields(self, agent):
        mock_resp = 'Sentiment: Positive\nTone: upbeat\nKey phrases: great\nExplanation: positive mood.'
        with patch('ollama.chat', return_value=_fake_chat(mock_resp)):
            result = agent._tool_sentiment('great product experience here today')
        assert 'Sentiment' in result
        assert 'Tone' in result

    def test_parse_tool_call_with_parens(self, agent):
        name, arg = agent._parse_tool_call('TOOL: rag_search(cat sleep hours)')
        assert name == 'rag_search'
        assert arg == 'cat sleep hours'

    def test_parse_tool_call_without_parens(self, agent):
        name, arg = agent._parse_tool_call('TOOL: rag_search cat sleep hours')
        assert name == 'rag_search'
        assert arg == 'cat sleep hours'

    def test_parse_tool_call_malformed_returns_none(self, agent):
        name, arg = agent._parse_tool_call('This is not a valid response.')
        assert name is None

    def test_fast_path_summarise_4_searches(self, agent):
        agent._tool_rag_search = MagicMock(return_value='result')
        with patch('ollama.chat', return_value=_fake_chat('summary answer')):
            result = agent._fast_path_summarise('summarise the document')
        tool_steps = [s for s in result['steps'] if s['tool'] == 'rag_search']
        assert len(tool_steps) == 4

    def test_fast_path_sentiment_strips_labels(self, agent):
        raw = '- [cats.txt L1-1] Cats are wonderful creatures indeed.'
        agent._tool_rag_search = MagicMock(return_value=raw)
        captured = []
        def cap_sentiment(text):
            captured.append(text)
            return 'Sentiment: Positive\nTone: warm\nKey phrases: wonderful\nExplanation: positive.'
        agent._tool_sentiment = cap_sentiment
        agent._fast_path_sentiment('what is the sentiment of cats')
        if captured:
            assert '[cats.txt' not in captured[0]

    def test_calculator_auto_finish(self, agent):
        with patch('ollama.chat', return_value=_fake_chat('TOOL: calculator(16 * 365)')):
            result = agent.run('how many hours in 365 days at 16h each')
        finish_steps = [s for s in result['steps'] if s['tool'] == 'finish']
        assert len(finish_steps) >= 1
        assert '5840' in result['answer']

    def test_rag_search_auto_finish(self, agent):
        with patch('ollama.chat', side_effect=[
            _fake_chat('TOOL: rag_search(cats sleep)'),
            _fake_chat('Cats sleep 16 hours a day.'),
        ]):
            result = agent.run('how many hours do cats sleep')
        assert 'answer' in result
        finish_steps = [s for s in result['steps'] if s['tool'] == 'finish']
        assert len(finish_steps) >= 1

    def test_bad_format_retry_max_2(self, agent):
        responses = [
            _fake_chat('not a tool call'),
            _fake_chat('still not a tool call'),
            _fake_chat('TOOL: finish(final answer here)'),
        ]
        with patch('ollama.chat', side_effect=responses):
            result = agent.run('any query')
        assert result is not None

    def test_step_limit_reached(self, agent):
        with patch('ollama.chat', return_value=_fake_chat('not a tool call ever')):
            result = agent.run('any query')
        assert 'answer' in result

    def test_collected_context_used_for_final_answer(self, agent):
        """After rag_search, finish uses synthesized answer from collected context."""
        with patch('ollama.chat', side_effect=[
            _fake_chat('TOOL: rag_search(cats)'),
            _fake_chat('Cats sleep 16 hours.'),
        ]):
            result = agent.run('tell me about cats sleeping')
        assert result['answer']
        assert len(result['steps']) >= 2

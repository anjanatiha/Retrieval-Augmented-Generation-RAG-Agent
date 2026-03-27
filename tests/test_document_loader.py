"""test_document_loader.py — Unit and integration tests for DocumentLoader.

Covers:
  - DocumentLoader.__init__ (state initialisation)
  - All 9 chunker functions from src.rag.chunkers (txt, md, pdf, docx, xlsx, xls, csv, pptx, html)
  - truncate_chunk from src.rag.chunkers
  - DocumentLoader.scan_all_files (misplaced file detection)
  - DocumentLoader.chunk_url (4-priority type detection, URL ingestion)

Mock strategy:
  - requests.get is mocked only in TestChunkUrl.
  - Never mock: fitz, python-docx, openpyxl, xlrd, python-pptx, beautifulsoup4.
  - Real libraries are always used for file format tests.
"""

import csv
import io
import json
import os
import tempfile
import textwrap
import pytest
from unittest.mock import MagicMock, patch

# ── Chunker functions are now standalone module-level functions ───────────────
from src.rag.chunkers import (
    chunk_txt, chunk_md, chunk_pdf, chunk_docx,
    chunk_xlsx, chunk_xls, chunk_csv, chunk_pptx,
    chunk_html, truncate_chunk,
)


# ---------------------------------------------------------------------------
# Fixture
# ---------------------------------------------------------------------------

@pytest.fixture
def loader():
    """Provide a fresh DocumentLoader instance for each test.

    Used by tests that exercise DocumentLoader state (init, scan_all_files,
    chunk_url, _dispatch_chunker). Chunker function tests use the module
    functions directly from src.rag.chunkers.
    """
    from src.rag.document_loader import DocumentLoader
    return DocumentLoader()


# ---------------------------------------------------------------------------
# __init__
# ---------------------------------------------------------------------------

class TestInit:
    """Tests that DocumentLoader initialises its state correctly from config."""

    def test_has_doc_folders(self, loader):
        """doc_folders attribute matches the DOC_FOLDERS constant from config."""
        from src.rag.config import DOC_FOLDERS
        assert loader.doc_folders == DOC_FOLDERS

    def test_has_ext_to_type(self, loader):
        """ext_to_type attribute matches the EXT_TO_TYPE constant from config."""
        from src.rag.config import EXT_TO_TYPE
        assert loader.ext_to_type == EXT_TO_TYPE

    def test_has_chunk_sizes(self, loader):
        """chunk_sizes attribute exists and contains the txt_chunk_size key."""
        assert hasattr(loader, 'chunk_sizes')
        assert 'txt_chunk_size' in loader.chunk_sizes


# ---------------------------------------------------------------------------
# truncate_chunk
# ---------------------------------------------------------------------------

class TestTruncateChunk:
    """Tests for truncate_chunk: 300-word and 1200-char limits (whichever shorter)."""

    def test_short_text_unchanged(self):
        """Short text well under both limits is returned unchanged."""
        text = "hello world"
        assert truncate_chunk(text) == text

    def test_truncates_at_300_words(self):
        """Text with 400 words is truncated so the result has at most 300 words."""
        text = ' '.join(['word'] * 400)
        result = truncate_chunk(text)
        assert len(result.split()) <= 300

    def test_truncates_at_1200_chars(self):
        """Text with 2000 identical chars is truncated to at most 1200 chars."""
        text = 'a' * 2000
        result = truncate_chunk(text)
        assert len(result) <= 1200

    def test_word_limit_takes_priority_over_char_if_shorter(self):
        """100 short words (600 chars total) hit neither limit and are returned as-is."""
        # 100 words each 5 chars → 600 chars total — neither limit is hit
        text = ' '.join(['hello'] * 100)
        result = truncate_chunk(text)
        assert result == text

    def test_char_limit_takes_priority_when_words_short_but_chars_long(self):
        """50 long tokens (≈1550 chars, under 300 words) are truncated at 1200 chars."""
        # 50 * 31 chars ≈ 1550 chars, but only 50 words — under word limit but over char limit
        text = ' '.join(['a' * 30] * 50)
        result = truncate_chunk(text)
        assert len(result) <= 1200


# ---------------------------------------------------------------------------
# chunk_txt
# ---------------------------------------------------------------------------

class TestChunkTxt:
    """Tests for chunk_txt: line-based chunking of plain-text files."""

    def test_basic_chunking(self, tmp_path):
        """Three non-empty lines produce three chunks with correct type and source."""
        f = tmp_path / 'test.txt'
        f.write_text('line one\nline two\nline three\n')
        chunks = chunk_txt(str(f), 'test.txt')
        assert len(chunks) == 3
        assert all(c['type'] == 'txt' for c in chunks)
        assert all(c['source'] == 'test.txt' for c in chunks)

    def test_empty_lines_skipped(self, tmp_path):
        """Blank lines between content lines are ignored; only two chunks returned."""
        f = tmp_path / 'test.txt'
        f.write_text('line one\n\n\nline two\n')
        chunks = chunk_txt(str(f), 'test.txt')
        assert len(chunks) == 2

    def test_chunk_has_start_end_line(self, tmp_path):
        """Each chunk carries a start_line key; first chunk starts at line 1."""
        f = tmp_path / 'test.txt'
        f.write_text('a\nb\nc\n')
        chunks = chunk_txt(str(f), 'test.txt')
        assert chunks[0]['start_line'] == 1


# ---------------------------------------------------------------------------
# chunk_md
# ---------------------------------------------------------------------------

class TestChunkMd:
    """Tests for chunk_md: Markdown chunking with markup stripping."""

    def test_strips_headings(self, tmp_path):
        """Heading markers (# ) are removed from chunk text."""
        f = tmp_path / 'test.md'
        f.write_text('# Title\nsome content\n')
        chunks = chunk_md(str(f), 'test.md')
        assert all('# ' not in c['text'] for c in chunks)

    def test_strips_bold(self, tmp_path):
        """Bold markers (**) are removed from chunk text."""
        f = tmp_path / 'test.md'
        f.write_text('**bold text** here\n')
        chunks = chunk_md(str(f), 'test.md')
        assert all('**' not in c['text'] for c in chunks)

    def test_type_is_md(self, tmp_path):
        """All chunks from a .md file have type set to 'md'."""
        f = tmp_path / 'test.md'
        f.write_text('some markdown content\n')
        chunks = chunk_md(str(f), 'test.md')
        assert all(c['type'] == 'md' for c in chunks)


# ---------------------------------------------------------------------------
# chunk_pdf
# ---------------------------------------------------------------------------

class TestChunkPdf:
    """Tests for chunk_pdf: page-level sentence-window chunking via PyMuPDF."""

    def test_returns_list(self, tmp_path):
        """A valid single-page PDF produces a non-empty list of chunks."""
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Hello world. This is a PDF. It has sentences.")
        path = str(tmp_path / 'test.pdf')
        doc.save(path)
        doc.close()
        chunks = chunk_pdf(path, 'test.pdf')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1

    def test_type_is_pdf(self, tmp_path):
        """All chunks from a PDF file have type set to 'pdf'."""
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Test sentence. Another sentence.")
        path = str(tmp_path / 'test.pdf')
        doc.save(path)
        doc.close()
        chunks = chunk_pdf(path, 'test.pdf')
        assert all(c['type'] == 'pdf' for c in chunks)

    def test_bad_file_returns_empty(self, tmp_path):
        """A corrupted or non-PDF file returns an empty list without raising."""
        path = str(tmp_path / 'bad.pdf')
        with open(path, 'w') as f:
            f.write('not a pdf')
        chunks = chunk_pdf(path, 'bad.pdf')
        assert chunks == []


# ---------------------------------------------------------------------------
# chunk_docx
# ---------------------------------------------------------------------------

class TestChunkDocx:
    """Tests for chunk_docx: paragraph chunking and table row extraction."""

    def _make_docx(self, tmp_path, paragraphs=None, table_rows=None):
        """Create a temporary .docx file with optional paragraphs and table rows."""
        from docx import Document
        doc = Document()
        for p in (paragraphs or []):
            doc.add_paragraph(p)
        if table_rows:
            table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
            for i, row_data in enumerate(table_rows):
                for j, cell_text in enumerate(row_data):
                    table.rows[i].cells[j].text = cell_text
        path = str(tmp_path / 'test.docx')
        doc.save(path)
        return path

    def test_basic_paragraphs(self, tmp_path):
        """Three paragraphs produce at least one chunk of type 'docx'."""
        path = self._make_docx(tmp_path, paragraphs=['Para one', 'Para two', 'Para three'])
        chunks = chunk_docx(path, 'test.docx')
        assert len(chunks) >= 1
        assert all(c['type'] == 'docx' for c in chunks)

    def test_table_rows_extracted(self, tmp_path):
        """Table cell values are included in chunk text."""
        path = self._make_docx(tmp_path, table_rows=[['Name', 'Alice'], ['Job', 'Engineer']])
        chunks = chunk_docx(path, 'test.docx')
        full_text = ' '.join(c['text'] for c in chunks)
        assert 'Alice' in full_text or 'Name' in full_text

    def test_merged_cells_deduplicated(self, tmp_path):
        """Repeated text from merged cells appears only once in the chunk output."""
        from docx import Document
        doc = Document()
        table = doc.add_table(rows=1, cols=3)
        table.rows[0].cells[0].text = 'Merged'
        table.rows[0].cells[1].text = 'Merged'   # simulate merged cell repeat
        table.rows[0].cells[2].text = 'Different'
        path = str(tmp_path / 'merged.docx')
        doc.save(path)
        chunks = chunk_docx(path, 'merged.docx')
        full_text = ' '.join(c['text'] for c in chunks)
        # Deduplication means "Merged" appears only once in the table row chunk
        assert full_text.count('Merged') == 1


# ---------------------------------------------------------------------------
# chunk_xlsx
# ---------------------------------------------------------------------------

class TestChunkXlsx:
    """Tests for chunk_xlsx: key=value pair chunking of Excel .xlsx files."""

    def test_row_chunks(self, tmp_path):
        """Two data rows (after header) produce two chunks containing cell values."""
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(['Name', 'Age'])
        ws.append(['Alice', 30])
        ws.append(['Bob', 25])
        path = str(tmp_path / 'test.xlsx')
        wb.save(path)
        chunks = chunk_xlsx(path, 'test.xlsx')
        assert len(chunks) == 2  # 2 data rows
        assert all(c['type'] == 'xlsx' for c in chunks)
        assert 'Alice' in chunks[0]['text']

    def test_empty_rows_skipped(self, tmp_path):
        """Rows where all cells are None are skipped; only one valid chunk returned."""
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(['Name', 'Age'])
        ws.append([None, None])
        ws.append(['Bob', 25])
        path = str(tmp_path / 'test.xlsx')
        wb.save(path)
        chunks = chunk_xlsx(path, 'test.xlsx')
        assert len(chunks) == 1


# ---------------------------------------------------------------------------
# chunk_xls
# ---------------------------------------------------------------------------

class TestChunkXls:
    """Tests for chunk_xls: legacy .xls file chunking via xlrd."""

    def test_xls_row_chunks(self, tmp_path):
        """One data row in a .xls file produces one chunk containing the cell value."""
        import xlwt
        wb = xlwt.Workbook()
        ws = wb.add_sheet('Sheet1')
        ws.write(0, 0, 'Name'); ws.write(0, 1, 'Score')
        ws.write(1, 0, 'Alice'); ws.write(1, 1, 95)
        path = str(tmp_path / 'test.xls')
        wb.save(path)
        chunks = chunk_xls(path, 'test.xls')
        assert len(chunks) == 1
        assert 'Alice' in chunks[0]['text']


# ---------------------------------------------------------------------------
# chunk_csv
# ---------------------------------------------------------------------------

class TestChunkCsv:
    """Tests for chunk_csv: DictReader key=value pair chunking of CSV files."""

    def test_csv_row_chunks(self, tmp_path):
        """Two data rows produce two chunks of type 'csv' with correct content."""
        f = tmp_path / 'test.csv'
        f.write_text('name,age\nAlice,30\nBob,25\n')
        chunks = chunk_csv(str(f), 'test.csv')
        assert len(chunks) == 2
        assert all(c['type'] == 'csv' for c in chunks)
        assert 'Alice' in chunks[0]['text']

    def test_empty_rows_skipped(self, tmp_path):
        """A row with all empty values is skipped; only one valid chunk returned."""
        f = tmp_path / 'test.csv'
        f.write_text('name,age\n,\nAlice,30\n')
        chunks = chunk_csv(str(f), 'test.csv')
        assert len(chunks) == 1


# ---------------------------------------------------------------------------
# chunk_pptx
# ---------------------------------------------------------------------------

class TestChunkPptx:
    """Tests for chunk_pptx: per-slide text extraction from PowerPoint files."""

    def test_slide_chunks(self, tmp_path):
        """A single slide with a textbox produces one chunk containing its text."""
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = 'Hello from slide one'
        path = str(tmp_path / 'test.pptx')
        prs.save(path)
        chunks = chunk_pptx(path, 'test.pptx')
        assert len(chunks) >= 1
        assert all(c['type'] == 'pptx' for c in chunks)
        assert 'Hello from slide one' in chunks[0]['text']


# ---------------------------------------------------------------------------
# chunk_html
# ---------------------------------------------------------------------------

class TestChunkHtml:
    """Tests for chunk_html: BeautifulSoup tag stripping and sentence windowing."""

    def test_strips_tags(self, tmp_path):
        """HTML tags are stripped; chunks contain plain text of type 'html'."""
        f = tmp_path / 'test.html'
        f.write_text('<html><body><p>Hello world.</p><p>Second sentence.</p></body></html>')
        chunks = chunk_html(str(f), 'test.html')
        assert len(chunks) >= 1
        assert all('<' not in c['text'] for c in chunks)
        assert all(c['type'] == 'html' for c in chunks)


# ---------------------------------------------------------------------------
# scan_all_files — misplaced detection
# ---------------------------------------------------------------------------

class TestScanAllFiles:
    """Tests for scan_all_files: detection of files placed in the wrong type folder."""

    def test_misplaced_file_detected(self, loader, tmp_path):
        """A .pdf file in the txts/ folder is flagged as misplaced with detected_type 'pdf'."""
        # Put a .pdf file in the txts/ folder to simulate a user putting it in the wrong place
        txts_dir = tmp_path / 'txts'
        txts_dir.mkdir(parents=True)
        (txts_dir / 'report.pdf').write_bytes(b'%PDF fake')

        fake_folders = {
            'pdf':  str(tmp_path / 'pdfs'),
            'txt':  str(txts_dir),
            'docx': str(tmp_path / 'docx'),
            'xlsx': str(tmp_path / 'xlsx'),
            'pptx': str(tmp_path / 'pptx'),
            'csv':  str(tmp_path / 'csv'),
            'md':   str(tmp_path / 'md'),
            'html': str(tmp_path / 'html'),
        }
        loader.doc_folders = fake_folders
        files = loader.scan_all_files()
        assert len(files) == 1
        assert files[0]['is_misplaced'] is True
        assert files[0]['detected_type'] == 'pdf'
        assert files[0]['filename'] == 'report.pdf'

    def test_correct_folder_not_misplaced(self, loader, tmp_path):
        """A .pdf file in the pdfs/ folder is not flagged as misplaced."""
        pdfs_dir = tmp_path / 'pdfs'
        pdfs_dir.mkdir(parents=True)
        (pdfs_dir / 'doc.pdf').write_bytes(b'%PDF fake')

        fake_folders = {
            'pdf':  str(pdfs_dir),
            'txt':  str(tmp_path / 'txts'),
            'docx': str(tmp_path / 'docx'),
            'xlsx': str(tmp_path / 'xlsx'),
            'pptx': str(tmp_path / 'pptx'),
            'csv':  str(tmp_path / 'csv'),
            'md':   str(tmp_path / 'md'),
            'html': str(tmp_path / 'html'),
        }
        loader.doc_folders = fake_folders
        files = loader.scan_all_files()
        assert len(files) == 1
        assert files[0]['is_misplaced'] is False


# ---------------------------------------------------------------------------
# URL ingestion — mock requests.get
# ---------------------------------------------------------------------------

class TestChunkUrl:
    """Tests for chunk_url: 4-priority type detection and dispatch for remote content."""

    def _mock_resp(self, content, content_type='text/html', encoding='utf-8'):
        """Build a mock requests.Response with given content, Content-Type, and encoding."""
        resp = MagicMock()
        resp.content = content if isinstance(content, bytes) else content.encode(encoding)
        resp.headers = {'Content-Type': content_type}
        resp.encoding = encoding
        resp.raise_for_status = MagicMock()
        return resp

    def test_html_webpage(self, loader):
        """URL with text/html Content-Type produces chunks of type 'html'."""
        html = b'<html><body><p>Hello world. Second sentence. Third one.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = loader.chunk_url('https://example.com')
        assert len(chunks) >= 1
        assert all(c['type'] == 'html' for c in chunks)

    def test_type_by_content_type_pdf(self, loader, tmp_path):
        """URL with application/pdf Content-Type is dispatched to the PDF chunker."""
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF content sentence. Another one.")
        buf = io.BytesIO()
        doc.save(buf)
        pdf_bytes = buf.getvalue()
        with patch('requests.get', return_value=self._mock_resp(pdf_bytes, 'application/pdf')):
            chunks = loader.chunk_url('https://example.com/file')
        assert isinstance(chunks, list)

    def test_type_by_extension_in_url(self, loader):
        """URL ending in .html is typed as 'html' even when Content-Type is octet-stream."""
        html = b'<html><body><p>Content. More content. Third.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'application/octet-stream')):
            chunks = loader.chunk_url('https://example.com/page.html')
        assert all(c['type'] == 'html' for c in chunks)

    def test_type_by_pdf_magic_bytes(self, loader):
        """Content starting with %PDF magic bytes is identified as PDF regardless of headers."""
        # content-type is octet-stream, no extension, but starts with %PDF magic bytes
        fake_pdf = b'%PDF-1.4 fake content'
        with patch('requests.get', return_value=self._mock_resp(fake_pdf, 'application/octet-stream')):
            # fitz.open will fail on fake bytes and return an empty list gracefully
            chunks = loader.chunk_url('https://example.com/nodoc')
        assert isinstance(chunks, list)

    def test_defaults_to_html(self, loader):
        """Unrecognised Content-Type with no extension and no magic bytes defaults to 'html'."""
        content = b'Hello world. This is a sentence. Third sentence here.'
        with patch('requests.get', return_value=self._mock_resp(content, 'application/octet-stream')):
            chunks = loader.chunk_url('https://example.com/nodoc')
        assert isinstance(chunks, list)

    def test_connection_error_returns_empty(self, loader):
        """A network exception during requests.get returns an empty list without raising."""
        with patch('requests.get', side_effect=Exception('timeout')):
            chunks = loader.chunk_url('https://example.com')
        assert chunks == []

    def test_source_label_truncated_60_chars(self, loader):
        """Very long URLs are truncated to at most 60 characters in the chunk source label."""
        long_url = 'https://example.com/' + 'a' * 100
        html = b'<html><body><p>Text. More. Here.</p></body></html>'
        with patch('requests.get', return_value=self._mock_resp(html, 'text/html')):
            chunks = loader.chunk_url(long_url)
        if chunks:
            assert len(chunks[0]['source']) <= 60

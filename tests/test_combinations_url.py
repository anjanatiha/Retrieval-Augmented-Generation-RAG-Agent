"""test_combinations_url.py — URL content-type, URL type detection, and chunker contract tests.

Split from test_combinations.py to keep each file under 500 lines (per CLAUDE.md).

Covers:
  - TestUrlFetchContentTypes:       chunk_url for every text and binary Content-Type
  - TestUrlTypeDetection:           detected 'type' key matches Content-Type priority
  - TestChunkerContractParametrized: 5-key contract across all 8 doc types (parametrized)

Mock strategy (per CLAUDE.md):
  Always mock: requests.get
  Never mock:  fitz, python-docx, openpyxl, python-pptx, beautifulsoup4
"""

import os
from unittest.mock import MagicMock, patch

import pytest

from src.rag.chunkers import (
    chunk_txt, chunk_md, chunk_csv, chunk_html,
)
from src.rag.binary_chunkers import (
    chunk_pdf, chunk_docx, chunk_xlsx, chunk_pptx,
)

# Map doc type string → chunker function (used instead of getattr on loader)
_CHUNKER_MAP = {
    'txt':  chunk_txt,
    'md':   chunk_md,
    'csv':  chunk_csv,
    'html': chunk_html,
    'pdf':  chunk_pdf,
    'docx': chunk_docx,
    'xlsx': chunk_xlsx,
    'pptx': chunk_pptx,
}


# ---------------------------------------------------------------------------
# Binary file fixture helpers (real parsers — never mocked)
# ---------------------------------------------------------------------------

def _make_pdf_file(tmp_path, filename='test.pdf'):
    """Write a real single-page PDF and return the file path."""
    import fitz
    doc  = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Cats sleep sixteen hours a day. They are nocturnal hunters.")
    path = str(tmp_path / filename)
    doc.save(path); doc.close()
    return path


def _make_docx_file(tmp_path, filename='test.docx'):
    """Write a real DOCX file and return the file path."""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Cats sleep sixteen hours a day.")
    doc.add_paragraph("They are excellent nocturnal hunters.")
    path = str(tmp_path / filename)
    doc.save(path)
    return path


def _make_xlsx_file(tmp_path, filename='test.xlsx'):
    """Write a real XLSX file and return the file path."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(['animal', 'fact'])
    ws.append(['cat', 'sleeps 16 hours a day'])
    path = str(tmp_path / filename)
    wb.save(path)
    return path


def _make_pptx_file(tmp_path, filename='test.pptx'):
    """Write a real PPTX file and return the file path."""
    from pptx import Presentation
    prs   = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text               = "Cat Facts"
    slide.placeholders[1].text_frame.text = "Cats sleep 16 hours a day."
    path  = str(tmp_path / filename)
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# 1. URL fetch × all URL content types
# ---------------------------------------------------------------------------

class TestUrlFetchContentTypes:
    """chunk_url returns non-empty chunks for every supported content-type."""

    @pytest.mark.parametrize("content_type,url_path,expected_type", [
        ("text/html",                "http://example.com/page",      "html"),
        ("text/plain",               "http://example.com/file.txt",  "txt"),
        ("text/csv",                 "http://example.com/data.csv",  "csv"),
        ("text/markdown",            "http://example.com/doc.md",    "md"),
        ("application/json",         "http://example.com/page",      "html"),  # fallback
    ])
    def test_text_content_type_produces_chunks(self, content_type, url_path,
                                               expected_type) -> None:
        """chunk_url for text-based content type returns chunks with correct shape.

        Args:
            content_type:  HTTP Content-Type header value.
            url_path:      URL string passed to chunk_url.
            expected_type: Expected 'type' value in returned chunks.
        """
        from src.rag.document_loader import DocumentLoader
        loader = DocumentLoader()

        mock_resp             = MagicMock()
        mock_resp.headers     = {'Content-Type': content_type}
        mock_resp.content     = (
            b"<html><body><p>Cats sleep 16 hours a day. "
            b"They are nocturnal hunters.</p></body></html>"
            if expected_type in ('html', 'webpage')
            else b"Cats sleep 16 hours a day.\nDogs are loyal.\n"
        )
        mock_resp.encoding         = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url(url_path)

        assert isinstance(chunks, list), f"chunk_url non-list for {content_type}"
        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"[{content_type}] chunk {i} missing key '{key}'"

    @pytest.mark.parametrize("content_type,url_path,expected_type", [
        ("application/pdf",          "http://example.com/doc",          "pdf"),
        ("application/octet-stream", "http://example.com/d.xlsx",       "xlsx"),
        ("application/octet-stream", "http://example.com/d.pptx",       "pptx"),
        ("application/octet-stream", "http://example.com/d.docx",       "docx"),
        ("application/octet-stream", "http://example.com/data.csv",     "csv"),
    ])
    def test_binary_content_type_produces_chunks(self, content_type, url_path,
                                                  expected_type, tmp_path) -> None:
        """chunk_url for binary Content-Type runs real parsers on actual file bytes.

        Args:
            content_type:  HTTP Content-Type header value.
            url_path:      URL string passed to chunk_url.
            expected_type: Expected doc type (selects fixture helper).
        """
        from src.rag.document_loader import DocumentLoader

        if expected_type == 'pdf':
            path = _make_pdf_file(tmp_path)
        elif expected_type == 'xlsx':
            path = _make_xlsx_file(tmp_path)
        elif expected_type == 'pptx':
            path = _make_pptx_file(tmp_path)
        elif expected_type == 'docx':
            path = _make_docx_file(tmp_path)
        elif expected_type == 'csv':
            path = str(tmp_path / 'data.csv')
            with open(path, 'w') as fh:
                fh.write("animal,fact\ncat,sleeps 16 hours\n")
        else:
            pytest.skip(f"No fixture for {expected_type}")

        with open(path, 'rb') as fh:
            real_bytes = fh.read()

        loader                     = DocumentLoader()
        mock_resp                  = MagicMock()
        mock_resp.headers          = {'Content-Type': content_type}
        mock_resp.content          = real_bytes
        mock_resp.encoding         = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url(url_path)

        assert isinstance(chunks, list), f"chunk_url non-list for {content_type}"
        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"[{content_type}] chunk {i} missing key '{key}'"


# ---------------------------------------------------------------------------
# 2. URL type detection × detection mechanism
# ---------------------------------------------------------------------------

class TestUrlTypeDetection:
    """URL type is correctly inferred from Content-Type header (priority 1)."""

    @pytest.mark.parametrize("content_type,url_path,expected_type", [
        ("text/html",            "http://example.com/page",      "html"),
        ("application/pdf",      "http://example.com/doc",       "pdf"),
        ("application/octet-stream", "http://example.com/d.csv", "csv"),
        ("text/plain",           "http://example.com/f.txt",     "txt"),
        ("application/octet-stream", "http://example.com/d.xlsx","xlsx"),
        ("application/octet-stream", "http://example.com/d.pptx","pptx"),
    ])
    def test_url_type_by_content_type(self, content_type, url_path,
                                       expected_type, tmp_path) -> None:
        """Detected URL 'type' matches expected_type for each Content-Type header.

        Args:
            content_type:  HTTP Content-Type to inject in the mock response.
            url_path:      URL string.
            expected_type: Expected 'type' in returned chunk dicts.
        """
        from src.rag.document_loader import DocumentLoader
        loader = DocumentLoader()

        # Build minimal real bytes for binary types so parsers don't fail
        if expected_type == 'pdf':
            path = _make_pdf_file(tmp_path)
            with open(path, 'rb') as fh: body = fh.read()
        elif expected_type == 'xlsx':
            path = _make_xlsx_file(tmp_path)
            with open(path, 'rb') as fh: body = fh.read()
        elif expected_type == 'pptx':
            path = _make_pptx_file(tmp_path)
            with open(path, 'rb') as fh: body = fh.read()
        elif expected_type in ('csv', 'txt'):
            body = b"animal,fact\ncat,sleeps 16 hours\n"
        else:
            body = b"<html><body><p>Cats sleep 16 hours.</p></body></html>"

        mock_resp                  = MagicMock()
        mock_resp.headers          = {'Content-Type': content_type}
        mock_resp.content          = body
        mock_resp.encoding         = 'utf-8'
        mock_resp.raise_for_status = MagicMock()

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url(url_path)

        if chunks:
            detected = chunks[0]['type']
            if expected_type == 'html':
                assert detected in ('html', 'webpage'), (
                    f"Expected html/webpage, got {detected}"
                )
            else:
                assert detected == expected_type, (
                    f"content_type={content_type!r}: "
                    f"expected {expected_type!r}, got {detected!r}"
                )


# ---------------------------------------------------------------------------
# 3. Chunker contract × every doc type (parametrized)
# ---------------------------------------------------------------------------

class TestChunkerContractParametrized:
    """Every chunker returns dicts with all 5 required keys — parametrized across types."""

    @pytest.mark.parametrize("doc_type,text_content,filename", [
        ("txt",  "Cats sleep 16 hours.\nDogs are loyal.",           "test.txt"),
        ("md",   "# Animals\nCats sleep 16 hours.\nDogs loyal.",    "test.md"),
        ("csv",  "animal,fact\ncat,sleeps 16 hours\ndog,loyal",     "test.csv"),
        ("html", "<html><body><p>Cats sleep 16 hours.</p></body></html>", "test.html"),
    ])
    def test_text_chunker_5_key_contract(self, doc_type: str, text_content: str,
                                          filename: str, tmp_path) -> None:
        """Text-format chunkers produce dicts with all 5 required keys.

        Args:
            doc_type:     One of txt, md, csv, html.
            text_content: Raw string written to the temp file.
            filename:     Filename passed to the chunker.
        """
        f = tmp_path / filename
        f.write_text(text_content)
        chunker = _CHUNKER_MAP[doc_type]
        chunks  = chunker(str(f), filename)

        assert isinstance(chunks, list), f"_chunk_{doc_type} did not return a list"
        assert len(chunks) >= 1,         f"_chunk_{doc_type} returned empty list"

        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"_chunk_{doc_type} chunk[{i}] missing key '{key}'"
            assert isinstance(c['text'],       str), f"chunk[{i}]['text'] not str"
            assert isinstance(c['source'],     str), f"chunk[{i}]['source'] not str"
            assert isinstance(c['start_line'], int), f"chunk[{i}]['start_line'] not int"
            assert isinstance(c['end_line'],   int), f"chunk[{i}]['end_line'] not int"
            assert isinstance(c['type'],       str), f"chunk[{i}]['type'] not str"

    @pytest.mark.parametrize("doc_type,make_fn", [
        ("pdf",  _make_pdf_file),
        ("docx", _make_docx_file),
        ("xlsx", _make_xlsx_file),
        ("pptx", _make_pptx_file),
    ])
    def test_binary_chunker_5_key_contract(self, doc_type: str, make_fn,
                                            tmp_path) -> None:
        """Binary-format chunkers produce dicts with all 5 required keys.

        Args:
            doc_type: One of pdf, docx, xlsx, pptx.
            make_fn:  Fixture helper that creates a real binary temp file.
        """
        path     = make_fn(tmp_path)
        filename = os.path.basename(path)
        chunker  = _CHUNKER_MAP[doc_type]
        chunks   = chunker(path, filename)

        assert isinstance(chunks, list), f"_chunk_{doc_type} did not return a list"
        assert len(chunks) >= 1,         f"_chunk_{doc_type} returned empty list"

        for i, c in enumerate(chunks):
            for key in ('text', 'source', 'start_line', 'end_line', 'type'):
                assert key in c, f"_chunk_{doc_type} chunk[{i}] missing key '{key}'"
            assert isinstance(c['text'],       str), f"chunk[{i}]['text'] not str"
            assert isinstance(c['source'],     str), f"chunk[{i}]['source'] not str"
            assert isinstance(c['start_line'], int), f"chunk[{i}]['start_line'] not int"
            assert isinstance(c['end_line'],   int), f"chunk[{i}]['end_line'] not int"
            assert isinstance(c['type'],       str), f"chunk[{i}]['type'] not str"

"""test_contracts.py — Contract tests for chunkers, URL ingestion, and search internals.

Contract tests assert output *shape*, required keys, and value types — not exact values.
This guards against accidental API breakage across refactors without tying tests to
specific LLM outputs.

Scope (this file): chunkers (_chunk_*), chunk_url(), _hybrid_retrieve(), _rerank().
run_pipeline() and agent.run() contracts are in test_contracts_pipeline.py.

Mock strategy (per CLAUDE.md):
  Always mock: ollama.embed, ollama.chat, chromadb → EphemeralClient
  Never mock:  BM25Okapi, chunk dicts, _classify_query, _source_label
"""

from unittest.mock import MagicMock, patch

import chromadb
import pytest
from rank_bm25 import BM25Okapi

from src.rag.binary_chunkers import chunk_docx, chunk_pdf, chunk_pptx, chunk_xls, chunk_xlsx
from src.rag.chunkers import chunk_csv, chunk_html, chunk_md, chunk_txt

# ---------------------------------------------------------------------------
# Chunk-dict contract helper
# ---------------------------------------------------------------------------

# Required keys and expected types for every chunk dict
CHUNK_KEYS = {
    'text':       str,
    'source':     str,
    'start_line': int,
    'end_line':   int,
    'type':       str,
}


def _assert_chunk_contract(chunk: dict, label: str = '') -> None:
    """Assert that *chunk* satisfies the 5-key chunk-dict contract.

    Args:
        chunk: The chunk dict to validate.
        label: Optional identifier printed in assertion messages.
    """
    prefix = f"[{label}] " if label else ""
    for key, expected_type in CHUNK_KEYS.items():
        assert key in chunk, f"{prefix}Missing key '{key}'"
        assert chunk[key] is not None, f"{prefix}Key '{key}' is None"
        assert isinstance(chunk[key], expected_type), (
            f"{prefix}Key '{key}' expected {expected_type.__name__}, "
            f"got {type(chunk[key]).__name__}"
        )


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_chunks():
    """Five minimal cat-facts txt chunks used across contract tests."""
    return [
        {
            'text':       f"Cats sleep about 16 hours a day chunk {i}.",
            'source':     'cats.txt',
            'start_line': i + 1,
            'end_line':   i + 1,
            'type':       'txt',
        }
        for i in range(5)
    ]


@pytest.fixture
def store(sample_chunks):
    """Return a VectorStore wired to an in-memory EphemeralClient.

    Uses real BM25Okapi and real ChromaDB so only ollama is mocked.
    """
    from src.rag.vector_store import VectorStore
    vs     = VectorStore()
    client = chromadb.EphemeralClient()
    col    = client.get_or_create_collection(
        'test_contracts', metadata={'hnsw:space': 'cosine'}
    )
    ids    = [f'c{i}' for i in range(len(sample_chunks))]
    texts  = [c['text'] for c in sample_chunks]
    metas  = [
        {'source': c['source'], 'start_line': c['start_line'],
         'end_line': c['end_line'], 'type': c['type']}
        for c in sample_chunks
    ]
    embeds = [[0.1, 0.2, 0.3, 0.4]] * len(sample_chunks)
    col.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
    vs.collection           = col
    vs.chunks               = sample_chunks
    vs.bm25_index           = BM25Okapi([c['text'].lower().split() for c in sample_chunks])
    vs.conversation_history = []
    return vs


@pytest.fixture
def loader():
    """Fresh DocumentLoader instance for chunker contract tests."""
    from src.rag.document_loader import DocumentLoader
    return DocumentLoader()


# ---------------------------------------------------------------------------
# 1. Chunker contracts — every _chunk_* method
# ---------------------------------------------------------------------------

class TestChunkerContracts:
    """Every chunker returns a list of dicts, each satisfying the 5-key contract."""

    def test_chunk_txt_contract(self, loader, tmp_path):
        """_chunk_txt chunks all satisfy the 5-key contract with correct types."""
        f = tmp_path / 'test.txt'
        f.write_text("Cats sleep 16 hours.\nDogs are loyal companions.\n")
        chunks = chunk_txt(str(f), 'test.txt')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"txt[{i}]")

    def test_chunk_md_contract(self, loader, tmp_path):
        """_chunk_md chunks all satisfy the 5-key contract with correct types."""
        f = tmp_path / 'test.md'
        f.write_text("# Animals\nCats sleep 16 hours.\nDogs are loyal.\n")
        chunks = chunk_md(str(f), 'test.md')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"md[{i}]")

    def test_chunk_pdf_contract(self, loader, tmp_path):
        """_chunk_pdf chunks all satisfy the 5-key contract with correct types."""
        import fitz
        doc  = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "Cats sleep sixteen hours a day. They hunt at night.")
        path = str(tmp_path / 'test.pdf')
        doc.save(path); doc.close()
        chunks = chunk_pdf(path, 'test.pdf')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"pdf[{i}]")

    def test_chunk_docx_contract(self, loader, tmp_path):
        """_chunk_docx chunks all satisfy the 5-key contract with correct types."""
        from docx import Document
        doc = Document()
        doc.add_paragraph("Cats sleep sixteen hours a day.")
        doc.add_paragraph("They are excellent nocturnal hunters.")
        path = str(tmp_path / 'test.docx')
        doc.save(path)
        chunks = chunk_docx(path, 'test.docx')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"docx[{i}]")

    def test_chunk_xlsx_contract(self, loader, tmp_path):
        """_chunk_xlsx chunks all satisfy the 5-key contract with correct types."""
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(['animal', 'fact'])
        ws.append(['cat', 'sleeps 16 hours'])
        path = str(tmp_path / 'test.xlsx')
        wb.save(path)
        chunks = chunk_xlsx(path, 'test.xlsx')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"xlsx[{i}]")

    def test_chunk_csv_contract(self, loader, tmp_path):
        """_chunk_csv chunks all satisfy the 5-key contract with correct types."""
        f = tmp_path / 'test.csv'
        f.write_text("animal,fact\ncat,sleeps 16 hours\ndog,loyal\n")
        chunks = chunk_csv(str(f), 'test.csv')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"csv[{i}]")

    def test_chunk_pptx_contract(self, loader, tmp_path):
        """_chunk_pptx chunks all satisfy the 5-key contract with correct types."""
        from pptx import Presentation
        from pptx.util import Pt
        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text           = "Cat Facts"
        slide.placeholders[1].text_frame.text = "Cats sleep 16 hours a day."
        path  = str(tmp_path / 'test.pptx')
        prs.save(path)
        chunks = chunk_pptx(path, 'test.pptx')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"pptx[{i}]")

    def test_chunk_html_contract(self, loader, tmp_path):
        """_chunk_html chunks all satisfy the 5-key contract with correct types."""
        f = tmp_path / 'test.html'
        f.write_text(
            "<html><body><p>Cats sleep 16 hours a day. "
            "They are nocturnal hunters.</p></body></html>"
        )
        chunks = chunk_html(str(f), 'test.html')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"html[{i}]")

    def test_chunk_xls_contract(self, loader, tmp_path):
        """_chunk_xls chunks all satisfy the 5-key contract with correct types."""
        import xlwt  # only for test fixture creation
        wb = xlwt.Workbook()
        ws = wb.add_sheet('Sheet1')
        ws.write(0, 0, 'animal'); ws.write(0, 1, 'fact')
        ws.write(1, 0, 'cat');    ws.write(1, 1, 'sleeps 16 hours')
        path = str(tmp_path / 'test.xls')
        wb.save(path)
        chunks = chunk_xls(path, 'test.xls')
        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"xls[{i}]")


# ---------------------------------------------------------------------------
# 2. chunk_url contract
# ---------------------------------------------------------------------------

class TestChunkUrlContract:
    """chunk_url() returns a list of dicts each satisfying the 5-key contract."""

    def test_chunk_url_html_contract(self, loader):
        """chunk_url for an HTML page returns chunk dicts with all required keys."""
        mock_resp = MagicMock()
        mock_resp.headers = {'Content-Type': 'text/html'}
        mock_resp.content = (
            b"<html><body><p>Cats sleep 16 hours a day. "
            b"They are nocturnal hunters.</p></body></html>"
        )
        mock_resp.encoding = 'utf-8'
        mock_resp.raise_for_status.return_value = None

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url('http://example.com/page')

        assert isinstance(chunks, list)
        assert len(chunks) >= 1
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"url_html[{i}]")

    def test_chunk_url_txt_contract(self, loader):
        """chunk_url for a plain-text URL returns chunk dicts with all required keys."""
        mock_resp = MagicMock()
        mock_resp.headers = {'Content-Type': 'text/plain'}
        mock_resp.content = b"Cats sleep 16 hours a day.\nDogs are loyal.\n"
        mock_resp.encoding = 'utf-8'
        mock_resp.raise_for_status.return_value = None

        with patch('requests.get', return_value=mock_resp):
            chunks = loader.chunk_url('http://example.com/file.txt')

        assert isinstance(chunks, list)
        # txt URL can return empty list for some content — only check contract when non-empty
        for i, c in enumerate(chunks):
            _assert_chunk_contract(c, f"url_txt[{i}]")


# ---------------------------------------------------------------------------
# 3. _hybrid_retrieve contract
# ---------------------------------------------------------------------------

class TestHybridRetrieveContract:
    """_hybrid_retrieve returns list of (dict, float) tuples with correct structure."""

    def test_returns_list(self, store):
        """_hybrid_retrieve always returns a list."""
        with patch('ollama.embed', return_value={'embeddings': [[0.1, 0.2, 0.3, 0.4]]}):
            result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        assert isinstance(result, list)

    def test_tuple_shape(self, store):
        """Each element is a 2-tuple of (dict, float)."""
        with patch('ollama.embed', return_value={'embeddings': [[0.1, 0.2, 0.3, 0.4]]}):
            result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        for item in result:
            assert len(item) == 2, f"Expected 2-tuple, got length {len(item)}"
            entry, score = item
            assert isinstance(entry, dict)
            assert isinstance(score, float)

    def test_entry_has_required_keys(self, store):
        """Each entry dict contains all 5 required chunk keys."""
        with patch('ollama.embed', return_value={'embeddings': [[0.1, 0.2, 0.3, 0.4]]}):
            result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        for entry, _ in result:
            _assert_chunk_contract(entry, '_hybrid_retrieve entry')

    def test_score_is_non_negative(self, store):
        """All returned hybrid scores are >= 0."""
        with patch('ollama.embed', return_value={'embeddings': [[0.1, 0.2, 0.3, 0.4]]}):
            result = store._hybrid_retrieve(['cats sleep'], top_n=3)
        for _, score in result:
            assert score >= 0.0, f"Score {score} is negative"


# ---------------------------------------------------------------------------
# 4. _rerank contract
# ---------------------------------------------------------------------------

class TestRerankContract:
    """_rerank returns list of (dict, float, float) triples with correct structure."""

    def test_returns_list(self, store):
        """_rerank always returns a list."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        with patch('ollama.chat', return_value={'message': {'content': '8'}}):
            result = store._rerank('cats sleep', candidates, top_n=2)
        assert isinstance(result, list)

    def test_triple_shape(self, store):
        """Each element is a 3-tuple of (dict, float, float)."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        with patch('ollama.chat', return_value={'message': {'content': '8'}}):
            result = store._rerank('cats sleep', candidates, top_n=2)
        for item in result:
            assert len(item) == 3, f"Expected 3-tuple, got length {len(item)}"
            entry, sim, llm_score = item
            assert isinstance(entry, dict)
            assert isinstance(sim, float)
            assert isinstance(llm_score, float)

    def test_entry_has_required_keys(self, store):
        """Each entry dict in rerank results has all 5 required chunk keys."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        with patch('ollama.chat', return_value={'message': {'content': '8'}}):
            result = store._rerank('cats sleep', candidates, top_n=2)
        for entry, _, _ in result:
            _assert_chunk_contract(entry, '_rerank entry')

    def test_scores_bounded(self, store):
        """sim and llm_score are both in [0, 1] range when LLM returns valid int."""
        candidates = [(c, 0.8) for c in store.chunks[:2]]
        with patch('ollama.chat', return_value={'message': {'content': '8'}}):
            result = store._rerank('cats sleep', candidates, top_n=2)
        for _, sim, llm_score in result:
            assert 0.0 <= sim <= 1.0, f"sim {sim} out of [0, 1]"
            assert 0.0 <= llm_score <= 1.0, f"llm_score {llm_score} out of [0, 1]"

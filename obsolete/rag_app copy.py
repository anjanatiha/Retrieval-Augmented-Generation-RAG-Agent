"""
RAG Chatbot — Full Enhanced Pipeline
=====================================
Features:
  1.  Sliding window chunking with overlap
  2.  Persistent vector DB (ChromaDB)
  3.  Hybrid search (BM25 + dense vector)
  4.  Query expansion
  5.  LLM reranker
  6.  Query classification
  7.  Confidence / hallucination filter
  8.  Source citation
  9.  Conversation memory
  10. Logging & analytics
  11. Streaming with typing indicator
  12. Benchmarking with before/after comparison
  13. Agent with tool calling (RAG, calculator, summarise)
  14. Streamlit UI
  15. PDF support      (./docs/pdfs/)   [NEW — subfolder]
  16. TXT support      (./docs/txts/)   [NEW — subfolder]
  17. DOCX support     (./docs/docx/)   [NEW]
  18. XLSX support     (./docs/xlsx/)   [NEW]
  19. PPTX support     (./docs/pptx/)   [NEW]
  20. CSV support      (./docs/csv/)    [NEW]
  21. MD support       (./docs/md/)     [NEW]
  22. HTML support     (./docs/html/)   [NEW]
  23. Smart mis-placed file detection   [NEW — wrong folder ✓ auto-corrected]

Folder structure (all created automatically):
  ./docs/
  ./docs/pdfs/   ← .pdf files
  ./docs/txts/   ← .txt files
  ./docs/docx/   ← .docx / .doc files
  ./docs/xlsx/   ← .xlsx / .xls / .csv files  (CSV can also go in csv/)
  ./docs/pptx/   ← .pptx / .ppt files
  ./docs/csv/    ← .csv files
  ./docs/md/     ← .md / .markdown files
  ./docs/html/   ← .html / .htm files

  Smart detection: drop ANY supported file into ANY subfolder — the loader
  detects the real type by extension and processes it correctly, printing a
  friendly [MISPLACED] notice so you know to reorganise later.

Run modes:
  python3 rag_app12.py                  # terminal chatbot
  python3 rag_app12.py --agent          # agent mode (terminal)
  python3 rag_app12.py --benchmark      # benchmark evaluation
  streamlit run rag_app12.py            # Streamlit UI

Required packages (install as needed):
  pip install rank-bm25 chromadb ollama
  pip install pymupdf                   # PDF
  pip install python-docx               # DOCX
  pip install openpyxl xlrd             # XLSX / XLS
  pip install python-pptx               # PPTX
  pip install beautifulsoup4            # HTML
  # CSV and MD are stdlib / no extra deps
"""

import ollama, json, os, sys, re, time, argparse
from datetime import datetime
from rank_bm25 import BM25Okapi
import chromadb
from chromadb.config import Settings

# ============================================================
# CONFIGURATION
# ============================================================
EMBEDDING_MODEL      = 'hf.co/CompendiumLabs/bge-base-en-v1.5-gguf'
LANGUAGE_MODEL       = 'hf.co/bartowski/Llama-3.2-3B-Instruct-GGUF'

# Root docs folder
DOCS_ROOT            = './docs'

# Subfolders — canonical home for each type
DOC_FOLDERS = {
    'pdf':  os.path.join(DOCS_ROOT, 'pdfs'),
    'txt':  os.path.join(DOCS_ROOT, 'txts'),
    'docx': os.path.join(DOCS_ROOT, 'docx'),
    'xlsx': os.path.join(DOCS_ROOT, 'xlsx'),
    'pptx': os.path.join(DOCS_ROOT, 'pptx'),
    'csv':  os.path.join(DOCS_ROOT, 'csv'),
    'md':   os.path.join(DOCS_ROOT, 'md'),
    'html': os.path.join(DOCS_ROOT, 'html'),
}

# Maps every supported file extension → canonical type key
EXT_TO_TYPE = {
    '.pdf':      'pdf',
    '.txt':      'txt',
    '.docx':     'docx',
    '.doc':      'docx',
    '.xlsx':     'xlsx',
    '.xls':      'xlsx',
    '.pptx':     'pptx',
    '.ppt':      'pptx',
    '.csv':      'csv',
    '.md':       'md',
    '.markdown': 'md',
    '.html':     'html',
    '.htm':      'html',
}

CHROMA_DIR           = './chroma_db'
CHROMA_COLLECTION    = 'rag_docs'
LOG_FILE             = 'rag_logs.json'
BENCHMARK_FILE       = 'benchmark_results.json'
SIMILARITY_THRESHOLD = 0.4
TOP_RETRIEVE         = 20
TOP_RERANK           = 3
# TXT / MD: 1 line per chunk — original behaviour, do not change
TXT_CHUNK_SIZE       = 1
TXT_CHUNK_OVERLAP    = 0
# New format chunk sizes — tuned per type
PDF_CHUNK_SENTENCES  = 5
DOCX_CHUNK_PARAS     = 3
PPTX_CHUNK_SLIDES    = 1
HTML_CHUNK_SENTENCES = 5

# ============================================================
# FOLDER SETUP — auto-create all subfolders
# ============================================================
def ensure_folders():
    os.makedirs(DOCS_ROOT, exist_ok=True)
    for folder in DOC_FOLDERS.values():
        if not os.path.exists(folder):
            os.makedirs(folder)
            print(f"  Created folder: {folder}/")

# ============================================================
# SMART FILE SCANNER
# Walks ALL subfolders, detects real type by extension.
# Files in a "wrong" folder still get processed correctly.
# ============================================================
def scan_all_files():
    """
    Returns a list of (filepath, detected_type, canonical_folder, is_misplaced).
    Scans every subfolder under DOCS_ROOT so no file is missed.
    """
    found = []
    for folder_type, folder_path in DOC_FOLDERS.items():
        if not os.path.isdir(folder_path):
            continue
        for fname in os.listdir(folder_path):
            ext = os.path.splitext(fname)[1].lower()
            detected_type = EXT_TO_TYPE.get(ext)
            if detected_type is None:
                continue  # unsupported extension, skip silently
            filepath      = os.path.join(folder_path, fname)
            canonical_dir = DOC_FOLDERS[detected_type]
            is_misplaced  = (folder_path != canonical_dir)
            found.append({
                'filepath':      filepath,
                'filename':      fname,
                'detected_type': detected_type,
                'found_in':      folder_path,
                'canonical_dir': canonical_dir,
                'is_misplaced':  is_misplaced,
            })
            if is_misplaced:
                print(f"  [MISPLACED] '{fname}' found in '{folder_path}/' "
                      f"— detected as '{detected_type.upper()}', "
                      f"canonical folder is '{canonical_dir}/'. Processing anyway.")
    return found

# ============================================================
# CHUNKERS — one per type
# ============================================================

# ---------- TXT ----------
def _chunk_txt(filepath, filename, chunk_size=TXT_CHUNK_SIZE, overlap=TXT_CHUNK_OVERLAP):
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        lines = [l.strip() for l in f.readlines() if l.strip()]
    step   = max(1, chunk_size - overlap)
    chunks = []
    for i in range(0, len(lines), step):
        window = lines[i: i + chunk_size]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'txt',
        })
    return chunks

# ---------- MD (Markdown) ----------
def _chunk_md(filepath, filename, chunk_size=TXT_CHUNK_SIZE, overlap=TXT_CHUNK_OVERLAP):
    """Strips markdown syntax and chunks like plain text."""
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        raw = f.read()
    # Strip common markdown markers for cleaner embedding
    clean = re.sub(r'#{1,6}\s*', '', raw)      # headings
    clean = re.sub(r'[*_`~]{1,3}', '', clean)  # bold/italic/code
    clean = re.sub(r'!\[.*?\]\(.*?\)', '', clean)  # images
    clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)  # links → text
    lines = [l.strip() for l in clean.splitlines() if l.strip()]
    step  = max(1, chunk_size - overlap)
    chunks = []
    for i in range(0, len(lines), step):
        window = lines[i: i + chunk_size]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'md',
        })
    return chunks

# ---------- PDF ----------
def _chunk_pdf(filepath, filename, sentences_per_chunk=PDF_CHUNK_SENTENCES):
    try:
        import fitz
    except ImportError:
        print("  [WARNING] pymupdf not installed — skipping PDF. Install: pip install pymupdf")
        return []

    chunks  = []
    page_num = 0
    try:
        doc = fitz.open(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    for page_num, page in enumerate(doc, start=1):
        raw = page.get_text()
        if not raw.strip():
            continue
        sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', raw) if s.strip()]
        for i in range(0, len(sentences), sentences_per_chunk):
            window = sentences[i: i + sentences_per_chunk]
            if not window:
                continue
            chunks.append({
                'text':       ' '.join(window),
                'source':     filename,
                'start_line': page_num,
                'end_line':   page_num,
                'type':       'pdf',
            })
    doc.close()
    return chunks

# ---------- DOCX ----------
def _chunk_docx(filepath, filename, paras_per_chunk=DOCX_CHUNK_PARAS):
    try:
        from docx import Document
    except ImportError:
        print("  [WARNING] python-docx not installed — skipping DOCX. Install: pip install python-docx")
        return []

    chunks = []
    try:
        doc = Document(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for i in range(0, len(paras), paras_per_chunk):
        window = paras[i: i + paras_per_chunk]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'docx',
        })
    return chunks

# ---------- XLSX / XLS / CSV (spreadsheet row chunking) ----------
def _chunk_xlsx(filepath, filename):
    """Each non-empty row becomes a chunk (joined as key=value pairs)."""
    try:
        import openpyxl
    except ImportError:
        print("  [WARNING] openpyxl not installed — skipping XLSX. Install: pip install openpyxl")
        return []

    chunks = []
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    for sheet in wb.sheetnames:
        ws     = wb[sheet]
        rows   = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else f"col{i}" for i, c in enumerate(rows[0])]
        for row_idx, row in enumerate(rows[1:], start=2):
            cells = [str(v) for v in row if v is not None and str(v).strip()]
            if not cells:
                continue
            pairs = '; '.join(
                f"{header[i] if i < len(header) else f'col{i}'}={str(row[i])}"
                for i in range(len(row)) if row[i] is not None and str(row[i]).strip()
            )
            chunks.append({
                'text':       f"[{sheet}] {pairs}",
                'source':     filename,
                'start_line': row_idx,
                'end_line':   row_idx,
                'type':       'xlsx',
            })
    wb.close()
    return chunks

def _chunk_xls(filepath, filename):
    """Fallback for old .xls via xlrd."""
    try:
        import xlrd
    except ImportError:
        print("  [WARNING] xlrd not installed — skipping XLS. Install: pip install xlrd")
        return []

    chunks = []
    try:
        wb = xlrd.open_workbook(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    for sheet in wb.sheets():
        if sheet.nrows < 2:
            continue
        header = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
        for row_idx in range(1, sheet.nrows):
            pairs = '; '.join(
                f"{header[c] if c < len(header) else f'col{c}'}={sheet.cell_value(row_idx, c)}"
                for c in range(sheet.ncols)
                if str(sheet.cell_value(row_idx, c)).strip()
            )
            if not pairs:
                continue
            chunks.append({
                'text':       f"[{sheet.name}] {pairs}",
                'source':     filename,
                'start_line': row_idx + 1,
                'end_line':   row_idx + 1,
                'type':       'xlsx',
            })
    return chunks

def _chunk_csv(filepath, filename):
    """Each CSV row becomes a chunk."""
    import csv
    chunks = []
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            reader = csv.DictReader(f)
            for row_idx, row in enumerate(reader, start=2):
                pairs = '; '.join(f"{k}={v}" for k, v in row.items() if v and str(v).strip())
                if not pairs:
                    continue
                chunks.append({
                    'text':       pairs,
                    'source':     filename,
                    'start_line': row_idx,
                    'end_line':   row_idx,
                    'type':       'csv',
                })
    except Exception as e:
        print(f"  [ERROR] Could not read '{filename}': {e}")
    return chunks

# ---------- PPTX ----------
def _chunk_pptx(filepath, filename, slides_per_chunk=PPTX_CHUNK_SLIDES):
    try:
        from pptx import Presentation
    except ImportError:
        print("  [WARNING] python-pptx not installed — skipping PPTX. Install: pip install python-pptx")
        return []

    chunks = []
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slide_texts.append(' '.join(texts))

    for i in range(0, len(slide_texts), slides_per_chunk):
        window = slide_texts[i: i + slides_per_chunk]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'pptx',
        })
    return chunks

# ---------- HTML ----------
def _chunk_html(filepath, filename, sentences_per_chunk=HTML_CHUNK_SENTENCES):
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        print("  [WARNING] beautifulsoup4 not installed — skipping HTML. Install: pip install beautifulsoup4")
        return []

    chunks = []
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
    except Exception as e:
        print(f"  [ERROR] Could not open '{filename}': {e}")
        return []

    text      = soup.get_text(separator=' ', strip=True)
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
    for i in range(0, len(sentences), sentences_per_chunk):
        window = sentences[i: i + sentences_per_chunk]
        if not window:
            continue
        chunks.append({
            'text':       ' '.join(window),
            'source':     filename,
            'start_line': i + 1,
            'end_line':   i + len(window),
            'type':       'html',
        })
    return chunks

# ============================================================
# DISPATCH TABLE — maps type key → chunker function
# ============================================================
def _dispatch_chunker(file_info):
    """
    Picks the right chunker based on detected_type (not folder).
    This ensures misplaced files are still processed correctly.
    """
    fp  = file_info['filepath']
    fn  = file_info['filename']
    ext = os.path.splitext(fn)[1].lower()
    t   = file_info['detected_type']

    if t == 'txt':
        return _chunk_txt(fp, fn)
    elif t == 'md':
        return _chunk_md(fp, fn)
    elif t == 'pdf':
        return _chunk_pdf(fp, fn)
    elif t == 'docx':
        return _chunk_docx(fp, fn)
    elif t == 'xlsx':
        if ext == '.xls':
            return _chunk_xls(fp, fn)
        elif ext == '.csv':
            return _chunk_csv(fp, fn)
        else:
            return _chunk_xlsx(fp, fn)
    elif t == 'csv':
        return _chunk_csv(fp, fn)
    elif t == 'pptx':
        return _chunk_pptx(fp, fn)
    elif t == 'html':
        return _chunk_html(fp, fn)
    else:
        print(f"  [SKIP] No chunker for type '{t}' ({fn})")
        return []

# ============================================================
# COMBINED CHUNKING — scans all subfolders, smart dispatch
# ============================================================
def chunk_all_documents():
    """
    Scans every subfolder under ./docs/, auto-detects file type,
    and routes each file to the correct chunker.
    Files in the wrong folder are still processed (with a notice).
    """
    print("\nLoading documents...")
    file_list = scan_all_files()

    if not file_list:
        print(f"\nNo supported documents found under '{DOCS_ROOT}/'")
        print("Supported types: PDF, TXT, DOCX, XLSX, XLS, PPTX, CSV, MD, HTML")
        print("Place files in the matching subfolder (or any subfolder — smart detection handles the rest).")
        sys.exit(1)

    all_chunks  = []
    type_counts = {}

    for file_info in file_list:
        t      = file_info['detected_type']
        chunks = _dispatch_chunker(file_info)
        all_chunks.extend(chunks)
        type_counts[t] = type_counts.get(t, 0) + len(chunks)
        flag = " [MISPLACED — processed anyway]" if file_info['is_misplaced'] else ""
        print(f"  [{t.upper()}] '{file_info['filename']}': {len(chunks)} chunks{flag}")

    print("\n  Chunk summary by type:")
    for t, count in sorted(type_counts.items()):
        print(f"    {t.upper():<8} {count} chunks")
    print(f"  Total: {len(all_chunks)} chunks\n")

    return all_chunks

# ============================================================
# 2. PERSISTENT VECTOR DB — ChromaDB
# ============================================================
def get_chroma_collection():
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    return client.get_or_create_collection(
        name=CHROMA_COLLECTION,
        metadata={"hnsw:space": "cosine"}
    )

def build_or_load_chroma(chunks):
    collection = get_chroma_collection()
    existing   = collection.count()

    if existing == len(chunks):
        print(f"ChromaDB loaded — {existing} chunks already stored.\n")
        return collection

    if existing > 0:
        print(f"ChromaDB has {existing} chunks but dataset has {len(chunks)} — rebuilding...")
        collection.delete(ids=collection.get()['ids'])

    print(f"Embedding {len(chunks)} chunks into ChromaDB...")
    batch_size = 50
    for i in range(0, len(chunks), batch_size):
        batch  = chunks[i: i + batch_size]
        ids    = [f"chunk_{i+j}" for j in range(len(batch))]
        texts  = [c['text'] for c in batch]
        metas  = [{'source': c['source'], 'start_line': c['start_line'],
                   'end_line': c['end_line'], 'type': c.get('type', 'txt')} for c in batch]
        embeds = [ollama.embed(model=EMBEDDING_MODEL, input=t)['embeddings'][0] for t in texts]
        collection.add(ids=ids, embeddings=embeds, documents=texts, metadatas=metas)
        print(f"  Stored {min(i+batch_size, len(chunks))}/{len(chunks)}", end='\r')

    print(f"\nChromaDB ready — {collection.count()} chunks stored.\n")
    return collection

# ============================================================
# 3. COSINE SIMILARITY
# ============================================================
def cosine_similarity(a, b):
    dot = sum(x * y for x, y in zip(a, b))
    na  = sum(x**2 for x in a)**0.5
    nb  = sum(x**2 for x in b)**0.5
    return dot / (na * nb) if na and nb else 0.0

# ============================================================
# 4. QUERY EXPANSION
# ============================================================
def expand_query(query):
    """
    Generates 2 alternative phrasings of the query using the LLM.
    Improves recall by covering synonyms and acronyms (e.g. NLP ↔ natural language processing).
    Falls back to the original query on any error.
    """
    prompt = (
        "Rewrite the following search query in 2 different ways to improve document retrieval. "
        "Use synonyms, acronyms, and related terms. Output ONLY the 2 rewrites, one per line, "
        "no numbering, no explanation.\n\n"
        f"Query: {query}"
    )
    try:
        resp = ollama.chat(
            model=LANGUAGE_MODEL,
            messages=[{'role': 'user', 'content': prompt}],
            options={"temperature": 0.3}
        )
        lines = [l.strip() for l in resp['message']['content'].strip().splitlines() if l.strip()]
        expansions = lines[:2]  # take at most 2
        return [query] + expansions  # original + 2 rewrites
    except Exception:
        return [query]  # graceful fallback

# ============================================================
# 5. HYBRID SEARCH (BM25 + ChromaDB Dense)
# ============================================================
def build_bm25_index(chunks):
    return BM25Okapi([c['text'].lower().split() for c in chunks])

def hybrid_retrieve(queries, collection, chunks, bm25_index, top_n=TOP_RETRIEVE, alpha=0.5):
    """
    True hybrid search: fuses BM25 (lexical) + ChromaDB (dense) scores across
    all expanded queries. Takes the best score per chunk across all queries.
    alpha=0.5 = equal weight. Increase alpha to favour dense semantic search.
    """
    fused = {}  # doc_text → (entry, best_score)

    for query in queries:
        # --- Dense retrieval ---
        q_emb   = ollama.embed(model=EMBEDDING_MODEL, input=query)['embeddings'][0]
        results = collection.query(query_embeddings=[q_emb], n_results=min(top_n * 2, collection.count()))

        dense_map = {}
        for doc, meta, dist in zip(results['documents'][0], results['metadatas'][0], results['distances'][0]):
            entry = {
                'text':       doc,
                'source':     meta.get('source', '?'),
                'start_line': meta.get('start_line', 0),
                'end_line':   meta.get('end_line', 0),
                'type':       meta.get('type', 'txt'),
            }
            dense_map[doc] = (entry, 1 - dist)

        # --- BM25 retrieval ---
        tokenized      = query.lower().split()
        bm25_scores_raw = bm25_index.get_scores(tokenized)
        bm25_max       = max(bm25_scores_raw) if max(bm25_scores_raw) > 0 else 1.0
        bm25_norm      = [s / bm25_max for s in bm25_scores_raw]

        # --- Fuse scores for this query ---
        for doc, (entry, dense_score) in dense_map.items():
            bm25_score = 0.0
            for idx, c in enumerate(chunks):
                if c['text'] == doc:
                    bm25_score = bm25_norm[idx]
                    break
            score = alpha * dense_score + (1 - alpha) * bm25_score
            # Keep the best score seen across all expanded queries
            if doc not in fused or score > fused[doc][1]:
                fused[doc] = (entry, score)

    return sorted(fused.values(), key=lambda x: x[1], reverse=True)[:top_n]

# ============================================================
# 6. LLM RERANKER
# ============================================================
def rerank(query, chunks, top_n=TOP_RERANK):
    scored = []
    for entry, sim in chunks:
        prompt = (
            f"On a scale of 1-10, how relevant is the following text to the query?\n"
            f"Query: {query}\nText: {entry['text']}\n"
            f"Reply with a single integer from 1 to 10 and nothing else."
        )
        try:
            resp = ollama.chat(
                model=LANGUAGE_MODEL,
                messages=[{"role": "user", "content": prompt}],
                options={"temperature": 0}
            )
            raw       = resp['message']['content'].strip()
            llm_score = float(re.search(r'\d+', raw).group()) / 10.0
        except Exception:
            llm_score = sim
        scored.append((entry, sim, llm_score))
    scored.sort(key=lambda x: x[2], reverse=True)
    return scored[:top_n]

# ============================================================
# 7. QUERY CLASSIFICATION
# ============================================================
def classify_query(query):
    """
    Classifies query as factual / comparison / general using keyword matching.
    Drives smart_top_n — comparison queries retrieve more candidates.
    """
    q = query.lower()
    comparison_signals = ['compare', 'difference', 'vs', 'versus', 'better', 'worse',
                          'pros and cons', 'which is', 'how does', 'contrast']
    factual_signals    = ['what is', 'what are', 'who is', 'who are', 'when did',
                          'where is', 'how many', 'how much', 'does', 'did', 'has',
                          'have', 'list', 'name', 'define', 'tell me']
    if any(s in q for s in comparison_signals):
        return 'comparison'
    if any(s in q for s in factual_signals):
        return 'factual'
    return 'general'

def smart_top_n(qtype):
    return {'factual': 5, 'comparison': 15, 'general': 10}.get(qtype, TOP_RETRIEVE)

# ============================================================
# 8. CONFIDENCE CHECK
# ============================================================
def check_confidence(chunks):
    if not chunks: return False, 0.0
    best = chunks[0][1]
    return best >= SIMILARITY_THRESHOLD, best

# ============================================================
# 9. LOGGING
# ============================================================
def log_interaction(query, qtype, chunks_used, sim_scores, response):
    entry = {
        'timestamp':       datetime.now().isoformat(),
        'query':           query,
        'query_type':      qtype,
        'chunks_used':     chunks_used,
        'top_similarity':  round(sim_scores[0], 4) if sim_scores else 0,
        'avg_similarity':  round(sum(sim_scores)/len(sim_scores), 4) if sim_scores else 0,
        'response_length': len(response),
    }
    logs = []
    if os.path.exists(LOG_FILE):
        with open(LOG_FILE, 'r') as f:
            try: logs = json.load(f)
            except Exception: logs = []
    logs.append(entry)
    with open(LOG_FILE, 'w') as f:
        json.dump(logs, f, indent=2)

# ============================================================
# 10. STREAMING WITH TYPING INDICATOR
# ============================================================
def stream_response(stream):
    print("\nChatbot: ", end='', flush=True)
    for _ in range(3):
        print('.', end='', flush=True)
        time.sleep(0.3)
    print('\r' + ' ' * 30 + '\r', end='', flush=True)
    print("Chatbot: ", end='', flush=True)
    full = ''
    for chunk in stream:
        c = chunk['message']['content']
        sys.stdout.write(c)
        sys.stdout.flush()
        full += c
    print()
    return full

# ============================================================
# 11. SOURCE LABEL HELPER
# ============================================================
def _source_label(entry):
    """Returns a consistent location label for any doc type."""
    t = entry.get('type', 'txt')
    if t == 'pdf':
        return f"p{entry['start_line']}"
    elif t in ('xlsx', 'csv'):
        return f"row{entry['start_line']}"
    elif t == 'pptx':
        return f"slide{entry['start_line']}"
    elif t == 'html':
        return f"s{entry['start_line']}"
    else:
        return f"L{entry['start_line']}-{entry['end_line']}"

# ============================================================
# 12. CORE PIPELINE
# ============================================================
def run_pipeline(query, collection, chunks, bm25_index, conversation_history, streamlit_mode=False):
    qtype     = classify_query(query)
    top_n     = smart_top_n(qtype)
    queries   = expand_query(query)
    retrieved = hybrid_retrieve(queries, collection, chunks, bm25_index, top_n=top_n)
    is_confident, best_score = check_confidence(retrieved)
    reranked  = rerank(query, retrieved, top_n=TOP_RERANK)

    context_lines = []
    for e, _, _ in reranked:
        label = _source_label(e)
        context_lines.append(f" - [{e['source']} {label}] {e['text']}")
    context = '\n'.join(context_lines)

    instruction_prompt = (
        "You are a helpful chatbot.\n"
        "Use only the following context to answer the question.\n"
        "Do not make up new information.\n"
        "Cite sources at the end of your answer.\n\n"
        f"{context}"
    )

    conversation_history.append({'role': 'user', 'content': query})
    stream = ollama.chat(
        model=LANGUAGE_MODEL,
        messages=[{'role': 'system', 'content': instruction_prompt}, *conversation_history],
        stream=True,
    )

    full_response = (''.join(c['message']['content'] for c in stream)
                     if streamlit_mode else stream_response(stream))

    conversation_history.append({'role': 'assistant', 'content': full_response})

    sim_scores = [s for _, s, _ in reranked]
    log_interaction(query, qtype, len(reranked), sim_scores, full_response)

    return {
        'response':     full_response,
        'query_type':   qtype,
        'queries':      queries,
        'is_confident': is_confident,
        'best_score':   best_score,
        'retrieved':    retrieved,
        'reranked':     reranked,
    }

# ============================================================
# 13. AGENT — Tool Calling
# ============================================================
AGENT_SYSTEM_PROMPT = """You are an AI agent. You must ONLY respond with tool calls — no explanations, no extra text.

Available tools:
1. rag_search - search the knowledge base for information
2. calculator - evaluate a math expression
3. summarise  - summarise a piece of text
4. finish     - return the final answer to the user

You MUST respond in EXACTLY this format with NO other text before or after:
TOOL: tool_name(your argument here)

Examples:
TOOL: rag_search(does the candidate have NLP experience)
TOOL: calculator(16 * 365)
TOOL: summarise(cats sleep a lot and are nocturnal hunters...)
TOOL: finish(Yes, the candidate has NLP experience including POS tagging and language modeling.)

Rules:
- Never write anything except a single TOOL: line
- Always end with TOOL: finish(your final answer)
- Use rag_search first to find information before answering
- Do not explain yourself or add any commentary
- The finish argument must be a clean, direct answer in plain English — NEVER paste raw bullet points or document chunks into finish
- Once you have enough information, call finish IMMEDIATELY
- For simple math questions, call calculator once then finish
- For simple factual questions, call rag_search once then finish
"""

def tool_rag_search(query, collection, chunks, bm25_index):
    """Returns retrieved chunks with source labels for grounded synthesis."""
    queries   = expand_query(query)
    retrieved = hybrid_retrieve(queries, collection, chunks, bm25_index, top_n=5)
    reranked  = rerank(query, retrieved, top_n=3)
    lines = []
    for e, sim, _ in reranked:
        label = _source_label(e)
        lines.append(f"- [{e['source']} {label}] {e['text']}")
    return '\n'.join(lines)

def _synthesize_rag_answer(question, raw_context):
    """
    Takes the raw retrieved context from tool_rag_search and asks the LLM
    to produce a clean, direct answer — instead of returning the chunks verbatim.
    This is the core fix for the agent dumping raw results to the user.
    """
    prompt = (
        "You are a helpful assistant. Answer the question below using ONLY the "
        "provided context. Be concise and direct. Do not repeat the context — "
        "just answer the question. Cite the source filename at the end.\n\n"
        f"Context:\n{raw_context}\n\n"
        f"Question: {question}\n\n"
        "Answer:"
    )
    try:
        resp = ollama.chat(
            model=LANGUAGE_MODEL,
            messages=[{'role': 'user', 'content': prompt}],
            options={"temperature": 0}
        )
        return resp['message']['content'].strip()
    except Exception as e:
        # Fallback: return the raw context if synthesis fails
        return raw_context

def tool_calculator(expression):
    try:
        allowed = set('0123456789+-*/(). ')
        if not all(c in allowed for c in expression):
            return "Error: unsafe expression"
        return str(eval(expression))
    except Exception as e:
        return f"Error: {e}"

def tool_summarise(text):
    resp = ollama.chat(
        model=LANGUAGE_MODEL,
        messages=[{'role': 'user', 'content': f"Summarise this in 2-3 sentences:\n{text}"}]
    )
    return resp['message']['content'].strip()

def parse_tool_call(response_text):
    match = re.search(r'(?i)TOOL:\s*(\w+)\s*\(\s*(.+?)\s*\)', response_text, re.DOTALL)
    if match:
        return match.group(1).strip().lower(), match.group(2).strip()
    match = re.search(r'(?i)TOOL:\s*(\w+)\s+(.+)', response_text)
    if match:
        return match.group(1).strip().lower(), match.group(2).strip()
    return None, None

def run_agent(user_query, collection, chunks, bm25_index, max_steps=8, streamlit_mode=False):
    messages = [
        {'role': 'system', 'content': AGENT_SYSTEM_PROMPT},
        {'role': 'user',   'content': user_query},
    ]

    steps            = []
    answer           = None
    bad_format_count = 0

    for step in range(max_steps):
        resp     = ollama.chat(model=LANGUAGE_MODEL, messages=messages)
        raw_text = resp['message']['content'].strip()
        tool_name, tool_arg = parse_tool_call(raw_text)

        if not tool_name:
            bad_format_count += 1
            if bad_format_count <= 2:
                if not streamlit_mode:
                    print(f"\n  [Agent] Bad format (attempt {bad_format_count}/2), retrying...")
                messages.append({'role': 'assistant', 'content': raw_text})
                messages.append({'role': 'user', 'content':
                    'Wrong format. You must respond with ONLY this format — nothing else:\n'
                    'TOOL: tool_name(argument)\n'
                    'Example: TOOL: rag_search(cat sleep hours)'})
                continue
            else:
                answer = raw_text
                steps.append({'step': step+1, 'tool': 'none', 'arg': '', 'result': raw_text})
                break

        bad_format_count = 0

        if tool_name == 'finish':
            answer = tool_arg
            steps.append({'step': step+1, 'tool': 'finish', 'arg': tool_arg, 'result': tool_arg})
            break

        if tool_name == 'rag_search':
            result = tool_rag_search(tool_arg, collection, chunks, bm25_index)
        elif tool_name == 'calculator':
            result = tool_calculator(tool_arg)
        elif tool_name == 'summarise':
            result = tool_summarise(tool_arg)
        else:
            result = f"Unknown tool '{tool_name}'. Available: rag_search, calculator, summarise, finish"

        steps.append({'step': step+1, 'tool': tool_name, 'arg': tool_arg, 'result': result})

        if not streamlit_mode:
            print(f"\n  [Agent Step {step+1}] {tool_name}({tool_arg[:60]}...)"
                  if len(tool_arg) > 60 else f"\n  [Agent Step {step+1}] {tool_name}({tool_arg})")
            print(f"  → {result[:120]}..." if len(result) > 120 else f"  → {result}")

        if tool_name == 'calculator' and not result.startswith('Error'):
            answer = f"{tool_arg} = {result}"
            steps.append({'step': step+2, 'tool': 'finish', 'arg': answer, 'result': answer})
            if not streamlit_mode:
                print(f"\n  [Agent Step {step+2}] finish({answer})")
            break

        if tool_name == 'rag_search':
            # Synthesize a clean answer from the raw retrieved chunks
            answer = _synthesize_rag_answer(user_query, result)
            steps.append({'step': step+2, 'tool': 'finish', 'arg': answer, 'result': answer})
            if not streamlit_mode:
                print(f"\n  [Agent Step {step+2}] finish({answer[:120]}..." if len(answer) > 120 else f"\n  [Agent Step {step+2}] finish({answer})")
            break

        messages.append({'role': 'assistant', 'content': raw_text})
        messages.append({'role': 'user', 'content':
            f"Tool result: {result}\n\n"
            f"Original task: {user_query}\n\n"
            f"If you now have enough information to answer the original task, call:\n"
            f"TOOL: finish(your answer)\n\n"
            f"Otherwise call the next tool. Respond ONLY with a single TOOL: line."})

    if answer is None:
        answer = "Agent reached max steps without a final answer."

    return {'answer': answer, 'steps': steps}

# ============================================================
# 14. BENCHMARKING
# ============================================================
DEFAULT_TEST_CASES = [
    {'question': 'How many hours do cats sleep per day?',    'expected_keywords': ['sleep', '16']},
    {'question': 'Can cats see in dim light?',               'expected_keywords': ['dim', 'light', 'see']},
    {'question': 'How many toes do cats have on front paws?','expected_keywords': ['five', 'toes', 'front']},
    {'question': 'How many whiskers does a cat have?',       'expected_keywords': ['whiskers', '12']},
    {'question': 'Can cats taste sweet food?',               'expected_keywords': ['sweet', 'taste']},
]

def score_faithfulness(response, reranked):
    context = ' '.join(e['text'] for e, _, _ in reranked)
    stopwords = {'a','an','the','is','are','was','were','do','does','it','its',
                 'to','of','in','for','and','or','not','with','on','at','by',
                 'this','that','be','as','i','you','we','they','but','so','if'}
    context_words  = set(w for w in context.lower().split()  if w not in stopwords)
    response_words = set(w for w in response.lower().split() if w not in stopwords)
    if not response_words: return 0.0
    return min(len(response_words & context_words) / max(len(response_words), 1), 1.0)

def score_relevancy(question, response):
    stopwords = {'a','an','the','is','are','was','were','do','does','did','have',
                 'has','can','what','how','why','when','where','who','to','of','in',
                 'it','its','for','and','or','not','with','on','at','by','from'}
    q_words = set(question.lower().split()) - stopwords
    r_words = set(response.lower().split()) - stopwords
    if not q_words: return 0.0
    precision = len(q_words & r_words) / max(len(r_words), 1)
    recall    = len(q_words & r_words) / max(len(q_words), 1)
    if precision + recall == 0: return 0.0
    return min(2 * precision * recall / (precision + recall), 1.0)

def score_keyword_recall(response, keywords):
    if not keywords: return 1.0
    rl = response.lower()
    return sum(1 for kw in keywords if kw.lower() in rl) / len(keywords)

def score_context_relevance(reranked):
    if not reranked: return 0.0
    scores = [sim for _, sim, _ in reranked[:TOP_RERANK]]
    return sum(scores) / len(scores)

def run_benchmark(collection, chunks, bm25_index, test_cases=None):
    test_cases = test_cases or DEFAULT_TEST_CASES
    print("\n" + "="*70)
    print("  BENCHMARKING RAG PIPELINE")
    print("="*70)
    results = []

    for i, tc in enumerate(test_cases):
        q, kw = tc['question'], tc.get('expected_keywords', [])
        print(f"\n[{i+1}/{len(test_cases)}] {q}")

        queries   = expand_query(q)
        retrieved = hybrid_retrieve(queries, collection, chunks, bm25_index, top_n=TOP_RETRIEVE)
        reranked  = rerank(q, retrieved, top_n=TOP_RERANK)

        context = '\n'.join(f" - {e['text']}" for e, _, _ in reranked)
        stream  = ollama.chat(
            model=LANGUAGE_MODEL,
            messages=[
                {'role': 'system', 'content':
                    f"You are a factual assistant. Answer in 1-2 sentences "
                    f"using ONLY the facts below.\n\nFacts:\n{context}"},
                {'role': 'user', 'content': q},
            ], stream=True)
        response = ''.join(c['message']['content'] for c in stream)

        faith = score_faithfulness(response, reranked)
        rel   = score_relevancy(q, response)
        kwr   = score_keyword_recall(response, kw)
        ctx   = score_context_relevance(reranked)
        ovr   = (faith + rel + kwr + ctx) / 4

        results.append({'question': q, 'faithfulness': round(faith,3),
                        'answer_relevancy': round(rel,3), 'keyword_recall': round(kwr,3),
                        'context_relevance': round(ctx,3), 'overall': round(ovr,3)})
        print(f"  faith={faith:.2f} rel={rel:.2f} kw={kwr:.2f} ctx={ctx:.2f} overall={ovr:.2f}")

    def avg(k): return sum(r[k] for r in results) / len(results)
    summary = {k: round(avg(k), 3) for k in
               ['faithfulness','answer_relevancy','keyword_recall','context_relevance','overall']}
    def bar(s): return '[' + '█'*int(s*20) + '░'*(20-int(s*20)) + ']'

    print("\n" + "="*70 + "\n  SUMMARY\n" + "="*70)
    for k, v in summary.items():
        print(f"  {k:<25} {v:>6.3f}  {bar(v)}")

    runs = []
    if os.path.exists(BENCHMARK_FILE):
        with open(BENCHMARK_FILE) as f:
            try: runs = json.load(f)
            except: runs = []
    if runs:
        prev = runs[-1]['summary']
        print("\n  vs PREVIOUS RUN")
        for k in summary:
            d = summary[k] - prev.get(k, 0)
            print(f"  {k:<25} {prev.get(k,0):>6.3f} → {summary[k]:>6.3f}  "
                  f"{'▲' if d>0 else '▼' if d<0 else '─'}{abs(d):.3f}")

    runs.append({'timestamp': datetime.now().isoformat(), 'summary': summary, 'results': results})
    with open(BENCHMARK_FILE, 'w') as f:
        json.dump(runs, f, indent=2)
    print(f"\n  Saved to '{BENCHMARK_FILE}'\n" + "="*70)
    return summary

# ============================================================
# 15. TERMINAL CHATBOT (loop)
# ============================================================
def run_terminal(collection, chunks, bm25_index):
    conv = []
    print("="*60 + "\n  RAG Chatbot — Full Pipeline\n" + "="*60)
    print("Commands: 'exit' quit | 'agent: <q>' use agent mode\n")

    while True:
        try:
            query = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nGoodbye!"); break

        if not query: continue
        if query.lower() in ['exit','quit','bye']:
            print("Goodbye!"); break

        if query.lower().startswith('agent:'):
            q = query[6:].strip()
            print("\n[Agent mode]")
            result = run_agent(q, collection, chunks, bm25_index)
            print(f"\nAgent answer: {result['answer']}")
        else:
            result = run_pipeline(query, collection, chunks, bm25_index, conv)
            if not result['is_confident']:
                print(f"[Warning] Low confidence ({result['best_score']:.2f})")
            print(f"\n[type:{result['query_type']} | expanded:{len(result['queries'])} queries]")
            print(f"Before rerank: {len(result['retrieved'])} chunks | After: {TOP_RERANK} chunks")
        print("-"*60)

# ============================================================
# 16. STREAMLIT UI
# ============================================================
def run_streamlit(collection, chunks, bm25_index):
    import streamlit as st

    st.set_page_config(page_title="RAG Chatbot", page_icon="🐱", layout="wide")
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Mono:wght@400;600&family=IBM+Plex+Sans:wght@300;400;600&display=swap');
    html,body,[class*="css"]{font-family:'IBM Plex Sans',sans-serif;background:#0d0d0d;color:#e8e8e8}
    .stApp{background:#0d0d0d}
    .rag-title{font-family:'IBM Plex Mono',monospace;font-size:2rem;font-weight:600;color:#f0c040;letter-spacing:-.02em}
    .rag-sub{font-family:'IBM Plex Mono',monospace;font-size:.75rem;color:#444;margin-bottom:1.5rem}
    .msg-user{background:#1a1a1a;border-left:3px solid #f0c040;padding:.8rem 1rem;margin:.4rem 0;border-radius:0 8px 8px 0}
    .msg-bot{background:#141414;border-left:3px solid #3a9ad9;padding:.8rem 1rem;margin:.4rem 0;border-radius:0 8px 8px 0;line-height:1.6}
    .msg-agent{background:#0f1a0f;border-left:3px solid #4caf50;padding:.8rem 1rem;margin:.4rem 0;border-radius:0 8px 8px 0}
    .msg-label{font-family:'IBM Plex Mono',monospace;font-size:.65rem;color:#444;margin-bottom:.2rem;text-transform:uppercase;letter-spacing:.1em}
    .chunk{background:#111;border:1px solid #1e1e1e;border-radius:6px;padding:.5rem .7rem;margin:.25rem 0;font-family:'IBM Plex Mono',monospace;font-size:.7rem;color:#888}
    .cs{color:#f0c040;font-weight:600}.src{color:#3a9ad9}
    .step{background:#0a1a0a;border:1px solid #1a2a1a;border-radius:6px;padding:.5rem .7rem;margin:.2rem 0;font-family:'IBM Plex Mono',monospace;font-size:.7rem;color:#4caf50}
    .badge{display:inline-block;font-family:'IBM Plex Mono',monospace;font-size:.65rem;padding:.15rem .4rem;border-radius:3px;margin:.1rem}
    .b-fact{background:#1a3a1a;color:#4caf50}.b-comp{background:#1a2a3a;color:#3a9ad9}.b-gen{background:#2a1a2a;color:#ce93d8}
    .b-ok{background:#1a3a1a;color:#4caf50}.b-warn{background:#3a2a00;color:#f0c040}
    .stat{font-family:'IBM Plex Mono',monospace;font-size:.72rem;color:#444;padding:.25rem 0;border-bottom:1px solid #151515;display:flex;justify-content:space-between}
    .sv{color:#f0c040}
    .stTextInput>div>div>input{background:#1a1a1a!important;border:1px solid #2a2a2a!important;color:#e8e8e8!important;font-family:'IBM Plex Mono',monospace!important;border-radius:6px!important}
    .stButton>button{background:#f0c040!important;color:#0d0d0d!important;font-family:'IBM Plex Mono',monospace!important;font-weight:600!important;border:none!important;border-radius:6px!important}
    [data-testid="stSidebar"]{background:#0a0a0a!important;border-right:1px solid #151515}
    hr{border-color:#1a1a1a!important}
    </style>
    """, unsafe_allow_html=True)

    for k,v in [('conv',[]),('display',[]),('total',0),('last',None),('mode','chat')]:
        if k not in st.session_state: st.session_state[k] = v

    col_main, col_side = st.columns([3,1])

    with col_main:
        st.markdown('<div class="rag-title">// RAG Chatbot</div>', unsafe_allow_html=True)
        st.markdown(
            '<div class="rag-sub">chunking · hybrid search · reranking · agent · '
            'PDF · TXT · DOCX · XLSX · PPTX · CSV · MD · HTML</div>',
            unsafe_allow_html=True
        )

        mode = st.radio("Mode:", ["Chat", "Agent"], horizontal=True,
                        index=0 if st.session_state.mode=='chat' else 1)
        st.session_state.mode = mode.lower()
        st.markdown("---")

        for msg in st.session_state.display:
            css = {'user':'msg-user','assistant':'msg-bot','agent':'msg-agent'}.get(msg['role'],'msg-bot')
            lbl = msg['role']
            st.markdown(f'<div class="{css}"><div class="msg-label">{lbl}</div>{msg["content"]}</div>',
                        unsafe_allow_html=True)

        st.markdown("---")
        with st.form('chat', clear_on_submit=True):
            placeholder = "Ask a question..." if st.session_state.mode == 'chat' else "Give the agent a task..."
            user_input = st.text_input("Input:", placeholder=placeholder, label_visibility='collapsed')
            submitted  = st.form_submit_button("Send →")

        if submitted and user_input.strip():
            st.session_state.display.append({'role':'user','content': user_input})
            if st.session_state.mode == 'agent':
                with st.spinner("Agent thinking..."):
                    res = run_agent(user_input, collection, chunks, bm25_index, streamlit_mode=True)
                steps_html = ''.join(
                    f'<div class="step">Step {s["step"]}: {s["tool"]}({s["arg"][:50]}...) → {s["result"][:80]}...</div>'
                    if len(s["arg"])>50 else
                    f'<div class="step">Step {s["step"]}: {s["tool"]}({s["arg"]}) → {s["result"][:80]}</div>'
                    for s in res['steps']
                )
                content = f"{steps_html}<br/><strong>Answer:</strong> {res['answer']}"
                st.session_state.display.append({'role':'agent','content': content})
                st.session_state.last = {'type':'agent','data': res}
            else:
                with st.spinner("Thinking..."):
                    res = run_pipeline(user_input, collection, chunks, bm25_index,
                                       st.session_state.conv, streamlit_mode=True)
                st.session_state.display.append({'role':'assistant','content': res['response']})
                st.session_state.last = {'type':'chat','data': res}
            st.session_state.total += 1
            st.rerun()

    with col_side:
        st.markdown("### Pipeline")
        if st.session_state.last:
            d = st.session_state.last['data']
            if st.session_state.last['type'] == 'chat':
                qt        = d['query_type']
                badge_cls = {'factual':'b-fact','comparison':'b-comp','general':'b-gen'}.get(qt,'b-gen')
                st.markdown(f'<span class="badge {badge_cls}">{qt}</span>', unsafe_allow_html=True)
                cc = 'b-ok' if d['is_confident'] else 'b-warn'
                cl = f"conf:{d['best_score']:.2f}" if d['is_confident'] else f"low:{d['best_score']:.2f}"
                st.markdown(f'<span class="badge {cc}">{cl}</span>', unsafe_allow_html=True)
                st.markdown("---")
                st.markdown("**Before rerank**")
                for e,s in d['retrieved'][:4]:
                    label = _source_label(e)
                    st.markdown(
                        f'<div class="chunk"><span class="cs">{s:.3f}</span> '
                        f'<span class="src">[{e["source"]} {label}]</span><br/>'
                        f'{e["text"][:55]}...</div>',
                        unsafe_allow_html=True
                    )
                st.markdown("**After rerank**")
                for e,sim,rs in d['reranked']:
                    label = _source_label(e)
                    st.markdown(
                        f'<div class="chunk"><span class="cs">sim:{sim:.2f} re:{rs:.0f}</span> '
                        f'<span class="src">[{e["source"]} {label}]</span><br/>'
                        f'{e["text"][:55]}...</div>',
                        unsafe_allow_html=True
                    )
            else:
                st.markdown("**Agent Steps**")
                for s in d['steps']:
                    st.markdown(f'<div class="step">{s["step"]}. {s["tool"]}</div>', unsafe_allow_html=True)
            st.markdown("---")

        st.markdown("**Session**")
        st.markdown(f'<div class="stat">Queries <span class="sv">{st.session_state.total}</span></div>', unsafe_allow_html=True)
        st.markdown(f'<div class="stat">Memory <span class="sv">{len(st.session_state.conv)//2} turns</span></div>', unsafe_allow_html=True)
        st.markdown(f'<div class="stat">Chunks <span class="sv">{len(chunks)}</span></div>', unsafe_allow_html=True)
        st.markdown(f'<div class="stat">Mode <span class="sv">{st.session_state.mode}</span></div>', unsafe_allow_html=True)
        st.markdown("---")

        # Show doc type breakdown in sidebar
        type_counts = {}
        for c in chunks:
            t = c.get('type','?')
            type_counts[t] = type_counts.get(t, 0) + 1
        st.markdown("**Document Types**")
        for t, cnt in sorted(type_counts.items()):
            st.markdown(f'<div class="stat">{t.upper()} <span class="sv">{cnt}</span></div>', unsafe_allow_html=True)
        st.markdown("---")

        if st.button("Clear Chat"):
            st.session_state.conv=[]; st.session_state.display=[]
            st.session_state.last=None; st.session_state.total=0
            st.rerun()

# ============================================================
# ENTRY POINT
# ============================================================
if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='RAG Chatbot')
    parser.add_argument('--benchmark', action='store_true', help='Run benchmark')
    parser.add_argument('--agent',     action='store_true', help='Agent mode in terminal')
    args = parser.parse_args()

    print("="*60)
    print("  Initializing RAG Pipeline")
    print("="*60)
    print(f"  Docs root:  {DOCS_ROOT}/")
    for t, folder in DOC_FOLDERS.items():
        print(f"    {t.upper():<8} → {folder}/")
    print(f"  Vector DB:  ChromaDB (persistent @ {CHROMA_DIR})")
    print(f"  Reranker:   LLM-based")
    print(f"  Smart mis-placed file detection: ENABLED")
    print("="*60 + "\n")

    ensure_folders()
    chunks     = chunk_all_documents()
    collection = build_or_load_chroma(chunks)
    bm25       = build_bm25_index(chunks)

    if args.benchmark:
        run_benchmark(collection, chunks, bm25)
    elif args.agent:
        print("Agent mode — type your task:\n")
        while True:
            try: q = input("Task: ").strip()
            except (EOFError, KeyboardInterrupt): print("\nGoodbye!"); break
            if not q: continue
            if q.lower() in ['exit','quit']: print("Goodbye!"); break
            res = run_agent(q, collection, chunks, bm25)
            print(f"\nFinal answer: {res['answer']}\n" + "-"*60)
    else:
        run_terminal(collection, chunks, bm25)
